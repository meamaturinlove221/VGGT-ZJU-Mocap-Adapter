from __future__ import annotations

import json
import hashlib
import math
import textwrap
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable
from urllib.request import urlretrieve

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = ROOT / "logs" / "modal_phase5" / "reports"
OUT_PPTX = REPORT_DIR / "orig_vggt_local_code_mods_presentation_latest.pptx"
OUT_MD = REPORT_DIR / "orig_vggt_local_code_mods_presentation_outline_latest.md"
OUT_CODE_PPTX = REPORT_DIR / "orig_vggt_code_mods_and_finetune_presentation_latest.pptx"
OUT_CODE_MD = REPORT_DIR / "orig_vggt_code_mods_and_finetune_presentation_outline_latest.md"
OUT_PROBE_PPTX = REPORT_DIR / "orig_vggt_one_step_stepcurve_presentation_latest.pptx"
OUT_PROBE_MD = REPORT_DIR / "orig_vggt_one_step_stepcurve_presentation_outline_latest.md"
OUT_MASK_PPTX = REPORT_DIR / "orig_vggt_mask_boost_ghost_presentation_latest.pptx"
OUT_MASK_MD = REPORT_DIR / "orig_vggt_mask_boost_ghost_presentation_outline_latest.md"
OUT_PROMPTS_MD = REPORT_DIR / "orig_vggt_ppt_image_prompts_latest.md"
CACHE_DIR = REPORT_DIR / "_orig_vggt_local_code_mods_presentation_cache"
OFFICIAL_ARCH_URL = "https://vgg-t.github.io/resources/architecture_v4.png"
OFFICIAL_ARCH_PNG = CACHE_DIR / "official_vggt_architecture.png"
LOCAL_ARCH_PNG = CACHE_DIR / "current_local_vggt_architecture.png"
ORIGINAL_REPO_ROOT = Path(r"G:\项目备份\vggt\vggt-main")

VIEWCOUNT_JSON = REPORT_DIR / "orig_vggt_viewcount_summary_latest.json"
ONE_STEP_JSON = REPORT_DIR / "orig_vggt_one_step_probe_summary_latest.json"
STEPCURVE_JSON = REPORT_DIR / "orig_vggt_stepcurve_probe_summary_latest.json"
MASK_BOOST_JSON = REPORT_DIR / "orig_vggt_mask_boost_probe_summary_latest.json"
MASK_BOOST_EXPLAINER_MD = REPORT_DIR / "orig_vggt_mask_boost_and_ghost_explainer_latest.md"

STEPCURVE_POINT_GRID = REPORT_DIR / "orig_vggt_stepcurve_point_support_grid_latest.png"
STEPCURVE_GHOST_GRID = REPORT_DIR / "orig_vggt_stepcurve_ghost_grid_latest.png"
MASK_BOOST_POINT_GRID = REPORT_DIR / "orig_vggt_mask_boost_point_support_grid_latest.png"
MASK_BOOST_GHOST_GRID = REPORT_DIR / "orig_vggt_mask_boost_ghost_grid_latest.png"

POINT_TAXONOMY_EXAMPLE = (
    ROOT
    / "logs"
    / "modal_phase5"
    / "orig_vggt_stepcurve_probe"
    / "12src_nested"
    / "step0004"
    / "compare"
    / "cat_weight_pred_tgt_subject_bbox.png"
)
GHOST_TAXONOMY_EXAMPLE = (
    ROOT
    / "logs"
    / "modal_phase5"
    / "orig_vggt_stepcurve_probe"
    / "12src_nested"
    / "step0004"
    / "compare"
    / "cat_fg_mask_pred_tgt_step000000.png"
)

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)
TITLE_COLOR = RGBColor(24, 36, 54)
ACCENT = RGBColor(199, 88, 25)
SUBTLE = RGBColor(88, 102, 120)
TEXT = RGBColor(35, 35, 35)
BG = RGBColor(248, 246, 241)
BOX = RGBColor(237, 231, 219)

Image.MAX_IMAGE_PIXELS = None
plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


@dataclass
class SlideSpec:
    title: str
    bullets: list[str] = field(default_factory=list)
    formulas: list[str] = field(default_factory=list)
    code_refs: list[str] = field(default_factory=list)
    speaker_notes: str = ""
    image_paths: list[Path] = field(default_factory=list)
    image_caption: str = ""
    layout: str = "bullets"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def format_float(value: float | int | None, digits: int = 4) -> str:
    if value is None:
        return "n/a"
    try:
        value = float(value)
    except Exception:
        return str(value)
    if math.isnan(value) or math.isinf(value):
        return "n/a"
    return f"{value:.{digits}f}"


def wrap_lines(lines: Iterable[str], width: int = 46) -> list[str]:
    out: list[str] = []
    for line in lines:
        if not line:
            out.append("")
            continue
        chunks = textwrap.wrap(line, width=width, break_long_words=False, replace_whitespace=False)
        out.extend(chunks or [""])
    return out


def safe_stem(text: str, limit: int = 48) -> str:
    keep = []
    for ch in text:
        if ch.isalnum() or ch in "._-":
            keep.append(ch)
        else:
            keep.append("_")
    base = "".join(keep).strip("._") or "slide"
    base = base[:limit]
    digest = hashlib.sha1(text.encode("utf-8")).hexdigest()[:8]
    return f"{base}_{digest}"


def ensure_cache_image(src: Path, stem: str, max_w: int = 1800, max_h: int = 1000) -> Path:
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    dst = CACHE_DIR / f"{stem}{src.suffix.lower()}"
    with Image.open(src) as img:
        img = img.convert("RGB")
        scale = min(max_w / img.width, max_h / img.height, 1.0)
        if scale < 1.0:
            img = img.resize((int(img.width * scale), int(img.height * scale)), Image.LANCZOS)
        img.save(dst)
    return dst


def add_textbox(slide, left, top, width, height, text="", font_size=20, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return box


def add_bullets(frame, bullets: list[str], font_size: int = 18, color=TEXT):
    frame.clear()
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = f"- {bullet}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = f"- {bullet}"
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = add_textbox(slide, Inches(0.6), Inches(0.32), Inches(12.1), Inches(0.75), title, 28, True, TITLE_COLOR)
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.95), Inches(12.0), Inches(0.35), subtitle, 12, False, SUBTLE)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), WIDE_W, Inches(0.15))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.color.rgb = ACCENT


def fit_image(slide, img_path: Path, left, top, width, height):
    with Image.open(img_path) as img:
        aspect = img.width / img.height
    box_w = width.inches
    box_h = height.inches
    box_aspect = box_w / box_h
    if aspect >= box_aspect:
        pic_w = width
        pic_h = Inches(box_w / aspect)
        pic_left = left
        pic_top = top + Inches((box_h - pic_h.inches) / 2.0)
    else:
        pic_h = height
        pic_w = Inches(box_h * aspect)
        pic_left = left + Inches((box_w - pic_w.inches) / 2.0)
        pic_top = top
    slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_w, height=pic_h)


def picture_exists(path: Path) -> bool:
    return path.exists() and path.is_file()


def profile_map(rows: list[dict], key: str) -> dict[str, dict]:
    return {row[key]: row for row in rows}


def step_row(profile_summary: dict, step: int) -> dict | None:
    for row in profile_summary.get("steps", []):
        if int(row.get("step", -1)) == int(step):
            return row
    return None


def mask_table_row(profile_summary: dict, alpha: int, step: int) -> dict | None:
    for row in profile_summary.get("table_rows", []):
        if int(row.get("alpha", -1)) == int(alpha) and int(row.get("step", -1)) == int(step):
            return row
    return None


def collect_data() -> dict:
    viewcount = load_json(VIEWCOUNT_JSON)
    one_step = load_json(ONE_STEP_JSON)
    stepcurve = load_json(STEPCURVE_JSON)
    mask_boost = load_json(MASK_BOOST_JSON)
    explainer_text = MASK_BOOST_EXPLAINER_MD.read_text(encoding="utf-8")
    return {
        "viewcount": viewcount,
        "viewcount_map": profile_map(viewcount["rows"], "label"),
        "one_step": one_step,
        "one_step_map": profile_map(one_step["rows"], "profile"),
        "stepcurve": stepcurve,
        "stepcurve_map": profile_map(stepcurve["profiles"], "profile"),
        "mask_boost": mask_boost,
        "mask_boost_map": profile_map(mask_boost["profiles"], "profile"),
        "mask_boost_explainer": explainer_text,
    }


def delta_vs_native_subject_psnr(data: dict, profile: str, step: int) -> tuple[float | None, float | None, float | None]:
    base = step_row(data["stepcurve_map"][profile], step)
    boost_profile = data["mask_boost_map"][profile]
    winner_alpha = int(boost_profile["winner_alpha"])
    boosted = mask_table_row(boost_profile, winner_alpha, step)
    if not base or not boosted:
        return None, None, None
    base_psnr = base.get("subject_psnr")
    boost_psnr = boosted.get("subject_psnr")
    if base_psnr is None or boost_psnr is None:
        return None, None, None
    return float(base_psnr), float(boost_psnr), float(boost_psnr) - float(base_psnr)


def cache_png(name: str) -> Path:
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    return CACHE_DIR / f"{name}.png"


def save_fig(fig, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#f8f6f1")
    plt.close(fig)
    return path


def save_with_fallback_bytes(path: Path, writer) -> Path:
    try:
        writer(path)
        return path
    except PermissionError:
        alt = path.with_name(f"{path.stem}_regenerated{path.suffix}")
        writer(alt)
        return alt


def valid_points(rows: list[dict], key: str) -> tuple[list[int], list[float]]:
    xs: list[int] = []
    ys: list[float] = []
    for row in rows:
        value = row.get(key)
        try:
            v = float(value)
        except Exception:
            continue
        if math.isnan(v) or math.isinf(v):
            continue
        xs.append(int(row["step"]))
        ys.append(v)
    return xs, ys


def draw_arch_box(
    ax,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    fc: str = "#ffffff",
    ec: str = "#555555",
    lw: float = 2.0,
    fontsize: int = 15,
    text_color: str = "#111111",
    radius: float = 0.02,
    alpha: float = 1.0,
    weight: str = "bold",
    linestyle: str = "-",
) -> FancyBboxPatch:
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle=f"round,pad=0.008,rounding_size={radius}",
        facecolor=fc,
        edgecolor=ec,
        linewidth=lw,
        linestyle=linestyle,
        alpha=alpha,
        zorder=2,
    )
    ax.add_patch(patch)
    ax.text(
        x + w / 2,
        y + h / 2,
        text,
        ha="center",
        va="center",
        fontsize=fontsize,
        color=text_color,
        fontweight=weight,
        zorder=3,
    )
    return patch


def draw_arch_arrow(ax, x1: float, y1: float, x2: float, y2: float, color: str = "#7a7a7a", lw: float = 2.0) -> None:
    arrow = FancyArrowPatch(
        (x1, y1),
        (x2, y2),
        arrowstyle="simple",
        mutation_scale=18,
        linewidth=0,
        color=color,
        alpha=0.9,
        zorder=1,
    )
    ax.add_patch(arrow)


def draw_token_stack(
    ax,
    x: float,
    y: float,
    color: str,
    *,
    count: int = 4,
    w: float = 0.018,
    h: float = 0.048,
    gap: float = 0.008,
    add_camera_token: bool = False,
    camera_color: str = "#c96310",
) -> None:
    for idx in range(count):
        yy = y + idx * (h + gap)
        rect = Rectangle((x, yy), w, h, facecolor=color, edgecolor="#ffffff", linewidth=1.6, zorder=3)
        ax.add_patch(rect)
    if add_camera_token:
        yy = y + count * (h + gap) + gap * 0.6
        rect = Rectangle((x, yy), w, h, facecolor=camera_color, edgecolor="#ffffff", linewidth=1.6, zorder=3)
        ax.add_patch(rect)


def draw_input_panel(ax, x: float, y: float, w: float, h: float, tint: str) -> None:
    rect = Rectangle((x, y), w, h, facecolor=tint, edgecolor="#666666", linewidth=1.2, zorder=2)
    ax.add_patch(rect)
    cols, rows = 5, 4
    for i in range(1, cols):
        xx = x + w * i / cols
        ax.plot([xx, xx], [y, y + h], color="#ffffff", lw=2, alpha=0.9, zorder=3)
    for j in range(1, rows):
        yy = y + h * j / rows
        ax.plot([x, x + w], [yy, yy], color="#ffffff", lw=2, alpha=0.9, zorder=3)
    ax.plot([x + 0.02 * w, x + 0.42 * w, x + 0.82 * w], [y + 0.18 * h, y + 0.75 * h, y + 0.35 * h], color="#4b4b4b", lw=3, alpha=0.6, zorder=3)
    ax.plot([x + 0.12 * w, x + 0.68 * w], [y + 0.92 * h, y + 0.58 * h], color="#4b4b4b", lw=3, alpha=0.55, zorder=3)


def generate_current_architecture_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(20, 7.2))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    orange = "#c96310"
    peach = "#f4b27c"
    green = "#a7cf88"
    blue = "#8db3d9"
    purple = "#7b61c8"
    yellow_fill = "#fde7a2"
    blue_fill = "#dce7f7"
    green_fill = "#c5dfb7"
    shell_fill = "#f6f1e7"
    shell_edge = "#9a8d74"

    ax.text(0.03, 0.91, "Input Views", fontsize=20, fontweight="bold", color="#222222")
    draw_input_panel(ax, 0.02, 0.63, 0.07, 0.13, "#d8d0c5")
    draw_input_panel(ax, 0.02, 0.41, 0.07, 0.13, "#d2d8cf")
    draw_input_panel(ax, 0.02, 0.19, 0.07, 0.13, "#d3d8e4")

    for y in (0.66, 0.44, 0.22):
        draw_arch_arrow(ax, 0.10, y + 0.03, 0.15, y + 0.03)
    draw_token_stack(ax, 0.17, 0.60, peach)
    draw_token_stack(ax, 0.17, 0.38, green)
    draw_token_stack(ax, 0.17, 0.16, blue)
    ax.text(
        0.155,
        0.52,
        "Patch / image tokens",
        fontsize=14,
        fontweight="bold",
        color="#222222",
        bbox=dict(boxstyle="round,pad=0.18", fc="#fcfbf7", ec="none", alpha=0.95),
        zorder=5,
    )

    for y in (0.66, 0.44, 0.22):
        draw_arch_arrow(ax, 0.205, y + 0.03, 0.245, y + 0.03)
    draw_token_stack(ax, 0.26, 0.60, peach, add_camera_token=True, camera_color=orange)
    draw_token_stack(ax, 0.26, 0.38, green, add_camera_token=False)
    draw_token_stack(ax, 0.26, 0.16, blue, add_camera_token=False)
    ax.text(0.245, 0.91, "Add", fontsize=18, fontweight="bold", color="#222222")
    ax.text(0.238, 0.865, "camera token", fontsize=15, fontweight="bold", color=orange)
    draw_arch_arrow(ax, 0.269, 0.855, 0.269, 0.80, color="#8f8f8f")

    draw_arch_box(ax, 0.335, 0.11, 0.27, 0.80, "", fc="none", ec="#bbbbbb", lw=3.0, radius=0.03, linestyle=(0, (4, 4)))
    ax.text(0.41, 0.87, "Global\nAttention", ha="center", va="center", fontsize=19, fontweight="bold", color="#222222")
    ax.text(0.525, 0.87, "Frame\nAttention", ha="center", va="center", fontsize=19, fontweight="bold", color="#222222")

    draw_arch_arrow(ax, 0.288, 0.66, 0.33, 0.66)
    draw_arch_arrow(ax, 0.288, 0.44, 0.33, 0.44)
    draw_arch_arrow(ax, 0.288, 0.22, 0.33, 0.22)

    draw_arch_box(ax, 0.375, 0.18, 0.04, 0.66, "", fc=blue_fill, ec="#8aa5d8", lw=3.0, radius=0.02)
    draw_token_stack(ax, 0.386, 0.60, peach, w=0.018, h=0.045)
    draw_token_stack(ax, 0.386, 0.38, green, w=0.018, h=0.045)
    draw_token_stack(ax, 0.386, 0.16, blue, w=0.018, h=0.045)
    for y in (0.66, 0.44, 0.22):
        draw_arch_arrow(ax, 0.42, y, 0.47, y)

    draw_arch_box(ax, 0.485, 0.59, 0.045, 0.20, "", fc=yellow_fill, ec="#f2b300", lw=3.0, radius=0.02)
    draw_arch_box(ax, 0.485, 0.37, 0.045, 0.20, "", fc=yellow_fill, ec="#f2b300", lw=3.0, radius=0.02)
    draw_arch_box(ax, 0.485, 0.15, 0.045, 0.20, "", fc=yellow_fill, ec="#f2b300", lw=3.0, radius=0.02)
    draw_token_stack(ax, 0.497, 0.60, peach, w=0.018, h=0.043)
    draw_token_stack(ax, 0.497, 0.38, green, w=0.018, h=0.043)
    draw_token_stack(ax, 0.497, 0.16, blue, w=0.018, h=0.043)

    for y in (0.66, 0.44, 0.22):
        draw_arch_arrow(ax, 0.535, y, 0.59, y)

    draw_token_stack(ax, 0.61, 0.61, orange, w=0.02, h=0.052, gap=0.012)
    draw_token_stack(ax, 0.61, 0.39, peach, w=0.02, h=0.052, gap=0.012)
    draw_token_stack(ax, 0.61, 0.17, green, w=0.02, h=0.052, gap=0.012)

    draw_arch_box(ax, 0.69, 0.73, 0.09, 0.10, "Camera\nHead", fc="#fff4ec", ec=orange, lw=2.5, fontsize=16)
    draw_arch_box(ax, 0.69, 0.36, 0.09, 0.25, "DPT\nHeads", fc=green_fill, ec="#4f7f2f", lw=2.8, fontsize=22)
    draw_arch_arrow(ax, 0.635, 0.74, 0.685, 0.78)
    draw_arch_arrow(ax, 0.635, 0.49, 0.685, 0.49)

    draw_arch_box(ax, 0.83, 0.74, 0.10, 0.08, "Cameras", fc="#f1ebff", ec=purple, lw=2.5, fontsize=18)
    draw_arch_box(ax, 0.83, 0.55, 0.10, 0.08, "Depth", fc="#e7f3e2", ec="#4f7f2f", lw=2.5, fontsize=18)
    draw_arch_box(ax, 0.83, 0.42, 0.10, 0.08, "Point map", fc="#e7f3e2", ec="#4f7f2f", lw=2.5, fontsize=18)
    draw_arch_box(ax, 0.83, 0.29, 0.10, 0.08, "Confidence", fc="#fff2e5", ec=orange, lw=2.5, fontsize=16)
    draw_arch_arrow(ax, 0.782, 0.78, 0.825, 0.78)
    draw_arch_arrow(ax, 0.782, 0.56, 0.825, 0.59)
    draw_arch_arrow(ax, 0.782, 0.49, 0.825, 0.46)
    draw_arch_arrow(ax, 0.782, 0.42, 0.825, 0.33)

    draw_arch_box(ax, 0.22, 0.005, 0.74, 0.17, "", fc="none", ec="#b9b0a0", lw=2.5, radius=0.02, linestyle=(0, (4, 4)))
    ax.text(
        0.245,
        0.188,
        "Local training / analysis shell (new)",
        fontsize=14,
        fontweight="bold",
        color="#555555",
        bbox=dict(boxstyle="round,pad=0.12", fc="#fcfbf7", ec="none", alpha=0.95),
        zorder=5,
    )

    draw_arch_box(ax, 0.25, 0.055, 0.12, 0.08, "Geometry\nTeacher", fc="#fff1e8", ec=orange, lw=2.6, fontsize=15)
    draw_arch_box(ax, 0.40, 0.055, 0.11, 0.08, "Geometry\nCache", fc=shell_fill, ec=shell_edge, lw=2.4, fontsize=16)
    draw_arch_box(ax, 0.55, 0.035, 0.16, 0.10, "Supervision\nRouter", fc="#edf5e8", ec="#638d45", lw=2.6, fontsize=15)
    ax.text(0.63, 0.055, "confidence  |  foreground\nmultiview  |  camera / normal / reproj", ha="center", va="center", fontsize=10.0, color="#34512b", zorder=4)
    draw_arch_box(ax, 0.75, 0.035, 0.12, 0.10, "Compare\n& Metrics", fc="#edf1fb", ec="#6882b4", lw=2.6, fontsize=14)

    draw_arch_box(ax, 0.89, 0.103, 0.07, 0.037, "One-step", fc="#fbede2", ec=orange, lw=2.0, fontsize=12)
    draw_arch_box(ax, 0.89, 0.060, 0.07, 0.037, "Stepcurve", fc="#eef6ea", ec="#638d45", lw=2.0, fontsize=12)
    draw_arch_box(ax, 0.89, 0.017, 0.07, 0.037, "Mask-boost", fc="#eef0fb", ec="#6882b4", lw=2.0, fontsize=11)

    ax.plot([0.88, 0.88], [0.07, 0.135], color="#7a7a7a", lw=2, zorder=1)
    draw_arch_arrow(ax, 0.87, 0.095, 0.89, 0.122)
    draw_arch_arrow(ax, 0.87, 0.095, 0.89, 0.079)
    draw_arch_arrow(ax, 0.87, 0.095, 0.89, 0.036)

    ax.plot([0.88, 0.88], [0.33, 0.215], color="#7a7a7a", lw=2, zorder=1)
    draw_arch_arrow(ax, 0.88, 0.29, 0.31, 0.13)
    draw_arch_arrow(ax, 0.88, 0.55, 0.31, 0.13)
    draw_arch_arrow(ax, 0.88, 0.42, 0.31, 0.13)
    draw_arch_arrow(ax, 0.88, 0.74, 0.31, 0.13)

    draw_arch_arrow(ax, 0.37, 0.095, 0.40, 0.095)
    draw_arch_arrow(ax, 0.51, 0.095, 0.55, 0.095)
    draw_arch_arrow(ax, 0.71, 0.085, 0.75, 0.085)

    ax.text(0.03, 0.03, "Style reference: official VGGT architecture. Content: current local geometry-training-analysis pipeline.", fontsize=11, color="#666666")

    return save_fig(fig, LOCAL_ARCH_PNG)


def generate_code_touch_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.6))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    ax.text(0.06, 0.92, "Original repo: directly touched", fontsize=20, fontweight="bold", color="#222222")
    ax.text(0.54, 0.92, "Local repo: newly added shell", fontsize=20, fontweight="bold", color="#222222")

    left_boxes = [
        ("Geometry convention\nunprojection + pixel center", 0.78),
        ("Image loading\nmultithread reader", 0.64),
        ("COLMAP bridge\nregistration + BA compatibility", 0.50),
        ("Core structure notes\nclarification only", 0.36),
    ]
    for text, y in left_boxes:
        draw_arch_box(ax, 0.07, y, 0.31, 0.09, text, fc="#fff1e8", ec="#c96310", lw=2.3, fontsize=16)

    right_boxes = [
        ("Geometry teacher", 0.80),
        ("Reusable geometry cache", 0.68),
        ("Pseudo-geometry\nfinetune wrapper", 0.56),
        ("Unified visual compare", 0.44),
        ("Ghost measurement", 0.32),
        ("Probe automation", 0.20),
    ]
    for text, y in right_boxes:
        draw_arch_box(ax, 0.57, y, 0.31, 0.08, text, fc="#eef5ea", ec="#5b8742", lw=2.3, fontsize=15)

    draw_arch_box(ax, 0.43, 0.47, 0.10, 0.12, "Original\nVGGT\ncore", fc="#f3efe6", ec="#9a8d74", lw=2.5, fontsize=18)
    for _, y in left_boxes:
        draw_arch_arrow(ax, 0.38, y + 0.045, 0.43, 0.53)
    for _, y in right_boxes:
        draw_arch_arrow(ax, 0.53, 0.53, 0.57, y + 0.04)

    ax.text(0.50, 0.10, "Message: the backbone is mostly preserved; most work is added around it.", ha="center", fontsize=14, color="#555555")
    return save_fig(fig, cache_png("code_touch_map"))


def generate_teacher_npz_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.6))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.06, 0.60, 0.16, 0.16, "Original VGGT\nDepth / Conf /\nCamera", fc="#fff1e8", ec="#c96310", lw=2.5, fontsize=18)
    draw_arch_box(ax, 0.31, 0.58, 0.18, 0.20, "Geometry\nTeacher", fc="#f4efe4", ec="#9a8d74", lw=2.5, fontsize=20)
    draw_arch_box(ax, 0.58, 0.58, 0.16, 0.20, "Unproject /\nframe meta /\npoint map", fc="#eef5ea", ec="#5b8742", lw=2.5, fontsize=18)
    draw_arch_box(ax, 0.82, 0.54, 0.12, 0.28, "Geometry\ncache\nimages\ndepth\nconf\npointmap\ncamera\nmeta", fc="#eef0fb", ec="#6882b4", lw=2.5, fontsize=16)

    draw_arch_arrow(ax, 0.22, 0.68, 0.31, 0.68)
    draw_arch_arrow(ax, 0.49, 0.68, 0.58, 0.68)
    draw_arch_arrow(ax, 0.74, 0.68, 0.82, 0.68)

    ax.text(0.50, 0.34, "x_cam = (u + o - cx) * z / fx\ny_cam = (v + o - cy) * z / fy\nP_world = R^T * ([x_cam, y_cam, z]^T - t)", ha="center", va="center", fontsize=17, color="#222222")
    ax.text(0.50, 0.12, "Goal: turn original outputs into a stable geometry teacher and a reusable precompute contract.", ha="center", fontsize=14, color="#555555")
    return save_fig(fig, cache_png("teacher_npz_flow"))


def generate_loss_routing_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.8))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.08, 0.41, 0.18, 0.15, "Pseudo-geometry dataset\n(images + geometry cache)", fc="#f4efe4", ec="#9a8d74", lw=2.4, fontsize=16)
    draw_arch_box(ax, 0.36, 0.33, 0.26, 0.30, "Supervision Router\n(train wrapper)", fc="#eef5ea", ec="#5b8742", lw=2.8, fontsize=20)
    ax.text(0.49, 0.40, "confidence weighting\nloss routing\nmask gating", ha="center", va="center", fontsize=13, color="#35522a")
    draw_arch_arrow(ax, 0.26, 0.485, 0.36, 0.485)

    targets = [
        ("Depth", 0.74, 0.73, "#fff1e8", "#c96310"),
        ("Point", 0.86, 0.73, "#eef0fb", "#6882b4"),
        ("Point reproj", 0.74, 0.56, "#eef0fb", "#6882b4"),
        ("Normal", 0.86, 0.56, "#eef5ea", "#5b8742"),
        ("MV-depth", 0.74, 0.39, "#eef5ea", "#5b8742"),
        ("Camera", 0.86, 0.39, "#f7eefb", "#7b61c8"),
        ("FG boost", 0.74, 0.22, "#fff1e8", "#c96310"),
        ("Conf", 0.86, 0.22, "#f4efe4", "#9a8d74"),
    ]
    for text, x, y, fc, ec in targets:
        draw_arch_box(ax, x, y, 0.11, 0.10, text, fc=fc, ec=ec, lw=2.2, fontsize=14)
        draw_arch_arrow(ax, 0.62, 0.48, x, y + 0.05)

    ax.text(0.50, 0.08, "One supervision router collects geometry, camera, reprojection and multiview constraints around the original backbone.", ha="center", fontsize=14, color="#555555")
    return save_fig(fig, cache_png("loss_routing"))


def generate_fg_boost_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 5.6))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.08, 0.32, 0.16, 0.24, "Foreground\nmask", fc="#fff1e8", ec="#c96310", lw=2.5, fontsize=18)
    ax.text(0.16, 0.20, "M_fg", ha="center", fontsize=16, fontweight="bold", color="#c96310")
    draw_arch_arrow(ax, 0.24, 0.44, 0.34, 0.44)

    draw_arch_box(ax, 0.34, 0.32, 0.18, 0.24, "Erode interior", fc="#f4efe4", ec="#9a8d74", lw=2.5, fontsize=18)
    ax.text(0.43, 0.20, "M_interior = erode(M_fg, r)", ha="center", fontsize=14, color="#555555")
    draw_arch_arrow(ax, 0.52, 0.44, 0.63, 0.44)

    draw_arch_box(ax, 0.63, 0.32, 0.18, 0.24, "Foreground\nvalid gating", fc="#eef5ea", ec="#5b8742", lw=2.5, fontsize=18)
    ax.text(0.72, 0.20, "M_valid = M_valid_all * M_fg", ha="center", fontsize=13, color="#35522a")
    draw_arch_arrow(ax, 0.81, 0.44, 0.92, 0.44)

    draw_arch_box(ax, 0.83, 0.32, 0.14, 0.24, "Final\nweight", fc="#eef0fb", ec="#6882b4", lw=2.5, fontsize=20)
    ax.text(0.50, 0.08, "Current code semantics: foreground gating first, then interior boost on top of base weight.", ha="center", fontsize=14, color="#555555")
    return save_fig(fig, cache_png("fg_boost_semantics"))


def generate_compare_probe_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.4))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.06, 0.42, 0.14, 0.14, "Checkpoint", fc="#f4efe4", ec="#9a8d74", lw=2.4, fontsize=17)
    draw_arch_box(ax, 0.28, 0.42, 0.16, 0.14, "Finetune\nstep", fc="#eef5ea", ec="#5b8742", lw=2.4, fontsize=18)
    draw_arch_box(ax, 0.52, 0.42, 0.18, 0.14, "Unified visual\ncompare", fc="#eef0fb", ec="#6882b4", lw=2.4, fontsize=18)
    draw_arch_box(ax, 0.78, 0.54, 0.15, 0.11, "Ghost\nmeasurement", fc="#fff1e8", ec="#c96310", lw=2.3, fontsize=16)
    draw_arch_box(ax, 0.78, 0.37, 0.15, 0.11, "Point-support\nmeasurement", fc="#eef5ea", ec="#5b8742", lw=2.3, fontsize=15)
    draw_arch_box(ax, 0.78, 0.20, 0.15, 0.11, "Probe\nautomation", fc="#f7eefb", ec="#7b61c8", lw=2.3, fontsize=16)

    draw_arch_arrow(ax, 0.20, 0.49, 0.28, 0.49)
    draw_arch_arrow(ax, 0.44, 0.49, 0.52, 0.49)
    draw_arch_arrow(ax, 0.70, 0.49, 0.78, 0.595)
    draw_arch_arrow(ax, 0.70, 0.49, 0.78, 0.425)
    draw_arch_arrow(ax, 0.70, 0.49, 0.78, 0.255)

    ax.text(0.855, 0.10, "one-step\nstepcurve\nmask-boost", ha="center", fontsize=13, color="#5c4a87")
    ax.text(0.50, 0.08, "Probe scripts are an evaluation shell around the same original-VGGT line, not a new model family.", ha="center", fontsize=14, color="#555555")
    return save_fig(fig, cache_png("compare_probe_framework"))


def generate_visual_assets(data: dict) -> dict[str, Path]:
    assets: dict[str, Path] = {}
    view = data["viewcount_map"]
    one = data["one_step_map"]
    stepcurve = data["stepcurve_map"]
    mask_boost = data["mask_boost_map"]

    assets["current_architecture"] = generate_current_architecture_diagram()
    assets["code_touch_map"] = generate_code_touch_diagram()
    assets["teacher_npz_flow"] = generate_teacher_npz_diagram()
    assets["loss_routing"] = generate_loss_routing_diagram()
    assets["fg_boost_semantics"] = generate_fg_boost_diagram()
    assets["compare_probe_framework"] = generate_compare_probe_diagram()

    fig, ax = plt.subplots(figsize=(13, 7))
    ax.set_axis_off()
    nodes = {
        "原版VGGT本地主线": (0.50, 0.92, "#c75819"),
        "Geometry teacher": (0.18, 0.72, "#d9c7a7"),
        "Precompute NPZ": (0.18, 0.46, "#d9c7a7"),
        "Finetune wrapper": (0.50, 0.72, "#d9c7a7"),
        "Compare & Metrics": (0.82, 0.72, "#d9c7a7"),
        "One-step probe": (0.35, 0.22, "#e8dcc5"),
        "Stepcurve probe": (0.50, 0.22, "#e8dcc5"),
        "Mask-boost probe": (0.65, 0.22, "#e8dcc5"),
        "Ghost": (0.74, 0.48, "#efe7d7"),
        "Point-support": (0.90, 0.48, "#efe7d7"),
        "Masked PSNR/L1": (0.90, 0.28, "#efe7d7"),
    }
    edges = [
        ("原版VGGT本地主线", "Geometry teacher"),
        ("原版VGGT本地主线", "Finetune wrapper"),
        ("原版VGGT本地主线", "Compare & Metrics"),
        ("Geometry teacher", "Precompute NPZ"),
        ("Precompute NPZ", "Finetune wrapper"),
        ("Finetune wrapper", "One-step probe"),
        ("Finetune wrapper", "Stepcurve probe"),
        ("Finetune wrapper", "Mask-boost probe"),
        ("Compare & Metrics", "Ghost"),
        ("Compare & Metrics", "Point-support"),
        ("Compare & Metrics", "Masked PSNR/L1"),
    ]
    for src, dst in edges:
        x1, y1, _ = nodes[src]
        x2, y2, _ = nodes[dst]
        ax.plot([x1, x2], [y1, y2], color="#7a6a58", lw=2, alpha=0.8)
    for label, (x, y, color) in nodes.items():
        ax.text(
            x,
            y,
            label,
            ha="center",
            va="center",
            fontsize=16 if label == "原版VGGT本地主线" else 13,
            fontweight="bold" if label == "原版VGGT本地主线" else "normal",
            bbox=dict(boxstyle="round,pad=0.45", fc=color, ec="#7a6a58", lw=1.5),
        )
    ax.text(0.5, 0.05, "思维导图: teacher / loss / compare / automation 四层围绕原版 VGGT 组织", ha="center", fontsize=14, color="#233")
    assets["mindmap"] = save_fig(fig, cache_png("mindmap"))

    fig, ax = plt.subplots(figsize=(13, 3.8))
    ax.set_axis_off()
    stages = [
        ("原版 VGGT", "depth/conf/camera"),
        ("Geometry teacher", "pointmap + frame/meta"),
        ("Geometry cache", "reusable supervision"),
        ("Supervision router", "loss routing"),
        ("Visual compare", "pred/tgt/support/ghost"),
        ("Probe summaries", "one-step / stepcurve / mask-boost"),
    ]
    xs = np.linspace(0.08, 0.92, len(stages))
    for i, ((title, subtitle), x) in enumerate(zip(stages, xs)):
        ax.text(x, 0.56, title, ha="center", va="center", fontsize=15, fontweight="bold",
                bbox=dict(boxstyle="round,pad=0.45", fc="#e8dcc5", ec="#7a6a58", lw=1.3))
        ax.text(x, 0.38, subtitle, ha="center", va="center", fontsize=11, color="#344")
        if i < len(stages) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.06, 0.56), xytext=(x + 0.06, 0.56),
                        arrowprops=dict(arrowstyle="->", lw=2, color="#c75819"))
    ax.text(0.5, 0.14, "流程图: 原模型输出先被固化为 teacher supervision，再进入受控训练、统一 compare 与自动化汇总", ha="center", fontsize=13)
    assets["pipeline"] = save_fig(fig, cache_png("pipeline"))

    labels = ["6src_hist", "12src_nested", "23cam_fullset"]
    coverage = [view[k]["coverage_ratio"] for k in labels]
    ghost = [view[k]["ghost_visual_score"] for k in labels]
    psnr = [view[k]["native_psnr"] for k in labels]
    fig, axes = plt.subplots(1, 3, figsize=(14, 4.2))
    for ax, vals, ttl, color in zip(
        axes,
        [coverage, ghost, psnr],
        ["Coverage", "Ghost Visual Score", "Native PSNR"],
        ["#648fff", "#ff7c43", "#2f9e44"],
    ):
        bars = ax.bar(labels, vals, color=color, alpha=0.85)
        ax.set_title(ttl, fontsize=14, fontweight="bold")
        ax.tick_params(axis="x", rotation=15)
        for b, v in zip(bars, vals):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height(), f"{v:.3f}", ha="center", va="bottom", fontsize=10)
    fig.suptitle("Baseline view-count comparison", fontsize=18, fontweight="bold")
    assets["baseline_bars"] = save_fig(fig, cache_png("baseline_bars"))

    fig, axes = plt.subplots(1, 2, figsize=(13, 4.6))
    profiles = ["6src_hist", "12src_nested", "23cam_fullset"]
    pre_g = [one[p]["pre_ghost_visual_score"] for p in profiles]
    post_g = [one[p]["post_ghost_visual_score"] for p in profiles]
    pre_p = [one[p]["pre_native_psnr"] for p in profiles]
    post_p = [one[p]["post_native_psnr"] for p in profiles]
    x = np.arange(len(profiles))
    w = 0.34
    axes[0].bar(x - w / 2, pre_g, width=w, color="#c9b07f", label="pre")
    axes[0].bar(x + w / 2, post_g, width=w, color="#c75819", label="post")
    axes[0].set_xticks(x, profiles, rotation=15)
    axes[0].set_title("One-step ghost", fontsize=15, fontweight="bold")
    axes[0].legend()
    axes[1].bar(x - w / 2, pre_p, width=w, color="#8db3e2", label="pre")
    axes[1].bar(x + w / 2, post_p, width=w, color="#2f9e44", label="post")
    axes[1].set_xticks(x, profiles, rotation=15)
    axes[1].set_title("One-step native PSNR", fontsize=15, fontweight="bold")
    axes[1].legend()
    fig.suptitle("One-step probe: pre vs post", fontsize=18, fontweight="bold")
    assets["one_step_pairs"] = save_fig(fig, cache_png("one_step_pairs"))

    for profile in ["12src_nested", "6src_hist"]:
        rows = stepcurve[profile]["steps"]
        fig, axes = plt.subplots(2, 2, figsize=(13, 8))
        for ax, key, title, color in [
            (axes[0, 0], "ghost_visual_score", "Ghost Visual Score", "#c75819"),
            (axes[0, 1], "subject_psnr", "Subject PSNR", "#2f9e44"),
            (axes[1, 0], "largest_component_share", "Largest Component Share", "#648fff"),
            (axes[1, 1], "secondary_component_mass", "Secondary Component Mass", "#8f5ae8"),
        ]:
            xs, ys = valid_points(rows, key)
            ax.plot(xs, ys, marker="o", lw=2.5, color=color)
            ax.set_title(title, fontsize=14, fontweight="bold")
            ax.set_xlabel("Step horizon")
            ax.grid(alpha=0.25)
        fig.suptitle(f"Stepcurve metrics: {profile}", fontsize=19, fontweight="bold")
        assets[f"stepcurve_{profile}"] = save_fig(fig, cache_png(f"stepcurve_{profile}"))

    for profile in ["12src_nested", "6src_hist"]:
        boost_profile = mask_boost[profile]
        winner = int(boost_profile["winner_alpha"])
        native_rows = {int(r["step"]): r for r in stepcurve[profile]["steps"]}
        winner_rows = [r for r in boost_profile["table_rows"] if int(r["alpha"]) == winner and r.get("available")]
        fig, axes = plt.subplots(1, 2, figsize=(13, 4.8))
        xs = [int(r["step"]) for r in winner_rows]
        subj = [float(r["subject_psnr"]) for r in winner_rows]
        ghost_vals = [float(r["ghost_visual_score"]) for r in winner_rows]
        native_xs_subj = [s for s in xs if s in native_rows and not math.isnan(float(native_rows[s]["subject_psnr"]))]
        native_subj = [float(native_rows[s]["subject_psnr"]) for s in native_xs_subj]
        native_xs_ghost = [s for s in xs if s in native_rows and not math.isnan(float(native_rows[s]["ghost_visual_score"]))]
        native_ghost = [float(native_rows[s]["ghost_visual_score"]) for s in native_xs_ghost]
        axes[0].plot(native_xs_subj, native_subj, marker="o", lw=2.3, color="#2f9e44", label="native")
        axes[0].plot(xs, subj, marker="o", lw=2.3, color="#c75819", label=f"alpha={winner}")
        axes[0].set_title("Subject PSNR", fontsize=15, fontweight="bold")
        axes[0].grid(alpha=0.25)
        axes[0].legend()
        axes[1].plot(native_xs_ghost, native_ghost, marker="o", lw=2.3, color="#648fff", label="native")
        axes[1].plot(xs, ghost_vals, marker="o", lw=2.3, color="#c75819", label=f"alpha={winner}")
        axes[1].set_title("Ghost Visual Score", fontsize=15, fontweight="bold")
        axes[1].grid(alpha=0.25)
        axes[1].legend()
        fig.suptitle(f"Mask-boost winner vs native: {profile}", fontsize=18, fontweight="bold")
        assets[f"mask_boost_{profile}"] = save_fig(fig, cache_png(f"mask_boost_{profile}"))

    return assets


def file_groups() -> list[tuple[str, list[str]]]:
    return [
        (
            "Geometry / teacher / precompute",
            [
                "vggt/utils/load_fn.py",
                "vggt/utils/geometry.py",
                "vggt_geom.py",
                "precompute_zju_vggt_geom.py",
                "modal_run_train.py",
            ],
        ),
        (
            "Train wrapper / objective routing",
            [
                "finetune_vggt_pseudo.py",
            ],
        ),
        (
            "Render / scoring / metrics",
            [
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
            ],
        ),
        (
            "Automation / experiments",
            [
                "scripts/orig_vggt_one_step_probe/*",
                "scripts/orig_vggt_stepcurve_probe/*",
                "scripts/orig_vggt_mask_boost_probe/*",
            ],
        ),
    ]


def build_slide_specs(data: dict) -> list[SlideSpec]:
    view = data["viewcount_map"]
    one = data["one_step_map"]
    stepcurve = data["stepcurve_map"]
    mask_boost = data["mask_boost_map"]
    assets = data.get("assets", {})
    step_12 = stepcurve["12src_nested"]
    step_6 = stepcurve["6src_hist"]
    mb_12 = mask_boost["12src_nested"]
    mb_6 = mask_boost["6src_hist"]

    p12_s1 = step_row(step_12, 1)
    p12_s4 = step_row(step_12, 4)
    p12_s16 = step_row(step_12, 16)
    p6_s1 = step_row(step_6, 1)
    p6_s16 = step_row(step_6, 16)

    d12_s1 = delta_vs_native_subject_psnr(data, "12src_nested", 1)
    d12_s2 = delta_vs_native_subject_psnr(data, "12src_nested", 2)
    d12_s4 = delta_vs_native_subject_psnr(data, "12src_nested", 4)
    d12_s16 = delta_vs_native_subject_psnr(data, "12src_nested", 16)
    d6_s1 = delta_vs_native_subject_psnr(data, "6src_hist", 1)
    d6_s2 = delta_vs_native_subject_psnr(data, "6src_hist", 2)
    d6_s16 = delta_vs_native_subject_psnr(data, "6src_hist", 16)

    groups = file_groups()
    group_lines = [
        f"{title}: " + ", ".join(paths)
        for title, paths in groups
    ]

    slides: list[SlideSpec] = [
        SlideSpec(
            title="基于当前本地代码的原版 VGGT 改动讲解",
            bullets=[
                "范围: 只覆盖 original VGGT + model.pt 这条本地实验线，以及围绕它的 teacher / loss / compare / probe 自动化改动。",
                "不覆盖: AnySplat / Gaussian 新 family、PixelSplat 外部实现、view-decoder ablation 等非原版 VGGT 主线。",
                "目标: 让导师能按“改了什么、为什么改、公式是什么、代码在哪、实验结论如何”逐层理解整条线。",
            ],
            speaker_notes=(
                "开场先把边界讲清楚。这份 PPT 不是介绍所有 3D 重建路线，而是只讲当前本地仓库里围绕原版 VGGT 做过的所有有效改动。"
                "重点是告诉导师，我们没有替换 VGGT 主干，而是在 teacher 构造、伪几何预计算、训练损失路由、可视化比较和自动化 probing 上逐层加了壳层。"
            ),
        ),
        SlideSpec(
            title="讲解范围与主问题",
            bullets=[
                "主问题 1: 原版 VGGT 直接出点云和继续短程训练后，单主体是否更稳、更集中、重影是否更少。",
                "主问题 2: 如果只对人体 mask 内部加强监督，能否让主体区域更准，尤其是 mask 内 PSNR / L1 更好。",
                "主问题 3: 当前本地代码里各项改动彼此如何衔接，哪些是训练时改动，哪些只是分析与自动化。",
            ],
            formulas=[
                "总链路 = 原版 VGGT -> 几何 teacher / precompute -> finetune wrapper -> raw compare -> ghost / point-support -> one-step / stepcurve / mask-boost probe",
            ],
            speaker_notes=(
                "这一页把导师最关心的三层问题一次说清。第一层是模型行为，第二层是 mask 内增强是否对主体重建真有帮助，第三层是代码架构到底加了哪些部件。"
            ),
        ),
        SlideSpec(
            title="改动地图总览",
            bullets=group_lines,
            speaker_notes=(
                "这里建议按四大层来讲。第一层是 teacher 和预计算，第二层是训练损失和监督路由，第三层是 raw compare、ghost 和 point-support 后处理，第四层是 one-step、stepcurve、mask-boost 三条自动化实验线。"
            ),
        ),
        SlideSpec(
            title="文件索引 1/2",
            bullets=[
                "vggt/utils/load_fn.py: 负责多视图图像与几何输入加载，加速 threaded image loading。",
                "vggt/utils/geometry.py: 统一 unproject_depth_map_to_point_map，并支持不同 unproject 实现。",
                "vggt_geom.py: 定义 VGGTGeomTeacher，负责把原版 VGGT 输出整理成 depth / conf / pointmap / camera supervision。",
                "precompute_zju_vggt_geom.py: 把 teacher 输出固化为 NPZ，供后续 finetune 数据集直接读。",
                "modal_run_train.py: 把 precompute 和 train 运行参数打包到 Modal 任务里。",
            ],
            speaker_notes="这页先讲 teacher 到预计算链路，强调这些改动是给原版 VGGT 外面加数据壳，不是改 backbone 结构。",
        ),
        SlideSpec(
            title="文件索引 2/2",
            bullets=[
                "finetune_vggt_pseudo.py: 本地原版 VGGT 主训练包装器，包含 PseudoGeomDataset、loss routing、mask boost、MV-depth 等改动。",
                "scripts/orig_vggt_viewcount/render_raw_compare.py: 生成 pred / tgt / point-support / ghost triplet 等统一对比图。",
                "tools/score_ghosting_from_cat_pred.py: 从 ghost triplet 里计算 ghost_visual_score、peak_count 等启发式指标。",
                "scripts/orig_vggt_one_step_probe/*: 一步微调实验自动化。",
                "scripts/orig_vggt_stepcurve_probe/*: step0000/1/2/4/8/16 曲线实验、prefix 审计、point-support 后处理、汇总。",
                "scripts/orig_vggt_mask_boost_probe/*: mask 内增强实验的短程筛选与 winner 延长。",
            ],
            speaker_notes="这页进入训练包装器、可视化和三套自动化探针。导师如果只问‘做了哪些实验脚手架’，这一页基本能覆盖。",
        ),
    ]
    slides.extend(
        [
            SlideSpec(
                title="几何 Teacher: 从深度反投影到世界点",
                bullets=[
                    "VGGTGeomTeacher 把原版 VGGT 的 depth / conf / camera 输出整理成可监督的 pointmap、相机位姿和元数据。",
                    "核心操作是把每个像素深度反投影到相机坐标，再用外参变换到世界坐标。",
                    "本地代码支持不同 unproject 实现，并记录 unproject_impl、pointmap_source、pointmap_frame 进入 NPZ 元数据。",
                    "point_head_frame=auto 时，会根据自重投影误差选择更稳的一套 point-head 参考系。",
                ],
                formulas=[
                    "x_cam = (u - cx) * z / fx",
                    "y_cam = (v - cy) * z / fy",
                    "P_world = R^T * ([x_cam, y_cam, z]^T - t)",
                    "pixel offset o in {0, 0.5} depends on unproject_impl",
                ],
                code_refs=[
                    "vggt_geom.py: VGGTGeomTeacher / _unproject_depth_to_world_batched / forward_prepared_batch",
                    "vggt/utils/geometry.py: _resolve_unproject_impl / unproject_depth_map_to_point_map",
                ],
                speaker_notes=(
                    "这一页建议直接讲公式。核心不是新网络，而是把原版 VGGT 的深度输出可靠地变成世界点监督。"
                    "其中一个容易被忽略的细节是像素中心偏移，也就是 pixel-center offset。不同 unproject 约定会导致点图整体有亚像素级别偏差，所以代码里显式保留 unproject_impl。"
                ),
            ),
            SlideSpec(
                title="预计算链路: 把 teacher 输出固化成 NPZ",
                bullets=[
                    "precompute_zju_vggt_geom.py 先跑 teacher，再把 depth、depth_conf、pointmap、相机内外参和 pointmap 元信息写入每帧 NPZ。",
                    "这样 finetune 阶段不用每次重新调用原版 VGGT 前向，训练读的是固定几何 teacher 结果。",
                    "load_fn.py 加 threaded image loading，减少多视图图像 I/O 开销。",
                    "modal_run_train.py 把 pointmap_source、precompute_unproject_impl、MV 支持参数一起写到远端运行配置。",
                ],
                formulas=[
                    "saved NPZ = {images, depth, depth_conf, pointmap, extrinsic, intrinsic, pointmap_source, pointmap_frame, unproject_impl}",
                ],
                code_refs=[
                    "precompute_zju_vggt_geom.py: env resolution / teacher construction / NPZ save fields",
                    "vggt/utils/load_fn.py: threaded image loading",
                    "modal_run_train.py: precompute + train config routing",
                ],
                speaker_notes=(
                    "这里强调预计算的价值是把 teacher supervision 冻结住。后面 one-step、stepcurve、mask-boost 的所有训练，其实都是在同一个预计算监督场上做控制变量。"
                ),
            ),
            SlideSpec(
                title="训练包装器: 总损失不是单一项，而是可路由组合",
                bullets=[
                    "finetune_vggt_pseudo.py 定义了 PseudoGeomDataset，并把本地所有原版 VGGT 训练改动汇总在同一个 wrapper 里。",
                    "主干思路不是改 backbone，而是在训练阶段叠加 depth、point、reproj、normal、MV-depth、camera 等监督。",
                    "很多开关都是参数化的，因此可以做严格控制变量实验。",
                ],
                formulas=[
                    "L_total = λ_depth L_depth + λ_point L_point + λ_point_reproj L_point_reproj + λ_normal L_normal + λ_mvdepth L_mvdepth + λ_mvmask L_mvmask + λ_conf L_conf + λ_fgpresence L_fgpresence + λ_geom L_geom + λ_cam L_cam + ...",
                ],
                code_refs=[
                    "finetune_vggt_pseudo.py: PseudoGeomDataset / total loss accumulation",
                ],
                speaker_notes=(
                    "导师如果问‘你们到底改了原版 VGGT 什么’，这页要明确回答：主要是训练期 supervision routing 和 loss engineering，而不是把原模型结构推翻重来。"
                ),
            ),
            SlideSpec(
                title="置信度加权: 先决定哪些像素值得信",
                bullets=[
                    "深度和点监督不会对所有像素一视同仁，而是先根据 teacher 置信度构造 base weight。",
                    "支持按连续权重衰减，也支持 per-view quantile mask，只保留每个视角最可信的一部分像素。",
                    "这样能减少伪几何噪声把训练推偏。",
                ],
                formulas=[
                    "w_conf(x) = valid(x) * clip((c(x) - t) / (1 - t), 0, 1)^γ",
                    "quantile mask: keep top-q confident valid pixels per view",
                ],
                code_refs=[
                    "finetune_vggt_pseudo.py: _build_conf_weight / _build_per_view_conf_quantile_mask",
                ],
                speaker_notes=(
                    "这一步很关键，因为后面所有 boost、MV-depth、normal consistency 都是建立在 base weight 上再往外长的。"
                ),
            ),
            SlideSpec(
                title="前景 Mask 监督与 interior boost",
                bullets=[
                    "本地代码把人体 mask 接进训练路由，可以做 foreground gating，也可以进一步只增强前景内部区域。",
                    "当前 mask-boost probe 的真实语义不是“全图 loss + 轻微前景增强”，而是更强的“FG-only supervision + interior boost”。",
                    "原因是实验里 use_fg_mask=on，且 fg_supervision_bg_floor=0，所以背景通道实际被关掉了。",
                ],
                formulas=[
                    "M_valid = M_valid_all * (M_fg + b * (1 - M_fg))",
                    "M_interior = erode(M_fg, r)",
                    "w_final = w_base * [1 + (β - 1) * M_interior]",
                ],
                code_refs=[
                    "finetune_vggt_pseudo.py: _apply_fg_supervision_boost / _build_fg_supervision_boost_mask / _erode_mask_tensor / train-loop routing",
                ],
                speaker_notes=(
                    "这一页一定要替导师拆解语义偏差。导师原话更像是‘全图 loss 保留，但人体内部再加权’。"
                    "而当前 probe 实现得更强，已经接近‘只在前景里监督，再把内部拉高’。这也是为什么 masked PSNR 没有系统性变好时，我们不能简单说导师思路错了，而要说实现语义更激进。"
                ),
            ),
            SlideSpec(
                title="点目标路由: depth_unproject / pointmap / blend / consensus",
                bullets=[
                    "点监督并不只有一种目标，代码支持直接用 teacher pointmap，也支持把 depth 反投影结果和伪点图做混合。",
                    "blend 模式会根据两者一致性自适应调权，不一致处更依赖更可靠的一端。",
                    "depth_consensus_unproject 进一步强调跨视图深度一致的区域。",
                ],
                formulas=[
                    "r(x) = exp(-||U(D_tgt)(x) - P_pseudo(x)|| / τ)",
                    "P_tgt(x) = α(x) U(D_tgt)(x) + (1 - α(x)) P_pseudo(x)",
                ],
                code_refs=[
                    "finetune_vggt_pseudo.py: --point_target_mode / point target routing",
                ],
                speaker_notes=(
                    "这里的核心思想是，不把 teacher pointmap 当绝对真值，而是允许 depth 反投影和点图之间做互相校正。"
                ),
            ),
            SlideSpec(
                title="MV support 与 MV-depth: 用跨视图一致性压异常层",
                bullets=[
                    "多视图支持权重先看一个点在多少其他视角里还能对得上，再把这个 agree ratio 映射成 support weight。",
                    "MV-depth reprojection loss 会把 i 视角的点投到 j 视角深度图上比较，逼前后两层不要长期共存。",
                    "代码还支持 outlier emphasis，对明显不一致的区域额外加权，直接朝重影区域施压。",
                ],
                formulas=[
                    "agree_ratio(x) = #agree_views(x) / (V - 1)",
                    "s_mv(x) = floor + (1 - floor) * agree_ratio(x)",
                    "L_mvdepth = Avg_{i->j} [ Σ_x w_ij(x) ρ(|z_ij(x) - D_j(π_ij(x))|) / Σ_x w_ij(x) ]",
                    "w_ij <- w_ij * clip(1 + β_out * (1 - inlier), 1, cap)",
                ],
                code_refs=[
                    "finetune_vggt_pseudo.py: _map_support_weight / _point_multiview_support_weight / _point_multiview_depth_reproj_loss",
                ],
                speaker_notes=(
                    "如果导师问‘为什么你们总盯着 ghost、双层、厚壳’，这页就是数学化回答。MV-depth 的意义就是把前后两层不一致视为可优化的 reprojection error。"
                ),
            ),
            SlideSpec(
                title="辅助几何项: normal / geom / camera",
                bullets=[
                    "normal consistency 约束表面法向不要乱跳，能减轻局部壳层噪声。",
                    "geom consistency 让预测深度反投影和预测点图互相一致。",
                    "camera pose losses 把平移、旋转和视场约束都接进来，避免几何误差通过相机自由度泄漏。",
                ],
                formulas=[
                    "L_normal = Σ w(x) (1 - <n_pred, n_tgt>) / Σ w(x)",
                    "L_geom = Σ w_base ||U(D_pred) - P_pred|| / Σ w_base",
                    "L_cam = L_trans + w_rot (1 - |q_pred · q_tgt|) + w_fov L_fov",
                ],
                code_refs=[
                    "finetune_vggt_pseudo.py: _point_normal_consistency_loss / _camera_pose_losses / geometry-camera branch in train loop",
                ],
                speaker_notes="这一页可以讲成‘把局部表面、全局几何和相机都拴住’，否则单纯优化某一个重投影指标很容易出现投机解。",
            ),
            SlideSpec(
                title="raw compare 图谱: 哪张图在看什么",
                bullets=[
                    "weight_native.png / cat_weight_pred_tgt*.png 是 point-support 图，关注形状、集中度、连通性、多峰风险。",
                    "cat_fg_mask_pred_tgt_step000000.png 是 ghost triplet，左栏是 mask，中栏是预测，右栏是 target，用来做 ghost 评分。",
                    "训练 debug 里的 triplet_conf_tgt_pred_weight 右栏是 supervision weight，不是点云支持图，不能混读。",
                    "weight_native.png 是归一化后的相对支持图，不是绝对物理支持强度。",
                ],
                formulas=[
                    "weight01 = clip(weight_map / p99(weight_map[hit]), 0, 1)",
                ],
                code_refs=[
                    "scripts/orig_vggt_viewcount/render_raw_compare.py: _normalize_mask01 / _weight_to_rgb / weight normalization",
                ],
                image_paths=[POINT_TAXONOMY_EXAMPLE, GHOST_TAXONOMY_EXAMPLE],
                image_caption="左: point-support triplet; 右: ghost triplet。两类图不能混为一谈。",
                layout="two_images",
                speaker_notes=(
                    "这页一定要帮导师和组里把图谱 taxonomy 钉死。你们之前反复提到的误会，很多都来自把 support 图、ghost triplet 和训练 debug weight 图混读。"
                ),
            ),
        ]
    )
    slides.extend(
        [
            SlideSpec(
                title="Ghost 指标: 从 triplet 图像里量化重影",
                bullets=[
                    "ghost 不是训练 loss，而是后处理视觉启发式分数，输入固定是 cat_fg_mask_pred_tgt_step000000.png。",
                    "算法先从左栏得到 mask，再从中栏预测图里阈值出活跃区域，测宽度、面积、中心偏移和横向多峰。",
                    "最终 ghost_visual_score 越低越好；它同时惩罚过宽、过厚、多峰、偏移、太暗和塌缩。",
                ],
                formulas=[
                    "width_ratio = W_pred / W_mask",
                    "area_ratio = A_pred / A_mask",
                    "ghost_score = 1.0*max(0, width_ratio - 1.10) + 0.4*max(0, area_ratio - 1.30) + 0.6*max(0, peak_count - 1) + 0.5*max(0, center_offset_ratio - 0.22)",
                    "ghost_visual_score = ghost_score + dark_penalty + collapse_penalty",
                ],
                code_refs=[
                    "tools/score_ghosting_from_cat_pred.py: _count_peaks / _score_one / final penalties",
                ],
                speaker_notes=(
                    "ghost 的价值是把‘看起来有双层、太宽、太黑、太塌’这些感性描述转成可比数字。"
                    "但也要明确它是 2D 视觉启发式，不是 3D 真值指标，所以它在本地实验里始终是第二判据，不是唯一主判据。"
                ),
            ),
            SlideSpec(
                title="point-support 与 mask 内重建指标",
                bullets=[
                    "measure_point_support.py 固定从 weight_native.png 的绿色通道读取 support_map01，再结合 GT mask 计算主体性指标。",
                    "主判据看 subject_support_share、outside_subject_support_share、largest_component_share、secondary_component_mass、support_peak_count。",
                    "同时补充 mask 内重建指标: subject_psnr、subject_l1；优先使用 report.json 里的 GT mask 路径，必要时回退。",
                    "背景灰斑变多不一定代表绝对支持变强，常见原因是归一化把弱背景响应相对抬亮了。",
                ],
                formulas=[
                    "subject_support_share = Σ S M / Σ S",
                    "outside_subject_support_share = Σ S (1 - M) / Σ S",
                    "largest_component_share = max_c Σ_{x in c} S / Σ_{x in active} S",
                    "subject_psnr = -10 log10( Σ ((pred - tgt)^2 M) / (Σ M * C) )",
                    "subject_l1 = Σ |pred - tgt| M / (Σ M * C)",
                ],
                code_refs=[
                    "scripts/orig_vggt_stepcurve_probe/measure_point_support.py: _load_subject_mask / _masked_error_metrics / support_peak_count",
                ],
                speaker_notes=(
                    "这一页把导师最近追问的两个点都放一起。第一，mask 内 PSNR 用的是 GT mask 优先。第二，point-support 图是相对归一化形状指标，所以背景灰块不能直接解释成绝对能量飙升。"
                ),
            ),
            SlideSpec(
                title="one-step probe: 验证‘只训一步’是否有短程收益",
                bullets=[
                    f"6src_hist: ghost {format_float(one['6src_hist']['pre_ghost_visual_score'])} -> {format_float(one['6src_hist']['post_ghost_visual_score'])}, native PSNR {format_float(one['6src_hist']['pre_native_psnr'])} -> {format_float(one['6src_hist']['post_native_psnr'])}。",
                    f"12src_nested: ghost {format_float(one['12src_nested']['pre_ghost_visual_score'])} -> {format_float(one['12src_nested']['post_ghost_visual_score'])}, native PSNR {format_float(one['12src_nested']['pre_native_psnr'])} -> {format_float(one['12src_nested']['post_native_psnr'])}。",
                    f"23cam_fullset: ghost {format_float(one['23cam_fullset']['pre_ghost_visual_score'])} -> {format_float(one['23cam_fullset']['post_ghost_visual_score'])}, 基本无收益。",
                    "结论: 一步微调对 6src 和 12src 的 ghost 有短程帮助，但并不自动等价于画质或主体重建持续改善。",
                ],
                code_refs=[
                    "scripts/orig_vggt_one_step_probe/common.py / run_task.py",
                    "logs/modal_phase5/reports/orig_vggt_one_step_probe_summary_latest.json",
                ],
                speaker_notes="这一页用来说明为什么后来还要做 stepcurve。因为 one-step 只能告诉我们短程有没有方向性变化，不能回答会不会越训越好。",
            ),
            SlideSpec(
                title="stepcurve probe: 把‘第一次回弹在哪’自动化",
                bullets=[
                    "固定 horizon: step0000 / 0001 / 0002 / 0004 / 0008 / 0016；每个 horizon 都独立从 model.pt 出发，保证 prefix 可比。",
                    "audit_prefix.py 用显式 torch.Generator 复现 DataLoader 采样顺序，验证短 horizon 是长 horizon 的前缀。",
                    "profile 级冻结 support_threshold_abs = max(0.15, q75(step0000 非零 support))，避免阈值漂移污染曲线。",
                    f"12src_nested 结论: FIRST_REBOUND_STEP={step_12['FIRST_REBOUND_STEP']}, EXTEND_DECISION={step_12['EXTEND_DECISION']}, prefix_ok={step_12['prefix_consistent_through_step0016']}.",
                    f"6src_hist 结论: FIRST_REBOUND_STEP={step_6['FIRST_REBOUND_STEP']}, EXTEND_DECISION={step_6['EXTEND_DECISION']}, prefix_ok={step_6['prefix_consistent_through_step0016']}.",
                ],
                formulas=[
                    "first_rebound when ghost worsens, or at least two point-support main metrics worsen, or peak_count rises while largest_component_share does not improve",
                ],
                code_refs=[
                    "scripts/orig_vggt_stepcurve_probe/audit_prefix.py / run_task.py / summarize_runs.py",
                ],
                speaker_notes=(
                    "这一页要强调 task 化和 prefix 审计的意义。stepcurve 的可信度，不只是因为跑了 1/2/4/8/16，而是因为每个 horizon 都被证明真的是同一采样轨迹的前缀。"
                ),
            ),
            SlideSpec(
                title="stepcurve 结果解读: 短程有改善，但不支持继续深训",
                bullets=[
                    f"12src_nested: step0001 时 ghost {format_float(p12_s1['ghost_visual_score'])}，比 step0000 的 {format_float(step_row(step_12, 0)['ghost_visual_score'])} 更低；largest_component_share 从 {format_float(step_row(step_12, 0)['largest_component_share'])} 升到 {format_float(p12_s1['largest_component_share'])}，但 subject_psnr 从 {format_float(step_row(step_12, 0)['subject_psnr'])} 降到 {format_float(p12_s1['subject_psnr'])}。",
                    f"12src_nested: step0004 ghost 进一步到 {format_float(p12_s4['ghost_visual_score'])}，但 support_peak_count 回到 {p12_s4['support_peak_count']}，step0008 首次失败，step0016 largest_component_share 掉到 {format_float(p12_s16['largest_component_share'])}。",
                    f"6src_hist: step0001 ghost {format_float(p6_s1['ghost_visual_score'])} 低于 step0000 的 {format_float(step_row(step_6, 0)['ghost_visual_score'])}，但 first_rebound 仍被判为 {step_6['FIRST_REBOUND_STEP']}，因为 point-support 主指标没有形成持续单调改善。",
                    f"6src_hist: step0016 subject_psnr={format_float(p6_s16['subject_psnr'])}, largest_component_share={format_float(p6_s16['largest_component_share'])}, secondary_component_mass={format_float(p6_s16['secondary_component_mass'])}，仍不足以支持扩到 32。",
                ],
                code_refs=[
                    "logs/modal_phase5/reports/orig_vggt_stepcurve_probe_summary_latest.json",
                ],
                speaker_notes=(
                    "这页是 stepcurve 的核心结论页。口径要很稳：可以承认短程确实有改善，但不能把局部改善误说成会稳定收敛。"
                ),
            ),
            SlideSpec(
                title="mask-boost probe: 当前实现到底试了什么",
                bullets=[
                    "实验线不是新 family，而是基于现有 finetune_vggt_pseudo.py 加开关。",
                    "alpha 到 boost 的映射是 alpha=1/2/4 -> fg_supervision_boost=2/3/5。",
                    "短程先筛 1/2/4/8，再只延长 winner 到 12/16/24。",
                    f"12src_nested winner_alpha={mb_12['winner_alpha']}, winner_fg_supervision_boost={format_float(mb_12['winner_fg_supervision_boost'], 1)}。",
                    f"6src_hist winner_alpha={mb_6['winner_alpha']}, winner_fg_supervision_boost={format_float(mb_6['winner_fg_supervision_boost'], 1)}。",
                ],
                formulas=[
                    "alpha_to_fg_boost(alpha) = 1 + alpha",
                ],
                code_refs=[
                    "scripts/orig_vggt_mask_boost_probe/common.py / run_task.py",
                ],
                speaker_notes=(
                    "这页主要解释实验设计：不是无穷扫参，而是先做最小三档 alpha，再延长 winner。这样时间预算和结论强度之间更平衡。"
                ),
            ),
            SlideSpec(
                title="mask-boost 结果: point-support 有局部帮助，但 mask 内 PSNR 没有系统性变好",
                bullets=[
                    f"12src_nested vs native: step0001 delta_subject_psnr={format_float(d12_s1[2])}, step0002={format_float(d12_s2[2])}, step0004={format_float(d12_s4[2])}, step0016={format_float(d12_s16[2])}。",
                    f"6src_hist vs native: step0001 delta_subject_psnr={format_float(d6_s1[2])}, step0002={format_float(d6_s2[2])}, step0016={format_float(d6_s16[2])}。",
                    f"12src_nested 总结标签: POINT_SUPPORT_SINGLE_SUBJECT={mb_12['POINT_SUPPORT_SINGLE_SUBJECT']}, SUBJECT_RECON_IMPROVES={mb_12['SUBJECT_RECON_IMPROVES']}, GHOST_CONTINUES_DOWN={mb_12['GHOST_CONTINUES_DOWN']}。",
                    f"6src_hist 总结标签: POINT_SUPPORT_SINGLE_SUBJECT={mb_6['POINT_SUPPORT_SINGLE_SUBJECT']}, SUBJECT_RECON_IMPROVES={mb_6['SUBJECT_RECON_IMPROVES']}, GHOST_CONTINUES_DOWN={mb_6['GHOST_CONTINUES_DOWN']}。",
                    "解读: 6src 上 point-support 主体性有局部变好，但 mask 内 PSNR 不是稳定正收益；12src 则整体更差。",
                ],
                code_refs=[
                    "logs/modal_phase5/reports/orig_vggt_mask_boost_probe_summary_latest.json",
                    "logs/modal_phase5/reports/orig_vggt_mask_boost_and_ghost_explainer_latest.md",
                ],
                speaker_notes=(
                    "这是导师大概率会追问的一页。重点不是简单说‘mask boost 没用’，而是精确地说：当前实现语义更激进，结果表现为 point-support 局部可改善，但 subject_psnr 没有系统性提升。"
                ),
            ),
            SlideSpec(
                title="结果时间线: 当前本地原版 VGGT 线得到的最稳结论",
                bullets=[
                    f"原版直出 baseline: 6src coverage={format_float(view['6src_hist']['coverage_ratio'])}, ghost={format_float(view['6src_hist']['ghost_visual_score'])}; 12src coverage={format_float(view['12src_nested']['coverage_ratio'])}, ghost={format_float(view['12src_nested']['ghost_visual_score'])}; 23src 覆盖最高但 ghost 无明显更优。",
                    "one-step: 6src 与 12src 都有短程 ghost 改善。",
                    "stepcurve: 两条线都没有证明‘继续深训会稳定越训越好’，因此 stop at 16 是合理结论。",
                    "mask-boost: 当前实现没有把 mask 内 PSNR 系统性拉升，说明语义上还要更贴近导师原意再试一版。",
                ],
                speaker_notes=(
                    "这一页可以当总复盘。先有 baseline，再有 one-step，再有 stepcurve，最后才是 mask-boost。把故事线串起来后，导师就能看到每个实验不是孤立的。"
                ),
            ),
            SlideSpec(
                title="Stepcurve Point-Support 总览图",
                bullets=[
                    "按 step0000/1/2/4/8/16 排列的 point-support 总览图。",
                    "重点观察主体是否更集中、是否出现第二坨、support peak 是否减少、背景灰斑是否只是归一化抬亮。",
                ],
                image_paths=[STEPCURVE_POINT_GRID],
                image_caption="stepcurve point-support grid",
                layout="full_image",
                speaker_notes="讲这一页时建议用可视化先带一遍，再回到前面几页的数值定义。这样导师更容易把图感和指标对齐。",
            ),
            SlideSpec(
                title="Stepcurve Ghost 总览图",
                bullets=[
                    "按同样 step 顺序展示 ghost triplet。",
                    "重点观察主体横向是否变窄、是否出现双层、是否变暗或塌缩。",
                ],
                image_paths=[STEPCURVE_GHOST_GRID],
                image_caption="stepcurve ghost grid",
                layout="full_image",
                speaker_notes="ghost 图更适合回答‘视觉上重影有没有压下去’，但不要单独拿这一页当最终结论，必须和 point-support 主指标一起读。",
            ),
            SlideSpec(
                title="Mask-Boost Point-Support 总览图",
                bullets=[
                    "横向比较不同 alpha，纵向比较不同 step。",
                    "重点看 6src winner alpha=1 与 native 的差别，以及 12src 为什么后段出现主体碎裂或背景相对发灰。",
                ],
                image_paths=[MASK_BOOST_POINT_GRID],
                image_caption="mask-boost point-support grid",
                layout="full_image",
                speaker_notes="这页适合解释为什么 point-support 有时会‘看着像集中了一点’，但数值和 mask 内重建不一定同步跟上。",
            ),
            SlideSpec(
                title="Mask-Boost Ghost 总览图",
                bullets=[
                    "观察不同 alpha 下 ghost 是否继续下降，是否引入新的黑塌或结构破坏。",
                    "当前结果说明 winner alpha=1 也没有形成稳定 ghost 持续下降。",
                ],
                image_paths=[MASK_BOOST_GHOST_GRID],
                image_caption="mask-boost ghost grid",
                layout="full_image",
                speaker_notes="这页可以配合前面的 mask-boost 数字一起讲，结论会更稳：不是局部某一张图不好，而是整条线没有稳定收敛证据。",
            ),
            SlideSpec(
                title="最终结论与下一步建议",
                bullets=[
                    "当前本地代码对原版 VGGT 的主要贡献，不是换模型，而是把几何 teacher、损失路由、raw compare、ghost / point-support 指标和自动化 probing 全部补齐了。",
                    "实验层面最强结论: one-step 有短程收益；stepcurve 不支持继续深训到 32；当前 mask-boost 实现没有把 mask 内 PSNR 系统性拉升。",
                    "如果要更贴导师原意，下一版应尝试: 保留全图 loss，只对 GT human mask 的 eroded interior 叠加 boost，而不是 FG-only supervision。",
                    "因此，这份 PPT 最重要的信息不是‘某个数最好’，而是原版 VGGT 线已经被你们拆解到可以精确回答导师问题的程度。",
                ],
                speaker_notes=(
                    "收尾时要把代码价值和实验价值分开讲。代码层面，这套本地壳层已经非常完整。实验层面，则得出较克制但可信的结论：原版 VGGT 有短程改善，但当前还不足以证明深训收敛；mask 内增强要按更贴近导师原意的语义重做。"
                ),
            ),
        ]
    )
    slides[3:3] = build_visual_overview_slides(data)
    slides.extend(build_rich_appendix_slides(data))
    return slides


def build_visual_overview_slides(data: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})
    view = data["viewcount_map"]
    one = data["one_step_map"]
    return [
        SlideSpec(
            title="思维导图总览: 原版 VGGT 本地改动如何组织",
            bullets=[
                "中心节点始终是 original VGGT，本地改动围绕它往外长出 teacher、训练包装、compare/metrics 和自动化 probe 四层。",
                "这说明我们没有替换模型本体，而是在其外层建立了一整套可分析、可比较、可自动化复现的实验壳。",
                "导师如果只想先看全局，这一页就是最短路径。",
            ],
            formulas=[
                "Local original-VGGT line = model.pt + geometry teacher + finetune wrapper + compare metrics + taskized probes",
            ],
            code_refs=[
                "vggt_geom.py / finetune_vggt_pseudo.py / render_raw_compare.py / scripts/orig_vggt_*_probe/*",
            ],
            image_paths=[assets.get("mindmap")] if assets.get("mindmap") else [],
            image_caption="思维导图: 原版 VGGT 本地主线的代码层分解。",
            layout="full_image",
            speaker_notes="建议用这页先给导师建立地图感，再按后续章节逐层展开。",
        ),
        SlideSpec(
            title="全流程图: 从原版 VGGT 输出到最终报告",
            bullets=[
                "teacher 先把原版 VGGT 的 depth / conf / camera 输出整理成 pointmap 和几何元信息。",
                "precompute 把这些监督固化成 NPZ；finetune 再在此基础上做 loss routing；render_raw_compare 统一输出视觉对比图。",
                "one-step / stepcurve / mask-boost 三条 probe 线最终共用一套 compare 和后处理口径。",
            ],
            formulas=[
                "VGGT -> teacher -> NPZ -> finetune -> compare -> ghost/point-support/masked metrics -> summaries",
            ],
            code_refs=[
                "precompute_zju_vggt_geom.py / modal_run_train.py / finetune_vggt_pseudo.py / summarize_runs.py",
            ],
            image_paths=[assets.get('pipeline')] if assets.get('pipeline') else [],
            image_caption="流程图: 数据与监督在本地代码中的流动路径。",
            layout="full_image",
            speaker_notes="这一页帮助把代码模块和实验流程对应起来。",
        ),
        SlideSpec(
            title="Baseline 视角数对比图",
            bullets=[
                f"6src_hist: coverage={format_float(view['6src_hist']['coverage_ratio'])}, ghost={format_float(view['6src_hist']['ghost_visual_score'])}, native_psnr={format_float(view['6src_hist']['native_psnr'])}。",
                f"12src_nested: coverage={format_float(view['12src_nested']['coverage_ratio'])}, ghost={format_float(view['12src_nested']['ghost_visual_score'])}, native_psnr={format_float(view['12src_nested']['native_psnr'])}。",
                f"23cam_fullset: coverage={format_float(view['23cam_fullset']['coverage_ratio'])}, ghost={format_float(view['23cam_fullset']['ghost_visual_score'])}, native_psnr={format_float(view['23cam_fullset']['native_psnr'])}。",
                "覆盖率随视角数上升，但 ghost 和 native PSNR 并不自动单调更优，这也是后面需要 one-step 和 stepcurve 的原因。",
            ],
            formulas=[
                "Coverage != ghost quality != native PSNR; 多视角更多不等于单主体更稳",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_viewcount_summary_latest.json",
            ],
            image_paths=[assets.get('baseline_bars')] if assets.get('baseline_bars') else [],
            image_caption="baseline bar charts: coverage / ghost / native PSNR。",
            layout="full_image",
            speaker_notes="这一页用图直观说明为什么不能只用视角数解释好坏。",
        ),
        SlideSpec(
            title="One-Step 核心结果图",
            bullets=[
                f"6src_hist: ghost {format_float(one['6src_hist']['pre_ghost_visual_score'])} -> {format_float(one['6src_hist']['post_ghost_visual_score'])}，native_psnr {format_float(one['6src_hist']['pre_native_psnr'])} -> {format_float(one['6src_hist']['post_native_psnr'])}。",
                f"12src_nested: ghost {format_float(one['12src_nested']['pre_ghost_visual_score'])} -> {format_float(one['12src_nested']['post_ghost_visual_score'])}，native_psnr {format_float(one['12src_nested']['pre_native_psnr'])} -> {format_float(one['12src_nested']['post_native_psnr'])}。",
                "结论是: 一步微调确实有短程 ghost 收益，但 PSNR 并没有同步稳定抬升，所以不能直接外推成会长期收敛。",
            ],
            formulas=[
                "Δghost < 0 indicates short-term visual gain; ΔPSNR may still be <= 0",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_one_step_probe_summary_latest.json",
            ],
            image_paths=[assets.get('one_step_pairs')] if assets.get('one_step_pairs') else [],
            image_caption="one-step pre/post paired charts。",
            layout="full_image",
            speaker_notes="这页是从 baseline 过渡到 stepcurve 的桥。",
        ),
    ]


def build_rich_appendix_slides(data: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})
    stepcurve = data["stepcurve_map"]
    mask_boost = data["mask_boost_map"]
    step12 = stepcurve["12src_nested"]
    step6 = stepcurve["6src_hist"]
    mb12 = mask_boost["12src_nested"]
    mb6 = mask_boost["6src_hist"]
    s12_0 = step_row(step12, 0)
    s12_1 = step_row(step12, 1)
    s12_4 = step_row(step12, 4)
    s12_16 = step_row(step12, 16)
    s6_0 = step_row(step6, 0)
    s6_1 = step_row(step6, 1)
    s6_16 = step_row(step6, 16)
    return [
        SlideSpec(
            title="附录: 反投影公式为何是几何监督根基",
            bullets=[
                "几何 teacher 的根不是‘网络又预测了一张点图’，而是把 depth 和相机内外参组成一个可验证的几何映射。",
                "如果反投影约定错了，点云会整体偏移，后续 point loss、MV-depth 和 compare 都会被系统性污染。",
                "因此本地代码把 unproject_impl 和 pointmap_frame 都写进元数据，而不是隐式假设。",
            ],
            formulas=[
                "x_cam = (u - cx) * z / fx",
                "y_cam = (v - cy) * z / fy",
                "P_world = R^T * ([x_cam, y_cam, z]^T - t)",
            ],
            code_refs=[
                "vggt_geom.py: _unproject_depth_to_world_batched",
                "vggt/utils/geometry.py: unproject_depth_map_to_point_map",
            ],
            speaker_notes="要强调：这套几何不是拍脑袋近似，而是标准针孔模型反投影。",
        ),
        SlideSpec(
            title="附录: point_head_frame=auto 的意义",
            bullets=[
                "原版 VGGT 的点头输出和 depth 反投影不一定天然处在同一个最稳参考系里。",
                "point_head_frame=auto 会比较候选参考系的自重投影误差，选择误差更小的一套。",
                "这一步的目标不是追求数学花样，而是让 teacher 更少把坐标系偏差误当成重建误差。",
            ],
            formulas=[
                "choose frame* = argmin_frame reprojection_error(frame)",
            ],
            code_refs=[
                "vggt_geom.py: forward_prepared_batch",
            ],
            speaker_notes="这是一个很工程但很关键的稳健性修补。",
        ),
        SlideSpec(
            title="附录: NPZ 预计算到底保存了什么",
            bullets=[
                "每帧 NPZ 不是只存 pointmap，而是把多视图图像、depth、depth_conf、pointmap、extrinsic、intrinsic 以及 pointmap_source / frame / unproject_impl 全部存下来。",
                "这样后续训练才知道每个监督来自哪里，也方便 probe 汇总时解释差异来源。",
                "预计算把 expensive teacher 前向和 cheap finetune 解耦，是整条本地主线能高频试验的前提。",
            ],
            formulas=[
                "NPZ(frame) = {img, depth, conf, pointmap, K, [R|t], pointmap_source, pointmap_frame, unproject_impl}",
            ],
            code_refs=[
                "precompute_zju_vggt_geom.py: save fields near NPZ write",
                "modal_run_train.py: pointmap_source / precompute_unproject_impl",
            ],
            speaker_notes="这页适合回答‘为什么不直接在线算 teacher’。",
        ),
        SlideSpec(
            title="附录: 置信度加权如何决定监督强弱",
            bullets=[
                "teacher 的每个像素不是等价样本，conf weight 先把低置信度区抑掉，再让高置信度像素主导梯度。",
                "如果配合 per-view quantile mask，就相当于每个视角只挑前 q 分位的可信像素进入监督。",
                "这能减少 teacher 噪声把训练往错误几何上推。",
            ],
            formulas=[
                "w_conf(x) = valid(x) * clip((c(x)-t)/(1-t), 0, 1)^γ",
                "quantile mask(view v) = 1[c_v(x) in top-q among valid pixels]",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: _build_conf_weight / _build_per_view_conf_quantile_mask",
            ],
            speaker_notes="这里可以补一句：后续大部分权重增强都是在 w_base 上叠加，不是另起炉灶。",
        ),
        SlideSpec(
            title="附录: FG gating 和 interior boost 不是一回事",
            bullets=[
                "FG gating 决定哪些像素还能留下来当 valid supervision；interior boost 决定留下来的前景内部权重再乘多大。",
                "当前 mask-boost probe 把这两件事同时打开了，所以它比导师原话更激进。",
                "这也是解读结果时必须谨慎的原因。",
            ],
            formulas=[
                "M_valid = M_valid_all * (M_fg + b * (1 - M_fg))",
                "w_final = w_base * [1 + (β - 1) * M_interior]",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: foreground gating + _apply_fg_supervision_boost",
            ],
            speaker_notes="这一页必须明确告诉导师：当前实现偏离了‘全图 loss 保留，只在人体内部增强’的温和语义。",
        ),
        SlideSpec(
            title="附录: 为什么要先 erosion 再 boost",
            bullets=[
                "直接用原始 mask 全区域 boost，边界壳层会被一起抬高，容易把轮廓误差也放大。",
                "先 erosion 2px 的思路是尽量只留下主体 interior，让增强更多作用在人体内部而不是边缘。",
                "这本质上是在做一个形态学安全边界。",
            ],
            formulas=[
                "M_interior = erode(M_fg, r)",
                "erode(x) = 1 - dilate(1 - x)",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: _erode_mask_tensor / _build_fg_supervision_boost_mask",
            ],
            speaker_notes="虽然当前结果没有系统性提高 mask 内 PSNR，但 erosion 本身的设计动机是合理的。",
        ),
        SlideSpec(
            title="附录: point_target_mode 的四种路由",
            bullets=[
                "pointmap: 直接把 teacher pointmap 当目标。",
                "depth_unproject: 只用 depth 反投影的点作为目标。",
                "blend: 根据 depth-unproject 和 pseudo pointmap 的一致性做软混合。",
                "depth_consensus_unproject: 进一步偏向跨视图也能自洽的 depth 反投影区域。",
            ],
            formulas=[
                "r(x) = exp(-||U(D_tgt)(x) - P_pseudo(x)|| / τ)",
                "P_tgt(x) = α(x) U(D_tgt)(x) + (1 - α(x)) P_pseudo(x)",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: --point_target_mode and point target routing",
            ],
            speaker_notes="这一页强调：本地代码没有把某一种几何来源神化，而是允许它们互相校正。",
        ),
        SlideSpec(
            title="附录: MV support 与 MV-depth 为什么对 ghost 对症",
            bullets=[
                "重影、双层、厚壳，本质都是一个点在不同视角里无法一致解释成单一前表面。",
                "MV support 先用 agree ratio 给每个点一个‘跨视图有多少人认同你’的权重。",
                "MV-depth reprojection 则直接惩罚投到别的视角后深度对不上的点。",
            ],
            formulas=[
                "agree_ratio(x) = #agree_views(x) / (V - 1)",
                "s_mv(x) = floor + (1 - floor) * agree_ratio(x)",
                "L_mvdepth = Avg_{i->j}[Σ_x w_ij(x) ρ(|z_ij(x) - D_j(π_ij(x))|) / Σ_x w_ij(x)]",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: _point_multiview_support_weight / _point_multiview_depth_reproj_loss",
            ],
            speaker_notes="这页是 ghost 问题的几何层解释。",
        ),
        SlideSpec(
            title="附录: raw compare 的点云图为什么会出现背景灰斑",
            bullets=[
                "weight_native.png 不是绝对支持质量图，而是把 hit 区的 weight_map 用 p99 做归一化后再裁到 [0,1]。",
                "因此当主体支持分布整体形状变化时，背景一些弱响应也可能被相对抬亮成灰斑。",
                "这说明它适合看形状、集中度、连通性，不适合直接解读成绝对物理能量。",
            ],
            formulas=[
                "weight01 = clip(weight_map / p99(weight_map[hit]), 0, 1)",
            ],
            code_refs=[
                "scripts/orig_vggt_viewcount/render_raw_compare.py: weight normalization",
            ],
            speaker_notes="这页正面回应你们最近关于灰斑的疑问。",
        ),
        SlideSpec(
            title="附录: ghost 分数的分解项",
            bullets=[
                "ghost score 先看 width overflow、area overflow、center shift 和 peak_count。",
                "再叠加 dark penalty 与 collapse penalty，惩罚太暗和过度塌缩。",
                "这让它既能抓双层外扩，也能抓‘全黑了所以看着没 ghost’这种假好结果。",
            ],
            formulas=[
                "ghost_score = width overflow + area overflow + multi-peak + center shift",
                "ghost_visual_score = ghost_score + dark_penalty + collapse_penalty",
            ],
            code_refs=[
                "tools/score_ghosting_from_cat_pred.py: _score_one",
            ],
            speaker_notes="这里要提醒导师：ghost 只是 heuristic，不是 3D 真值。",
        ),
        SlideSpec(
            title="附录: point-support 主指标与 mask 内重建指标如何分工",
            bullets=[
                "point-support 主指标回答‘点云支持图是否更像单主体’。",
                "masked PSNR / L1 回答‘人体区域重建到底有没有更准’。",
                "这两组指标在当前实验里并不总是同向，这正是实验最有信息量的地方。",
            ],
            formulas=[
                "subject_support_share = ΣSM / ΣS",
                "largest_component_share = max_c Σ_{x in c}S / Σ_{x in active}S",
                "subject_psnr = -10log10(Σ((pred-tgt)^2M)/(ΣM*C))",
            ],
            code_refs=[
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
            ],
            speaker_notes="建议这里明确说：GT mask 优先来自 report.json.meta.tgt_mask_path。",
        ),
        SlideSpec(
            title="附录图: 12src_nested 的 stepcurve 曲线",
            bullets=[
                f"step0000 -> step0001: ghost {format_float(s12_0['ghost_visual_score'])} -> {format_float(s12_1['ghost_visual_score'])}，largest_component_share {format_float(s12_0['largest_component_share'])} -> {format_float(s12_1['largest_component_share'])}。",
                f"step0004 时 ghost 降到 {format_float(s12_4['ghost_visual_score'])}，但 support_peak_count={s12_4['support_peak_count']}，主体并非持续更单峰。",
                f"step0016 时 largest_component_share 掉到 {format_float(s12_16['largest_component_share'])}，secondary_component_mass 升到 {format_float(s12_16['secondary_component_mass'])}。",
                "所以 12src 的结论是短程能变好，但不支持继续深训。",
            ],
            formulas=[
                "FIRST_REBOUND_STEP = step0008",
                "EXTEND_DECISION = NO_STOP_AT_16",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_stepcurve_probe_summary_latest.json",
            ],
            image_paths=[assets.get('stepcurve_12src_nested')] if assets.get('stepcurve_12src_nested') else [],
            image_caption="12src stepcurve metrics chart。",
            layout="full_image",
            speaker_notes="这页比总览图更适合讲时间趋势。",
        ),
        SlideSpec(
            title="附录图: 6src_hist 的 stepcurve 曲线",
            bullets=[
                f"step0000 -> step0001: ghost {format_float(s6_0['ghost_visual_score'])} -> {format_float(s6_1['ghost_visual_score'])}，但 subject_psnr {format_float(s6_0['subject_psnr'])} -> {format_float(s6_1['subject_psnr'])} 只微降。",
                f"step0016 时 subject_psnr={format_float(s6_16['subject_psnr'])}，largest_component_share={format_float(s6_16['largest_component_share'])}，secondary_component_mass={format_float(s6_16['secondary_component_mass'])}。",
                "6src 比 12src 稍稳，但同样没有形成可以自信扩到 32 的单调收敛曲线。",
            ],
            formulas=[
                "FIRST_REBOUND_STEP = step0001",
                "worth_extending_to_32 = False",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_stepcurve_probe_summary_latest.json",
            ],
            image_paths=[assets.get('stepcurve_6src_hist')] if assets.get('stepcurve_6src_hist') else [],
            image_caption="6src stepcurve metrics chart。",
            layout="full_image",
            speaker_notes="和 12src 对照着讲，会很清楚。",
        ),
        SlideSpec(
            title="附录图: 12src_nested 的 mask-boost winner 曲线",
            bullets=[
                f"winner_alpha={mb12['winner_alpha']}，winner_fg_supervision_boost={format_float(mb12['winner_fg_supervision_boost'], 1)}。",
                "当前 winner 只是相对不那么差，不代表它真的优于 native。",
                "尤其在 subject PSNR 上，12src 的 boost 版相对 native 是持续负增益。",
            ],
            formulas=[
                "alpha=1 -> fg_supervision_boost=2.0",
                "12src: SUBJECT_RECON_IMPROVES = NO",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_mask_boost_probe_summary_latest.json",
            ],
            image_paths=[assets.get('mask_boost_12src_nested')] if assets.get('mask_boost_12src_nested') else [],
            image_caption="12src mask-boost winner vs native chart。",
            layout="full_image",
            speaker_notes="这页需要明确告诉导师：当前实现下，12src 的 masked reconstruction 是系统性变差的。",
        ),
        SlideSpec(
            title="附录图: 6src_hist 的 mask-boost winner 曲线",
            bullets=[
                f"winner_alpha={mb6['winner_alpha']}，winner_fg_supervision_boost={format_float(mb6['winner_fg_supervision_boost'], 1)}。",
                "6src 在 point-support 主体性上有局部正向迹象，但 subject PSNR 只是晚段局部回升，不是系统性稳定优势。",
                "所以 6src 的标签才会是 POINT_SUPPORT_SINGLE_SUBJECT=YES，但 SUBJECT_RECON_IMPROVES=NO。",
            ],
            formulas=[
                "6src: POINT_SUPPORT_SINGLE_SUBJECT = YES",
                "6src: SUBJECT_RECON_IMPROVES = NO",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_mask_boost_probe_summary_latest.json",
            ],
            image_paths=[assets.get('mask_boost_6src_hist')] if assets.get('mask_boost_6src_hist') else [],
            image_caption="6src mask-boost winner vs native chart。",
            layout="full_image",
            speaker_notes="这是最容易被误读的一页，要把 point-support gain 和 masked reconstruction gain 分开。",
        ),
        SlideSpec(
            title="附录: 当前失败样例与工程结论",
            bullets=[
                "部分 deeper step 或某些 alpha 分支失败，不是脚本没跑完，而是 compare / alignment 阶段出现 sim3 alignment rmse_after too high。",
                "例如 12src step0008 和部分 6src alpha4 分支，都被明确记录为 failed 或 blocked。",
                "这些失败本身也是结论的一部分: 当前训练语义下，越往后推越容易出现几何不稳定。",
            ],
            formulas=[
                "failure trigger example: rmse_after > threshold => compare task failed",
            ],
            code_refs=[
                "logs/modal_phase5/reports/orig_vggt_stepcurve_probe_summary_latest.json",
                "logs/modal_phase5/reports/orig_vggt_mask_boost_probe_summary_latest.json",
            ],
            speaker_notes="不要把失败当作脏数据删掉，导师更关心失败为什么发生。",
        ),
        SlideSpec(
            title="附录: 更贴导师原意的下一版实现建议",
            bullets=[
                "保留全图 loss，不关闭背景通道，只在 GT human mask 的 eroded interior 上乘一个温和 boost。",
                "先只做 alpha=1/2/4 的短程筛选，并且同时跟踪 point-support、ghost、subject_psnr、subject_l1 四组指标。",
                "如果 masked PSNR 仍旧没有系统性抬升，就说明问题不在‘前景内部权重不够’，而在更深层的几何表达与渲染机理。",
            ],
            formulas=[
                "recommended next loss: w_final = w_base * [1 + α * erode(M_gt, r)]  with background path kept alive",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: foreground gating / boost path",
            ],
            speaker_notes="这页能把当前实验结论自然过渡到下一步，而不是只停在否定。",
        ),
    ]


def add_panel(slide, left, top, width, height, fill=BOX):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def render_bullets_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, spec.title)
    left_box = add_panel(slide, Inches(0.55), Inches(1.25), Inches(7.35), Inches(5.7))
    left_frame = left_box.text_frame
    add_bullets(left_frame, wrap_lines(spec.bullets, 48), font_size=18)

    right_box = add_panel(slide, Inches(8.1), Inches(1.25), Inches(4.65), Inches(5.7), fill=RGBColor(244, 239, 231))
    right_frame = right_box.text_frame
    right_frame.clear()
    right_frame.word_wrap = True
    right_frame.vertical_anchor = MSO_ANCHOR.TOP
    sections: list[tuple[str, list[str]]] = []
    if spec.formulas:
        sections.append(("公式 / 规则", spec.formulas))
    if spec.code_refs:
        sections.append(("代码位置", spec.code_refs))
    if spec.image_caption:
        sections.append(("图像说明", [spec.image_caption]))

    for s_idx, (header, lines) in enumerate(sections):
        p = right_frame.paragraphs[0] if s_idx == 0 else right_frame.add_paragraph()
        p.text = header
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = header
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        p.space_after = Pt(4)
        for line in wrap_lines(lines, 32):
            pp = right_frame.add_paragraph()
            pp.text = line
            if pp.runs:
                rr = pp.runs[0]
            else:
                rr = pp.add_run()
                rr.text = line
            rr.font.name = "Consolas" if any(ch in line for ch in "=<>[]{}_") else "Microsoft YaHei"
            rr.font.size = Pt(12)
            rr.font.color.rgb = TEXT
            pp.space_after = Pt(2)


def render_two_images_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, spec.title)
    top_box = add_panel(slide, Inches(0.55), Inches(1.15), Inches(12.2), Inches(1.1))
    add_bullets(top_box.text_frame, wrap_lines(spec.bullets, 120), font_size=16)
    if len(spec.image_paths) >= 2:
        left = ensure_cache_image(spec.image_paths[0], safe_stem(f"{spec.title}_left"))
        right = ensure_cache_image(spec.image_paths[1], safe_stem(f"{spec.title}_right"))
        fit_image(slide, left, Inches(0.7), Inches(2.45), Inches(5.8), Inches(3.7))
        fit_image(slide, right, Inches(6.85), Inches(2.45), Inches(5.8), Inches(3.7))
    add_textbox(slide, Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.45), spec.image_caption, 13, False, SUBTLE)
    if spec.formulas or spec.code_refs:
        meta = add_panel(slide, Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.55), fill=RGBColor(244, 239, 231))
        lines = []
        if spec.formulas:
            lines.append("公式: " + " | ".join(spec.formulas[:2]))
        if spec.code_refs:
            lines.append("代码: " + " | ".join(spec.code_refs[:2]))
        add_textbox(slide, Inches(0.9), Inches(6.78), Inches(11.5), Inches(0.28), "  ".join(lines), 11, False, TEXT)


def render_full_image_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, spec.title)
    info_box = add_panel(slide, Inches(0.55), Inches(1.05), Inches(12.2), Inches(1.0))
    add_bullets(info_box.text_frame, wrap_lines(spec.bullets, 115), font_size=15)
    if spec.image_paths:
        img = ensure_cache_image(spec.image_paths[0], safe_stem(spec.title))
        fit_image(slide, img, Inches(0.55), Inches(2.15), Inches(12.2), Inches(4.3))
    meta_lines = []
    if spec.formulas:
        meta_lines.append("公式: " + " | ".join(spec.formulas[:3]))
    if spec.code_refs:
        meta_lines.append("代码: " + " | ".join(spec.code_refs[:2]))
    if meta_lines:
        meta = add_panel(slide, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.45), fill=RGBColor(244, 239, 231))
        add_textbox(slide, Inches(0.75), Inches(6.67), Inches(11.8), Inches(0.2), "  ".join(meta_lines), 10, False, TEXT)
    if spec.image_caption:
        add_textbox(slide, Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.2), spec.image_caption, 10, False, SUBTLE)


def build_presentation(slides: list[SlideSpec]) -> Presentation:
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    for spec in slides:
        if spec.layout == "two_images":
            render_two_images_slide(prs, spec)
        elif spec.layout == "full_image":
            render_full_image_slide(prs, spec)
        else:
            render_bullets_slide(prs, spec)
    return prs


def download_official_architecture() -> Path:
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    if OFFICIAL_ARCH_PNG.exists() and OFFICIAL_ARCH_PNG.stat().st_size > 0:
        return OFFICIAL_ARCH_PNG
    try:
        urlretrieve(OFFICIAL_ARCH_URL, OFFICIAL_ARCH_PNG)
    except Exception:
        if not OFFICIAL_ARCH_PNG.exists():
            raise
    return OFFICIAL_ARCH_PNG


def existing_images(*paths: Path) -> list[Path]:
    return [path for path in paths if picture_exists(path)]


def collect_code_diff_summary() -> dict:
    return {
        "direct_original_touches": [
            "vggt/utils/geometry.py: 给深度反投影增加 unproject_impl / pixel_center_offset",
            "vggt/utils/load_fn.py: 加 threaded image loading，减少多视图读图开销",
            "demo_colmap.py: 补 pycolmap 注册兼容、调试输出和 BA 可用性兜底",
            "vggt/models/vggt.py: 主要是注释补充，不是结构性改写",
        ],
        "new_local_modules": [
            "vggt_geom.py: 把原版输出整理成几何 teacher",
            "precompute_zju_vggt_geom.py: 预计算 depth/conf/pointmap/camera NPZ",
            "modal_run_train.py: 远端预计算和训练入口",
            "finetune_vggt_pseudo.py: 微调 wrapper、loss routing、FG boost、MV-depth",
            "scripts/orig_vggt_viewcount/render_raw_compare.py: 统一生成对照图",
            "tools/score_ghosting_from_cat_pred.py: 把重影现象量化成 ghost_visual_score",
            "scripts/orig_vggt_one_step_probe/*: 一步微调探针",
            "scripts/orig_vggt_stepcurve_probe/*: stepcurve 探针",
            "scripts/orig_vggt_mask_boost_probe/*: mask-boost 探针",
        ],
        "secondary_touches": [
            "training/config/default.yaml",
            "training/config/default_dataset.yaml",
            "training/data/dynamic_dataloader.py",
        ],
        "core_message": "真正直接碰原版代码树的入口很少，主体工作是围绕原版 VGGT 新增 teacher、预计算、微调和 probe 壳层。",
    }


def build_official_compare_slide(data: dict) -> SlideSpec:
    assets = data.get("assets", {})
    arch = assets.get("current_architecture")
    return SlideSpec(
        title="当前本地架构图: 参考官方画法，但内容对应现在这套代码",
        bullets=[
            "这张图不再直接放官方原图，而是参考官方版式，重画了一张“当前本地原版 VGGT 主线”的架构图。",
            "上半部分还是原版 VGGT 核心：patch tokens、camera token、aggregator、camera head、DPT heads。",
            "下半部分是我们现在真正新增的本地壳层：VGGTGeomTeacher、NPZ precompute、finetune wrapper、compare/metrics 和三条 probe。",
        ],
        formulas=[
            "原版核心: Aggregator + Camera Head + DPT Heads",
            "本地新增: Teacher -> NPZ -> Finetune -> Compare -> Probe",
        ],
        code_refs=[
            f"风格参考: {OFFICIAL_ARCH_URL}",
            "vggt_geom.py",
            "precompute_zju_vggt_geom.py",
            "finetune_vggt_pseudo.py",
        ],
        image_paths=existing_images(arch) if arch else [],
        image_caption="参考官方架构图风格重画，内容对应当前本地原版 VGGT 代码与流程链路。",
        layout="full_image",
        speaker_notes="这一页先让导师看到结构位置。最重要的信息是：上半部分保留原版 VGGT 核心，下半部分是我们新增的训练和分析壳层。",
    )


def build_metrics_taxonomy_slide(title: str = "先看什么图、再读什么指标") -> SlideSpec:
    return SlideSpec(
        title=title,
        bullets=[
            "左图是 point-support，可直观看主体是否集中、是否裂成两团、第二峰是否变大。",
            "右图是 ghost triplet，可直观看重影是否变宽、变暗、出现双层或横向漂移。",
            "导师问“主体有没有更稳、重影有没有更少”时，先看这两类图，再回到对应数字。",
        ],
        formulas=[
            "subject_support_share = sum(S * M) / sum(S)",
            "subject_psnr = -10 log10(sum((pred - tgt)^2 * M) / (3 * sum(M)))",
            "ghost_visual_score 越低越好",
        ],
        code_refs=[
            "scripts/orig_vggt_viewcount/render_raw_compare.py",
            "tools/score_ghosting_from_cat_pred.py",
            "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
        ],
        image_paths=existing_images(POINT_TAXONOMY_EXAMPLE, GHOST_TAXONOMY_EXAMPLE),
        image_caption="左: point-support 示例。右: ghost triplet 示例。两类图各看一件事，不要混着解释。",
        layout="two_images",
        speaker_notes="这页的作用是先统一读图口径，避免汇报时把 support 图、ghost 图和训练期 debug weight 图混成一类。",
    )


def build_combined_deck(data: dict, diff_summary: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})

    return [
        SlideSpec(
            title="基于原版 VGGT 的本地代码改动总览",
            bullets=[
                "这套 PPT 只讲代码相对原版 VGGT 改了什么，重点放在结构、原理、公式和文件入口。",
                "核心口径是: 原版主干基本保留，新增工作集中在 geometry teacher、预计算、微调 wrapper、compare 和 probe 这几层壳。",
                "讲解顺序按“架构位置 -> 文件入口 -> 原理/公式 -> 脚手架职责”展开，方便导师快速抓重点。",
            ],
            speaker_notes="开场先定边界：这次只汇报相对原版的代码改造与结构补充，不展开任何结果判断。",
        ),
        build_official_compare_slide(data),
        SlideSpec(
            title="改动位置一页看完",
            bullets=[
                "左半边是直接改到原版代码树的入口，数量不多，重点是 geometry、load_fn 和 demo_colmap。",
                "右半边是围绕原版新增的本地壳层，包括 teacher、precompute、finetune、compare/metrics 和 probe 脚手架。",
                diff_summary["core_message"],
            ],
            formulas=[
                "原版主干基本不动",
                "新增重点在图外流程，不在图内 backbone",
            ],
            code_refs=diff_summary["direct_original_touches"] + diff_summary["new_local_modules"],
            image_paths=existing_images(assets.get("code_touch_map")),
            image_caption="左: 直接改到原版仓库的入口。右: 本地新增的外层代码壳。",
            layout="full_image",
            speaker_notes="重点是让导师先分清“改原版”与“加外壳”两类工作。",
        ),
        SlideSpec(
            title="直接改到原版代码树的地方",
            bullets=[
                "真正直接改动原版文件的入口很少，这说明工作重点不是重写主干，而是扩展外围流程。",
                "影响训练语义的关键改动主要在 `vggt/utils/geometry.py` 和 `vggt/utils/load_fn.py`。",
                "`demo_colmap.py` 解决的是 demo / BA 兼容性；`vggt/models/vggt.py` 主要是结构说明级补充。",
            ],
            formulas=[
                "原版仓库改动少",
                "主要新增在外层训练与分析流程",
            ],
            code_refs=[
                "vggt/utils/geometry.py",
                "vggt/utils/load_fn.py",
                "demo_colmap.py",
                "vggt/models/vggt.py",
                "training/config/default.yaml",
                "training/config/default_dataset.yaml",
                "training/data/dynamic_dataloader.py",
            ],
            speaker_notes="这页只回答“原版内部到底动了哪几处”。",
        ),
        SlideSpec(
            title="Geometry teacher 与 NPZ 预计算",
            bullets=[
                "新增 `vggt_geom.py`，把原版 VGGT 的 depth / conf / camera 输出整理成可复用的几何 teacher。",
                "新增 `precompute_zju_vggt_geom.py`，把几何 teacher 固化为 NPZ 数据契约，供后续训练直接读取。",
                "这样后续所有微调与 probe 都共享同一套 teacher 输入，不需要每次重新在线跑原版前向。",
            ],
            formulas=[
                "x_cam = (u + o - cx) * z / fx",
                "y_cam = (v + o - cy) * z / fy",
                "P_world = R^T * ([x_cam, y_cam, z]^T - t)",
            ],
            code_refs=[
                "vggt_geom.py",
                "precompute_zju_vggt_geom.py",
                "modal_run_train.py",
                "vggt/utils/geometry.py",
            ],
            image_paths=existing_images(assets.get("teacher_npz_flow")),
            image_caption="原版输出先被整理为 geometry teacher，再被固化为 NPZ 数据契约。",
            layout="full_image",
            speaker_notes="这页讲清楚 teacher 和 precompute 的职责分工。",
        ),
        SlideSpec(
            title="finetune wrapper: 把监督路由起来",
            bullets=[
                "新增 `finetune_vggt_pseudo.py`，把 depth、point、reproj、normal、MV-depth、camera 等监督统一收在一个 wrapper 里。",
                "这层代码不是改原版 backbone，而是把 teacher、dataset、loss 和 train loop 接成一个可控的微调入口。",
                "这样每个训练分支都能复用同一套数据入口和监督路由。",
            ],
            formulas=[
                "L_total = sum_k lambda_k L_k",
                "w_conf(x) = valid(x) * clip((c(x) - t) / (1 - t), 0, 1)^gamma",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py",
            ],
            image_paths=existing_images(assets.get("loss_routing")),
            image_caption="新增的 wrapper 统一管理数据集、权重构造和损失路由。",
            layout="full_image",
            speaker_notes="这里只讲结构与公式，不谈效果。",
        ),
        SlideSpec(
            title="FG mask boost 的当前代码语义",
            bullets=[
                "当前实现的重点不是结果，而是语义：先用前景 mask 进行 valid gating，再对前景内部做额外增权。",
                "这页的作用是把代码语义讲清楚，避免把当前实现误说成别的训练策略。",
                "它属于训练路由的一部分，而不是原版主干的一部分。",
            ],
            formulas=[
                "M_valid = M_valid_all * M_fg",
                "M_interior = erode(M_fg, r)",
                "w_final = w_base * [1 + (beta - 1) * M_interior]",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: foreground gating / boost path",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="FG mask boost 在代码里体现为 valid gating 与 interior boost 的组合。",
            layout="full_image",
            speaker_notes="这页只讲机制，不讲任何跑出来的优劣。",
        ),
        SlideSpec(
            title="Compare / Metrics / Probe 脚手架",
            bullets=[
                "新增 `render_raw_compare.py` 统一生成对照图，再由 ghost scorer 和 point-support 脚本统一读图与汇总。",
                "新增 one-step / stepcurve / mask-boost 三套 probe 脚本，把不同训练分支封装成可重复执行的自动化任务。",
                "这些脚本是围绕同一条原版 VGGT 主线搭的分析壳层，不是新的模型家族。",
            ],
            formulas=[
                "原版核心 + compare shell + probe shell",
            ],
            code_refs=[
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
                "scripts/orig_vggt_one_step_probe/*",
                "scripts/orig_vggt_stepcurve_probe/*",
                "scripts/orig_vggt_mask_boost_probe/*",
            ],
            image_paths=existing_images(assets.get("compare_probe_framework")),
            image_caption="compare / metrics / probe 是围绕同一条原版主线组织起来的自动化代码壳。",
            layout="full_image",
            speaker_notes="强调这是工程脚手架，不讲任何实验现象。",
        ),
        SlideSpec(
            title="最后一句话怎么向导师汇报",
            bullets=[
                "相对原版 VGGT，我们没有重写主干，而是补齐了 geometry teacher、NPZ precompute、finetune wrapper、compare 和 probe 这几层外壳。",
                "原版仓库内部真正被改动的入口很少，新增工作主要集中在外围训练与分析流程。",
                "因此这次汇报的重点不是“跑出来多好”，而是“现在这条原版 VGGT 本地代码线已经被拆解成可说明、可复用、可继续扩展的结构”。",
            ],
            speaker_notes="收尾只保留代码结构口径，不带任何结果判断。",
        ),
    ]


def build_code_and_finetune_deck(data: dict, diff_summary: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})

    return [
        SlideSpec(
            title="原版 VGGT 改动与微调: 只讲代码主线",
            bullets=[
                "这套只讲基于原版 VGGT 的最近代码改动和微调主线，重点是结构、公式和文件职责。",
                "核心信息只有两层: 主干结构没重写；训练与分析壳层已经补齐。",
                "汇报顺序按“当前架构 -> 原版文件改动 -> teacher/precompute -> finetune wrapper -> 训练语义”展开。",
            ],
            speaker_notes="这套单独讲给导师时，要突出“最近这次基于原版的代码改造”这个主题。",
        ),
        build_official_compare_slide(data),
        SlideSpec(
            title="直接改到原版文件的地方",
            bullets=[
                "最关键的原版改动只有少数入口，不是大面积改 backbone。",
                "真正改变训练语义的，是 `geometry.py` 的反投影约定和 `load_fn.py` 的多线程读图。",
                "`demo_colmap.py` 解决 demo / BA 兼容；`vggt/models/vggt.py` 主要是结构说明级补充。",
            ],
            formulas=[
                "原版代码树改动少",
                "核心新增在外层训练与评估流程",
            ],
            code_refs=diff_summary["direct_original_touches"] + diff_summary["secondary_touches"],
            image_paths=existing_images(assets.get("code_touch_map")),
            image_caption="左侧是直接修改原版仓库的入口，右侧是围绕原版新增的代码壳。",
            layout="full_image",
            speaker_notes="这里把“原版内改了哪里”讲干净，导师最容易接受。",
        ),
        SlideSpec(
            title="新增模块 1: geometry teacher + NPZ precompute",
            bullets=[
                "先用 `VGGTGeomTeacher` 把原版输出整理成 pointmap / depth_conf / camera supervision。",
                "再用 `precompute_zju_vggt_geom.py` 把这些监督写进 NPZ，后续训练直接读固定 teacher 结果。",
                "这样所有训练分支都共享同一套预计算输入契约。",
            ],
            formulas=[
                "Teacher -> NPZ -> Finetune",
                "P_world = R^T * ([x_cam, y_cam, z]^T - t)",
            ],
            code_refs=[
                "vggt_geom.py",
                "precompute_zju_vggt_geom.py",
                "modal_run_train.py",
            ],
            image_paths=existing_images(assets.get("teacher_npz_flow")),
            image_caption="几何 teacher 负责整理原版输出，precompute 负责把它们固化为 NPZ 契约。",
            layout="full_image",
            speaker_notes="这里强调 teacher 的价值在于把训练底座固定住，而不是在线反复跑原版前向。",
        ),
        SlideSpec(
            title="新增模块 2: finetune wrapper 把监督路由起来",
            bullets=[
                "`finetune_vggt_pseudo.py` 是这条主线真正的核心文件，负责数据集、loss routing、FG boost 和 MV-depth。",
                "从导师视角看，它做的是把“原版输出”变成“可控的微调入口”。",
                "后续所有训练分支都是在这层 wrapper 上继续加开关和参数。",
            ],
            formulas=[
                "L_total = sum_k lambda_k L_k",
                "w_conf(x) = valid(x) * clip((c(x) - t) / (1 - t), 0, 1)^gamma",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py",
            ],
            image_paths=existing_images(assets.get("loss_routing")),
            image_caption="wrapper 统一处理数据读取、权重构造和损失路由。",
            layout="full_image",
            speaker_notes="这里不要陷入太多细枝末节，抓住 routing 这个词就够了。",
        ),
        SlideSpec(
            title="这次微调里最关键的两处设计",
            bullets=[
                "第一处是前景权重路径: 当前实现更接近 foreground gating + interior boost。",
                "第二处是 MV-depth: 用跨视图深度重投影把多视图几何一致性接进训练。",
                "这两处都属于 wrapper 层的训练语义改动，不属于原版主干结构改动。",
            ],
            formulas=[
                "w_final = w_base * [1 + (beta - 1) * M_interior]",
                "L_mvdepth ~ Avg rho(|z_ij(x) - D_j(pi_ij(x))|)",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: foreground gating / boost path",
                "finetune_vggt_pseudo.py: multiview depth reprojection path",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="FG mask boost 的代码语义和权重构造方式。",
            layout="full_image",
            speaker_notes="这一页的价值在于把训练语义说准。",
        ),
        SlideSpec(
            title="这套代码改动的汇报口径",
            bullets=[
                "原版架构图里的核心模块没有被推翻。",
                "真正新增的是 teacher、NPZ、wrapper、compare 和 probe，这些东西把原版 VGGT 变成了一条可讲清楚的本地代码链。",
                "因此这套 PPT 只讲代码结构和原理，不讲任何跑出来的效果。",
            ],
            speaker_notes="最后回到最朴素的三句话，便于导师记住重点。",
        ),
    ]


def build_probe_deck(data: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})

    return [
        SlideSpec(
            title="One-step / Stepcurve: 自动化 probe 脚手架",
            bullets=[
                "这套只讲 one-step / stepcurve 这两支自动化脚手架在代码里做了什么，重点是任务编排、文件职责和输出契约。",
                "它们不是新的模型分支，而是围绕同一条原版 VGGT 主线组织出来的训练与汇总脚本。",
                "从导师视角看，这一支更像自动化实验外壳，而不是模型结构改造。",
            ],
            speaker_notes="开头先把 probe 的性质讲清楚：它是脚手架，不是新模型。",
        ),
        SlideSpec(
            title="Probe 脚手架的总体位置",
            bullets=[
                "probe 脚手架接在 finetune 任务之后，统一调用 `render_raw_compare`、ghost scorer、point-support 和汇总脚本。",
                "one-step、stepcurve、mask-boost 共用 compare / metrics 这层壳，只是任务组织方式不同。",
                "所以这套分支的价值是工程自动化，而不是模型结构创新。",
            ],
            code_refs=[
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
            ],
            image_paths=existing_images(assets.get("compare_probe_framework")),
            image_caption="one-step / stepcurve / mask-boost 共用同一层 compare 与 metrics 脚手架。",
            layout="full_image",
            speaker_notes="这页只讲结构位置，不讲任何数值。",
        ),
        SlideSpec(
            title="One-step 分支做了什么",
            bullets=[
                "`scripts/orig_vggt_one_step_probe/*` 负责把“只训练一步”这件事封装成可重复执行的任务。",
                "这套脚本会准备任务、调用训练入口、触发 compare、再把结果汇总成统一格式的 JSON/MD。",
                "从代码职责上看，它解决的是最小 horizon 的自动化编排问题。",
            ],
            formulas=[
                "one-step = prepare task -> run task -> compare -> summarize",
            ],
            code_refs=[
                "scripts/orig_vggt_one_step_probe/common.py",
                "scripts/orig_vggt_one_step_probe/run_task.py",
                "scripts/orig_vggt_one_step_probe/summarize_runs.py",
            ],
            speaker_notes="把 one-step 讲成最小任务单元即可，不需要讲效果。",
        ),
        SlideSpec(
            title="Stepcurve 分支做了什么",
            bullets=[
                "`scripts/orig_vggt_stepcurve_probe/*` 负责把多个训练 horizon 串成同一套任务清单，并统一汇总。",
                "`audit_prefix.py` 用来检查 prefix 一致性，保证短 horizon 与长 horizon 的任务顺序可以对齐比较。",
                "这套脚手架解决的是多 horizon 任务组织与审计问题。",
            ],
            formulas=[
                "stepcurve = horizon list + prefix audit + compare + summarize",
            ],
            code_refs=[
                "scripts/orig_vggt_stepcurve_probe/audit_prefix.py",
                "scripts/orig_vggt_stepcurve_probe/run_task.py",
                "scripts/orig_vggt_stepcurve_probe/summarize_runs.py",
            ],
            speaker_notes="这里重点讲工程设计，不讲曲线结果。",
        ),
        SlideSpec(
            title="这套脚手架最终产出什么",
            bullets=[
                "probe 脚本最终统一产出 compare 目录、summary JSON、summary MD 以及可供 PPT 引用的汇总材料。",
                "也就是说，这一支代码改动的核心贡献是“把训练分支组织成稳定可复用的任务管线”。",
                "这层能力本身独立于任何具体结果，可以单独向导师说明。",
            ],
            code_refs=[
                "logs/modal_phase5/reports/*.json",
                "logs/modal_phase5/reports/*.md",
                "tools/generate_orig_vggt_local_code_mods_presentation.py",
            ],
            speaker_notes="结尾只讲自动化壳层的职责。",
        ),
    ]


def build_mask_deck(data: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})

    return [
        SlideSpec(
            title="Mask-Boost / Ghost / Point-Support 代码分支",
            bullets=[
                "这套只讲前景增强、ghost 评分和 point-support 量化脚本这条代码分支。",
                "重点是说明这些模块在代码里负责什么、公式是什么、它们与原版 VGGT 的关系是什么。",
                "换句话说，这是一套“训练语义与评估脚手架”的代码讲解。",
            ],
            speaker_notes="先把边界压住：只讲代码和原理，不讲表现。",
        ),
        SlideSpec(
            title="当前 mask-boost 的真实语义",
            bullets=[
                "这条代码路径的核心是先用前景 mask 做 valid gating，再对前景内部区域继续增权。",
                "它属于 `finetune_vggt_pseudo.py` 里的监督权重构造逻辑，而不是原版 backbone 的一部分。",
                "因此汇报时应把它讲成“训练权重语义改造”，而不是“模型结构改造”。",
            ],
            formulas=[
                "M_valid = M_valid_all * M_fg",
                "M_interior = erode(M_fg, r)",
                "w_final = w_base * [1 + (beta - 1) * M_interior]",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: foreground gating / boost path",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="mask-boost 在代码里体现为前景 gating 与 interior boost 的组合。",
            layout="full_image",
            speaker_notes="这一页是整套 mask-boost 解释的前提，语义要讲准。",
        ),
        SlideSpec(
            title="Ghost scorer 是怎么接进来的",
            bullets=[
                "`tools/score_ghosting_from_cat_pred.py` 负责把 compare 目录里的 ghost triplet 图转换成统一评分。",
                "它本质上是一个后处理评分脚本，位置在原版主线之外，但服务于整条本地代码线。",
                "因此这部分属于“新增评估壳层”，不是原版 VGGT 内部的预测头。",
            ],
            formulas=[
                "width_ratio = W_pred / W_mask",
                "area_ratio = A_pred / A_mask",
                "ghost_visual_score = ghost_score + dark_penalty + collapse_penalty",
            ],
            code_refs=[
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
            ],
            speaker_notes="这里只讲 ghost scorer 的职责与公式，不讲分数高低。",
        ),
        SlideSpec(
            title="Point-support 与 masked error 指标是怎么接进来的",
            bullets=[
                "`measure_point_support.py` 负责从 support 图和 mask 中提取结构性指标。",
                "同时这层脚本也把 masked PSNR / masked L1 统一进同一套汇总口径。",
                "这部分依然属于新增分析壳层，用来把 compare 输出整理成结构化指标。",
            ],
            formulas=[
                "subject_support_share = sum(S * M) / sum(S)",
                "largest_component_share = max_c sum_{x in c} S / sum_{x in active} S",
                "subject_psnr = -10 log10(sum((pred - tgt)^2 * M) / (3 * sum(M)))",
                "subject_l1 = sum(|pred - tgt| * M) / (3 * sum(M))",
            ],
            code_refs=[
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
            ],
            speaker_notes="这页只说明指标脚本的来源与公式。",
        ),
        SlideSpec(
            title="这一分支该怎么汇报",
            bullets=[
                "相对原版 VGGT，这一分支新增了三类代码: 前景权重语义、ghost 评分脚本、point-support / masked metrics 脚本。",
                "它们共同组成一层围绕原版主线的训练与评估壳。",
                "因此向导师汇报时，应讲成“代码如何组织与定义指标”，而不是“结果如何”。",
            ],
            image_paths=existing_images(assets.get("compare_probe_framework")),
            image_caption="mask-boost、ghost 和 point-support 都接在同一层 compare / probe 框架上。",
            layout="full_image",
            speaker_notes="最后把这条代码分支重新放回整体架构里。",
        ),
    ]


def build_visual_prompts_md() -> str:
    return build_visual_prompts_md_v3()

    sections = [
        (
            "图 1. 当前本地架构图",
            "current_local_vggt_architecture.png",
            "用于总版和“改动与微调版”的核心架构页。要求风格参考官方 VGGT 架构图，但内容必须对应当前本地代码。",
            """Create a clean academic architecture diagram in 16:9 horizontal layout, warm white background, vector style, research-paper figure quality.

IMPORTANT COMPOSITION RULE:
This must NOT look like the official native VGGT architecture figure.
The main subject of the image must be OUR LOCAL CODE WORK built around original VGGT.
Original VGGT should appear only as a small left-side foundation block, around 20 percent of the visual weight.
The added local modules must dominate the figure, around 80 percent of the visual attention.

Title concept:
Current local original-VGGT code architecture

Required layout from left to right:

Block A on the far left, small:
- Original VGGT backbone
- short sublabels only: aggregator, camera/depth/point heads
- visually compact, clearly secondary

Then a large central-right pipeline occupying most of the canvas:
1. VGGTGeomTeacher
2. Geometry unprojection
3. Precompute NPZ contract
4. PseudoGeomDataset
5. finetune_vggt_pseudo wrapper
6. Loss routing
7. FG mask boost
8. MV-depth reprojection
9. render_raw_compare
10. ghost scorer
11. point-support metrics
12. automation probes
13. PPT / JSON / MD outputs

Arrows:
- Original VGGT backbone feeds VGGTGeomTeacher
- teacher feeds Geometry unprojection
- unprojection feeds Precompute NPZ contract
- NPZ contract feeds PseudoGeomDataset
- dataset feeds finetune_vggt_pseudo wrapper
- wrapper branches into Loss routing, FG mask boost, MV-depth reprojection
- wrapper and compare pipeline feed render_raw_compare
- render_raw_compare feeds ghost scorer and point-support metrics
- ghost scorer and point-support metrics feed automation probes
- automation probes feed PPT / JSON / MD outputs

Visual hierarchy:
- biggest boxes and strongest emphasis should be:
  VGGTGeomTeacher
  Precompute NPZ contract
  finetune_vggt_pseudo wrapper
  render_raw_compare
  automation probes
- Original VGGT backbone must be visibly smaller and less important

Visual style:
- soft orange only for Original VGGT backbone and geometry-origin blocks
- soft beige for contract / dataset blocks
- soft green for training and loss blocks
- soft blue for compare / metrics / outputs
- soft purple only for automation probe blocks
- rounded rectangles, thin gray arrows, lots of white space, paper-figure feeling
- no icons, no screenshots, no photorealism

This figure should look like a serious paper diagram explaining a software-and-training pipeline, not a native model architecture figure.
The viewer should immediately conclude:
the original VGGT is only the starting point, and most of the work is the added local code shell.""",
            "Do not replicate the official VGGT figure layout. No three-row stacked token columns copied from the official figure. No large central dashed attention box copied from the official figure. No official VGGT native composition. No photo collage, no 3D render, no dark background, no neon colors, no decorative icons, no excessive gradients, no cartoon style, no crowded text, no benchmark numbers, no result charts.",
        ),
        (
            "图 2. 代码改动分层图",
            "code_touch_map.png",
            "用于说明哪些是直接改原版文件，哪些是新增外层代码壳。",
            """Create a clean 16:9 vector diagram on a warm white background showing code changes relative to the original VGGT repository.

Layout:
- left column titled Original repo: directly touched
- center small box titled Original VGGT core
- right column titled Local repo: newly added shell

Left column should contain 4 rounded boxes:
- vggt/utils/geometry.py
- vggt/utils/load_fn.py
- demo_colmap.py
- vggt/models/vggt.py

Right column should contain 6 rounded boxes:
- vggt_geom.py
- precompute_zju_vggt_geom.py
- finetune_vggt_pseudo.py
- render_raw_compare.py
- score_ghosting_from_cat_pred.py
- probe scripts

Arrows:
- all left boxes point into Original VGGT core
- Original VGGT core points out to all right boxes

Style:
- left boxes: light orange
- center box: beige / neutral
- right boxes: light green
- simple thin gray arrows
- academic, tidy, easy to read

Goal:
make it visually obvious that only a few original files were directly changed, while most of the work was added as a local outer shell.""",
            "No data plots, no result numbers, no screenshots, no code text paragraphs, no flashy infographic style, no dark theme.",
        ),
        (
            "图 3. Geometry teacher 与 NPZ 预计算流程图",
            "teacher_npz_flow.png",
            "用于讲解 teacher、反投影公式和 NPZ 数据契约。",
            """Create a clean research-style flow diagram in 16:9 landscape, warm white background, vector figure.

Show a left-to-right pipeline with four large blocks:
1. Original VGGT Outputs
   labels inside: Depth / Conf / Camera
2. VGGTGeomTeacher
3. Unproject / frame meta / point map
4. NPZ contract
   labels inside: images, depth, depth_conf, pointmap, intrinsic, extrinsic, meta

Arrows go from left to right across all four blocks.

Under the blocks, place the core formulas in a clean centered formula area:
x_cam = (u + o - cx) * z / fx
y_cam = (v + o - cy) * z / fy
P_world = R^T * ([x_cam, y_cam, z]^T - t)

Visual style:
- Original VGGT block in soft orange
- VGGTGeomTeacher in beige
- Unproject / point map block in soft green
- NPZ contract block in soft blue
- thin gray arrows, rounded rectangles, clean paper figure look

The message should visually communicate:
original outputs are reorganized into a geometry teacher, then fixed into a reusable NPZ data contract.""",
            "No photos, no benchmark numbers, no result curves, no 3D objects, no dark mode, no crowded annotation cloud.",
        ),
        (
            "图 4. Loss routing 图",
            "loss_routing.png",
            "用于讲解 finetune_vggt_pseudo 如何把监督统一路由起来。",
            """Create a 16:9 research paper style diagram on a warm white background, vector style.

Layout from left to right:
- left box: PseudoGeomDataset (images + NPZ teacher)
- middle large box: finetune_vggt_pseudo (wrapper / train loop)
  small internal labels: conf weight, loss routing, mask gating
- right side: multiple small target boxes arranged in two columns:
  Depth
  Point
  Point reproj
  Normal
  MV-depth
  Camera
  FG boost
  Conf
- far right a box labeled L_total

Arrows:
- dataset box into wrapper
- wrapper fans out to all target boxes
- all target boxes conceptually merge into L_total

At the bottom include one compact formula area:
L_total = sum_k lambda_k L_k
w_conf(x) = valid(x) * clip((c(x)-t)/(1-t), 0, 1)^gamma

Color coding:
- dataset beige
- wrapper light green
- target boxes mixed soft orange / green / blue / purple but still restrained
- L_total light blue

The figure should clearly communicate:
the wrapper routes many supervision terms without changing the original backbone.""",
            "No charts, no screenshots, no photo textures, no neon colors, no corporate infographic style, no unnecessary icons.",
        ),
        (
            "图 5. FG boost 语义图",
            "fg_boost_semantics.png",
            "用于讲解 foreground gating 与 interior boost 的代码语义。",
            """Create a minimal, clean, academic diagram in 16:9 landscape with warm white background, vector style.

Show four rounded boxes in a left-to-right row:
1. GT FG mask
2. Erode interior
3. Valid-mask gating
4. Final weight

Below or near each box place a short formula label:
- M_fg
- M_interior = erode(M_fg, r)
- M_valid = M_valid_all * M_fg
- w_final

At the bottom center place a clean formula line:
w_final = w_base * [1 + (beta - 1) * M_interior]

Style:
- box 1 soft orange
- box 2 beige
- box 3 soft green
- box 4 soft blue
- thin gray arrows between boxes
- lots of white space
- paper-figure aesthetics, highly legible

The image must feel like a code-semantics explanation figure, not an experiment result figure.""",
            "No people, no masks from real photos, no segmentation screenshots, no charts, no numbers from experiments, no dark background.",
        ),
        (
            "图 6. Compare / Metrics / Probe 脚手架图",
            "compare_probe_framework.png",
            "用于讲解 compare、ghost scorer、point-support 和自动化 probe 脚本的关系。",
            """Create a 16:9 clean research workflow diagram on a warm white background, vector figure style.

Layout from left to right:
- box 1: model.pt or checkpoint
- box 2: Finetune task
- box 3: render_raw_compare
- on the right three stacked boxes:
  Ghost scoring
  Point-support
  Automation scripts

Arrows:
- checkpoint -> Finetune task -> render_raw_compare
- render_raw_compare branches to Ghost scoring, Point-support, Automation scripts

At the bottom near Automation scripts, add small labels:
one-step
stepcurve
mask-boost

Visual style:
- checkpoint beige
- finetune soft green
- render_raw_compare soft blue
- Ghost scoring soft orange
- Point-support soft green
- Automation scripts soft purple
- very clean rounded boxes and thin gray arrows

The message should be:
compare / metrics / probe are an engineering shell around the same original-VGGT line, not a new model family.""",
            "No benchmark plots, no screenshots of rendered results, no dark theme, no excessive gradient, no dashboard style.",
        ),
    ]

    lines = [
        "# 原版 VGGT PPT 插图生成 Prompt",
        "",
        "## 总原则",
        "",
        "- 所有 prompt 都必须围绕“你相对原版 VGGT 所作的工作”来写。",
        "- 图的主题必须是新增代码壳、训练壳、分析壳、自动化壳，而不是解释原生 VGGT。",
        "- 原版 VGGT 在图里只能作为起点、底座、对照对象，不能占主画面。",
        "- 如果图看起来像官方架构图复刻，说明 prompt 失败了。",
        "",
        "## 统一写法要求",
        "",
        "- prompt 里必须优先写 `relative to the original VGGT` 或 `built around original VGGT`。",
        "- prompt 里必须明确写 `original VGGT is a small foundation block`。",
        "- prompt 里必须明确写 `the added local modules dominate the figure`。",
        "- prompt 里必须禁止 `replicate the official VGGT layout`。",
        "",
        "说明：下面这些 prompt 对应当前 PPT 中使用的概念图，全部只讲代码结构与模块职责，不展开任何结果判断。",
        "建议统一使用 `16:9`、`research paper figure`、`clean vector diagram` 风格。",
        "如果图片模型文字质量不稳定，建议先生成“留白框图版”，再在 PPT 里覆盖文字。",
        "",
    ]
    for idx, (title, filename, usage, prompt, negative) in enumerate(sections, start=1):
        lines.extend(
            [
                f"## {idx}. {title}",
                "",
                f"- 建议文件名: `{filename}`",
                f"- 用途: {usage}",
                "",
                "### 正向 Prompt",
                "",
                "```text",
                prompt.strip(),
                "```",
                "",
                "### Negative Prompt",
                "",
                "```text",
                negative.strip(),
                "```",
                "",
            ]
        )
    return "\n".join(lines).strip() + "\n"


def build_visual_prompts_md_v2() -> str:
    sections = [
        {
            "title": "图 1. 当前本地架构图",
            "filename": "current_local_vggt_architecture.png",
            "usage": "用于总版和“改动与微调版”的核心架构页。风格只参考官方图的整洁感，但内容必须直接体现你本地代码相对原版 VGGT 的具体改造。",
            "work_items": [
                "原版 VGGT 只能作为左侧小底座出现，只保留 aggregator、camera head、depth / point heads 这类核心标签，不能占主画面。",
                "要把你在原版树里直接动过的入口缩成一组小标注：`vggt/utils/geometry.py` 的 `unproject_impl / pixel_center_offset`，`vggt/utils/load_fn.py` 的 threaded image loading，`demo_colmap.py` 的 pycolmap 注册与 BA 兼容。",
                "主画面必须是你新增的外层代码壳：`vggt_geom.py` 的 `VGGTGeomTeacher`，`precompute_zju_vggt_geom.py` 的 NPZ 预计算契约，`finetune_vggt_pseudo.py` 的 dataset + wrapper。",
                "在 wrapper 内部必须明确画出 `Loss routing`、`FG mask boost`、`MV-depth reprojection` 三个训练语义模块。",
                "训练后面的分析层必须明确画出 `render_raw_compare.py`、`score_ghosting_from_cat_pred.py`、`measure_point_support.py`，以及 `one-step / stepcurve / mask-boost` 三套 probe 自动化脚本。",
                "整张图传达的信息必须是：原版核心基本保留，你的主要工作是把 teacher、NPZ、wrapper、compare、metrics、probe 这些外层流程补齐。",
            ],
            "labels": [
                "Original VGGT core",
                "geometry.py: unproject_impl / pixel_center_offset",
                "load_fn.py: threaded image loading",
                "demo_colmap.py: pycolmap / BA compatibility",
                "VGGTGeomTeacher",
                "Precompute NPZ contract",
                "PseudoGeomDataset",
                "finetune_vggt_pseudo.py",
                "Loss routing",
                "FG mask boost",
                "MV-depth reprojection",
                "render_raw_compare.py",
                "ghost scorer",
                "point-support metrics",
                "one-step / stepcurve / mask-boost probes",
                "JSON / MD / PPT outputs",
            ],
            "prompt": """Create a clean 16:9 academic architecture diagram, warm white background, restrained vector style, paper-figure quality.

This figure is NOT a generic explanation of native VGGT. It must directly visualize MY ACTUAL LOCAL CODE MODIFICATIONS relative to the original VGGT repository.

Composition rule:
- original VGGT is only a small left-side foundation block, about 15 to 20 percent of the visual weight
- the added local modules dominate the figure, about 80 to 85 percent of the visual attention
- do not use the official VGGT token-column layout; use a software-pipeline / system-map composition instead

The figure must explicitly show these concrete code modifications:
- a small block called Original VGGT core with tiny sublabels aggregator, camera head, depth head, point head
- a small callout cluster for directly touched original files:
  - vggt/utils/geometry.py: unproject_impl, pixel_center_offset
  - vggt/utils/load_fn.py: threaded image loading
  - demo_colmap.py: pycolmap registration, bundle adjustment compatibility
- a dominant main pipeline for the added local shell:
  - vggt_geom.py -> VGGTGeomTeacher
  - geometry unprojection / point-map construction
  - precompute_zju_vggt_geom.py -> reusable NPZ contract
  - PseudoGeomDataset
  - finetune_vggt_pseudo.py -> main finetune wrapper
  - inside the wrapper, three visible internal branches: Loss routing, FG mask boost, MV-depth reprojection
  - scripts/orig_vggt_viewcount/render_raw_compare.py
  - tools/score_ghosting_from_cat_pred.py
  - scripts/orig_vggt_stepcurve_probe/measure_point_support.py
  - scripts/orig_vggt_one_step_probe/*, scripts/orig_vggt_stepcurve_probe/*, scripts/orig_vggt_mask_boost_probe/*
  - final outputs: compare folders, summary JSON, summary MD, PPT materials

Arrow logic:
- Original VGGT core feeds VGGTGeomTeacher
- geometry.py callout is visually attached to the geometry conversion stage
- teacher feeds unprojection and point-map construction
- point-map construction feeds NPZ contract
- NPZ contract feeds PseudoGeomDataset
- dataset feeds finetune_vggt_pseudo wrapper
- wrapper branches to Loss routing, FG mask boost, MV-depth reprojection
- wrapper then feeds compare / metrics / probe layers
- compare / metrics / probe feed report outputs

Add 2 or 3 small formula callouts near the relevant blocks:
- P_world = R^T * ([x_cam, y_cam, z]^T - t)
- L_total = sum_k lambda_k L_k
- w_final = w_base * [1 + (beta - 1) * M_interior]

Visual style:
- soft orange for original-core and geometry-origin blocks
- beige for dataset / contract blocks
- soft green for training wrapper blocks
- soft blue for compare / metric blocks
- soft purple for probe automation blocks
- rounded rectangles, thin gray arrows, generous whitespace, no icons, no photorealism

The viewer should immediately understand:
the original VGGT is only the starting point, while most of the actual work is the added local code shell around it.""",
            "negative": """Do not replicate the official VGGT architecture layout. No token-column composition, no large central attention box, no official figure with relabeled text. No screenshots, no benchmark curves, no result charts, no photo collage, no dark theme, no neon colors, no decorative icons, no dense dashboard style.""",
        },
        {
            "title": "图 2. 代码改动分层图",
            "filename": "code_touch_map.png",
            "usage": "用于说明哪些地方是直接改原版文件，哪些地方是围绕原版新增的外层代码壳。",
            "work_items": [
                "左侧必须是“直接改到原版仓库”的少数入口，而且每个入口都要带一句具体改了什么。",
                "`vggt/utils/geometry.py` 不能只写文件名，必须写成 `unproject_impl / pixel_center_offset` 这类具体改动。",
                "`vggt/utils/load_fn.py` 必须体现 threaded image loading；`demo_colmap.py` 必须体现 pycolmap 注册、debug 输出、BA 兼容；`vggt/models/vggt.py` 只作为结构说明级补充。",
                "右侧必须是你新增的外层壳层，并明确每个文件负责什么：teacher、NPZ、wrapper、compare、ghost scorer、probe 脚手架。",
                "整张图要让导师一眼看懂：原版内部真正改动很少，主体工作在外围新模块。",
            ],
            "labels": [
                "Original repo: directly touched",
                "vggt/utils/geometry.py: unproject_impl / pixel_center_offset",
                "vggt/utils/load_fn.py: threaded image loading",
                "demo_colmap.py: pycolmap / BA compatibility",
                "vggt/models/vggt.py: structure clarification",
                "Original VGGT core",
                "vggt_geom.py: geometry teacher",
                "precompute_zju_vggt_geom.py: NPZ precompute",
                "finetune_vggt_pseudo.py: dataset + loss routing",
                "render_raw_compare.py: compare renderer",
                "score_ghosting_from_cat_pred.py: ghost scoring",
                "probe scripts: one-step / stepcurve / mask-boost",
            ],
            "prompt": """Create a clean 16:9 vector diagram on a warm white background showing WHICH PARTS OF MY LOCAL REPO DIRECTLY MODIFY THE ORIGINAL VGGT TREE and WHICH PARTS ARE NEW OUTER-SHELL MODULES.

This is not a generic repository map. It must explicitly show the concrete file-level work.

Layout:
- left column: Original repo, directly touched
- center small block: Original VGGT core
- right column: Local repo, newly added shell

Left column must contain these concrete file boxes with one-line responsibility text:
- vggt/utils/geometry.py: unproject_impl, pixel_center_offset
- vggt/utils/load_fn.py: threaded image loading
- demo_colmap.py: pycolmap registration, debug output, BA compatibility
- vggt/models/vggt.py: structure clarification, not a backbone rewrite

Optional small auxiliary strip at the bottom for light support changes:
- training/config/default.yaml
- training/config/default_dataset.yaml
- training/data/dynamic_dataloader.py

Right column must contain these concrete new-shell modules with one-line responsibility text:
- vggt_geom.py: VGGTGeomTeacher
- precompute_zju_vggt_geom.py: NPZ geometry contract
- finetune_vggt_pseudo.py: PseudoGeomDataset, loss routing, FG boost, MV-depth
- scripts/orig_vggt_viewcount/render_raw_compare.py: unified compare output
- tools/score_ghosting_from_cat_pred.py: ghost metric
- scripts/orig_vggt_one_step_probe/*, scripts/orig_vggt_stepcurve_probe/*, scripts/orig_vggt_mask_boost_probe/*: automation shells

Arrow logic:
- left-column boxes point into Original VGGT core
- Original VGGT core points toward the right-column outer-shell modules
- visually emphasize that the left side is small and the right side is the main added work

Style:
- left boxes soft orange
- center block neutral beige
- right boxes soft green / blue
- thin gray arrows, tidy academic layout, easy to scan

The diagram should make one idea obvious:
only a few original files were directly touched, while most of the work was added as an outer local shell around original VGGT.""",
            "negative": """No result numbers, no screenshots, no code paragraph screenshots, no dashboard style, no dark theme, no busy infographic decoration, no official VGGT architecture composition.""",
        },
    ]

    sections.extend(
        [
            {
                "title": "图 3. Geometry teacher 与 NPZ 预计算流程图",
                "filename": "teacher_npz_flow.png",
                "usage": "用于讲清楚 `vggt_geom.py`、`vggt/utils/geometry.py`、`precompute_zju_vggt_geom.py` 这一层到底在做什么。",
                "work_items": [
                    "必须明确：你不是简单拿原版输出直接训，而是先把原版的 depth / conf / camera 输出整理成 geometry teacher。",
                    "必须明确写出 `vggt/utils/geometry.py` 里的 unprojection 约定，因为这一层是 pointmap / 世界坐标监督的来源。",
                    "必须明确写出 `precompute_zju_vggt_geom.py` 生成的是固定 NPZ 契约，里面要有 `images / depth / depth_conf / pointmap / intrinsic / extrinsic / meta` 这类键。",
                    "图里要体现“原版输出 -> teacher -> 几何反投影 -> NPZ 契约 -> 后续训练读取”这条链，而不是只画一个抽象框。",
                ],
                "labels": [
                    "Original VGGT outputs: depth / conf / camera",
                    "vggt_geom.py: VGGTGeomTeacher",
                    "vggt/utils/geometry.py: unproject_impl / pixel_center_offset",
                    "point map / frame meta / camera supervision",
                    "precompute_zju_vggt_geom.py",
                    "NPZ: images / depth / depth_conf / pointmap / intrinsic / extrinsic / meta",
                    "modal_run_train.py",
                    "P_world = R^T * ([x_cam, y_cam, z]^T - t)",
                ],
                "prompt": """Create a clean 16:9 research-style flow diagram that directly explains my geometry-teacher and NPZ-precompute work relative to original VGGT.

This figure must visualize the actual local code path, not a generic depth pipeline.

Required pipeline from left to right:
1. Original VGGT outputs
   show small labels: depth, conf, camera
2. vggt_geom.py
   main label: VGGTGeomTeacher
   meaning: reorganize original outputs into reusable geometry supervision
3. vggt/utils/geometry.py
   main label: unproject_impl / pixel_center_offset
   meaning: define the 3D unprojection convention and build point maps / world coordinates
4. geometry products
   show labels: pointmap, frame meta, intrinsic, extrinsic
5. precompute_zju_vggt_geom.py
   meaning: write a reusable NPZ contract
6. NPZ contract
   show labels: images, depth, depth_conf, pointmap, intrinsic, extrinsic, meta
7. downstream read path
   small label: finetune_vggt_pseudo.py / PseudoGeomDataset

Arrow logic:
- original VGGT outputs feed VGGTGeomTeacher
- teacher feeds geometry unprojection
- geometry unprojection feeds point-map / frame-meta products
- those products feed precompute_zju_vggt_geom.py
- precompute writes the NPZ contract
- NPZ contract feeds downstream training

Include these formulas in a clean formula area:
- x_cam = (u + o - cx) * z / fx
- y_cam = (v + o - cy) * z / fy
- P_world = R^T * ([x_cam, y_cam, z]^T - t)

Use a clean paper-figure look:
- original output block soft orange
- teacher block beige
- geometry / pointmap block soft green
- NPZ contract block soft blue
- thin gray arrows, rounded rectangles, generous whitespace

The message should be unmistakable:
I first convert original VGGT outputs into a geometry teacher, then freeze them into an NPZ data contract for later training.""",
                "negative": """No result curves, no 3D rendered scenes, no screenshots, no benchmark tables, no dark mode, no decorative clutter, no effect-oriented wording.""",
            },
            {
                "title": "图 4. Loss routing 图",
                "filename": "loss_routing.png",
                "usage": "用于讲解 `finetune_vggt_pseudo.py` 如何把 teacher、dataset、conf 权重和多种监督统一接成一个 wrapper。",
                "work_items": [
                    "必须把 `PseudoGeomDataset` 画出来，因为这层代表你不是直接喂原版输出，而是读固定 teacher NPZ。",
                    "必须把 `finetune_vggt_pseudo.py` 画成主控 wrapper，而不是普通小方框。",
                    "必须在 wrapper 内部体现 `conf weighting`、`loss routing`、`mask gating`、`MV-depth reprojection` 这些训练语义。",
                    "右侧监督项不能只写 generic loss，至少要落到 `Depth / Point / Point reproj / Normal / Camera / MV-depth` 这些实际分支。",
                    "底部要给出 `L_total` 和 `w_conf(x)` 两个核心公式，让图直接对上代码逻辑。",
                ],
                "labels": [
                    "PseudoGeomDataset",
                    "NPZ teacher + images",
                    "finetune_vggt_pseudo.py",
                    "conf weighting",
                    "loss routing",
                    "mask gating",
                    "MV-depth reprojection",
                    "Depth / Point / Point reproj / Normal / Camera / MV-depth",
                    "FG boost",
                    "L_total = sum_k lambda_k L_k",
                    "w_conf(x) = valid(x) * clip((c(x)-t)/(1-t), 0, 1)^gamma",
                ],
                "prompt": """Create a clean 16:9 research-paper diagram that explains the actual training wrapper implemented in finetune_vggt_pseudo.py.

This is not a generic multi-loss diagram. It must reflect my specific local code structure.

Layout:
- left block: PseudoGeomDataset
  small label: images + NPZ teacher
- center dominant block: finetune_vggt_pseudo.py
  this should be the largest block in the figure
  show internal sections or inner labels:
  - conf weighting
  - loss routing
  - mask gating
  - MV-depth reprojection
- right side: a fan-out of concrete supervision branches
  - Depth
  - Point
  - Point reproj
  - Normal
  - Camera
  - MV-depth
- a small side branch attached to the wrapper: FG boost
- far right final merge block: L_total

Arrow logic:
- PseudoGeomDataset feeds finetune_vggt_pseudo.py
- the wrapper fans out to all concrete supervision branches
- FG boost modulates the wrapper weight path, not the original backbone
- all supervision branches conceptually merge into L_total

At the bottom include these formulas exactly:
- L_total = sum_k lambda_k L_k
- w_conf(x) = valid(x) * clip((c(x)-t)/(1-t), 0, 1)^gamma

Visual hierarchy:
- dataset block beige
- wrapper block soft green and visually dominant
- supervision branches mixed restrained orange / blue / green
- final L_total block soft blue
- rounded rectangles, thin gray arrows, clean paper layout

The figure must communicate:
my main training work is the wrapper that routes teacher-based supervision, not a rewrite of the original VGGT backbone.""",
                "negative": """No result charts, no screenshots, no flashy infographic style, no neon colors, no dark background, no generic AI-pipeline icons, no benchmark numbers.""",
            },
        ]
    )

    sections.extend(
        [
            {
                "title": "图 5. FG boost 语义图",
                "filename": "fg_boost_semantics.png",
                "usage": "用于把当前代码里的 foreground gating + interior boost 语义讲准，避免被误解成别的训练策略。",
                "work_items": [
                    "必须体现这不是网络结构，而是 `finetune_vggt_pseudo.py` 里的权重构造逻辑。",
                    "必须体现顺序：先有 `M_valid_all`，再乘 `M_fg` 形成 `M_valid`，再对 `M_fg` 做 erosion 得到 `M_interior`，最后只对 interior 区域增权。",
                    "不能只画抽象 attention 热力图，必须画成离散 mask / weight 语义流程。",
                    "公式必须落到 `M_valid = M_valid_all * M_fg`、`M_interior = erode(M_fg, r)`、`w_final = w_base * [1 + (beta - 1) * M_interior]`。",
                ],
                "labels": [
                    "finetune_vggt_pseudo.py: foreground gating / boost path",
                    "M_valid_all",
                    "GT foreground mask M_fg",
                    "M_valid = M_valid_all * M_fg",
                    "M_interior = erode(M_fg, r)",
                    "w_base",
                    "w_final = w_base * [1 + (beta - 1) * M_interior]",
                    "foreground gating",
                    "interior boost",
                ],
                "prompt": """Create a minimal but very explicit 16:9 academic diagram that explains the ACTUAL FG mask weight semantics implemented in finetune_vggt_pseudo.py.

This is not a generic attention map and not a result figure. It must visualize the current code logic.

Required visual logic from left to right:
1. valid geometry mask
   label: M_valid_all
2. GT foreground mask
   label: M_fg
3. combine them
   label: M_valid = M_valid_all * M_fg
4. erode the foreground interior
   label: M_interior = erode(M_fg, r)
5. build final supervision weight
   labels: w_base, w_final

The figure must clearly separate two ideas:
- foreground gating
- interior boost

At the bottom show this final formula exactly:
- w_final = w_base * [1 + (beta - 1) * M_interior]

Style:
- abstract binary-mask / weight-map boxes, not real photos
- soft orange for foreground-mask related boxes
- beige for erosion / morphology step
- soft green for gating step
- soft blue for final weight step
- thin gray arrows, generous whitespace, highly legible paper-figure style

Also include a small caption-like label inside the figure:
finetune_vggt_pseudo.py: foreground gating / boost path

The viewer should understand:
this branch changes supervision weights inside training, not the original model structure.""",
                "negative": """No human figures, no real segmentation screenshots, no attention heatmaps from papers, no experiment numbers, no result plots, no dark theme, no photorealistic masks.""",
            },
            {
                "title": "图 6. Compare / Metrics / Probe 脚手架图",
                "filename": "compare_probe_framework.png",
                "usage": "用于说明 compare、ghost 评分、point-support 脚本和三套 probe 自动化脚本是如何串起来的。",
                "work_items": [
                    "必须把 `render_raw_compare.py` 画成统一 compare 入口，因为后面的 ghost / point-support 都依赖这层输出。",
                    "必须把 `score_ghosting_from_cat_pred.py` 和 `measure_point_support.py` 单独画出来，说明它们是后处理分析脚本，不是原版预测头。",
                    "必须把 `one-step / stepcurve / mask-boost` 三套脚手架画成自动化任务组织层，而不是三种新模型。",
                    "输出端要明确写 `compare/`、`summary.json`、`summary.md`、`PPT materials` 这类报告材料。",
                    "整张图传达的信息必须是：你新增了一整层围绕原版主线的 compare / metrics / probe 工程壳。",
                ],
                "labels": [
                    "checkpoint / model.pt",
                    "finetune_vggt_pseudo.py task",
                    "scripts/orig_vggt_viewcount/render_raw_compare.py",
                    "tools/score_ghosting_from_cat_pred.py",
                    "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
                    "orig_vggt_one_step_probe",
                    "orig_vggt_stepcurve_probe",
                    "orig_vggt_mask_boost_probe",
                    "compare/",
                    "summary.json / summary.md / PPT materials",
                ],
                "prompt": """Create a clean 16:9 research workflow diagram that explains the actual compare / metrics / probe shell added around the original-VGGT training line.

This is not a model-architecture figure. It must explicitly show the local engineering scripts and what they do.

Required layout from left to right:
- checkpoint / model.pt
- finetune_vggt_pseudo.py task
- scripts/orig_vggt_viewcount/render_raw_compare.py
- two explicit analysis boxes branching from render_raw_compare:
  - tools/score_ghosting_from_cat_pred.py
  - scripts/orig_vggt_stepcurve_probe/measure_point_support.py
- a probe automation layer with three concrete boxes:
  - orig_vggt_one_step_probe
  - orig_vggt_stepcurve_probe
  - orig_vggt_mask_boost_probe
- final output boxes:
  - compare/
  - summary.json
  - summary.md
  - PPT materials

Arrow logic:
- checkpoint feeds finetune task
- finetune task feeds render_raw_compare
- render_raw_compare feeds both ghost scoring and point-support measurement
- ghost scoring and point-support feed the three probe automation boxes
- probe automation boxes feed the final report outputs

Optional small formula / metric labels near the analysis boxes:
- ghost_visual_score
- subject_support_share
- largest_component_share
- subject_psnr / subject_l1

Visual style:
- checkpoint beige
- finetune block soft green
- render_raw_compare soft blue
- ghost scoring soft orange
- point-support soft green
- probe automation soft purple
- report outputs neutral blue-gray
- rounded rectangles, thin gray arrows, no screenshots, generous whitespace

The diagram should make one idea clear:
compare, metrics, and probes are a new engineering shell around the same original-VGGT line, not a new model family.""",
                "negative": """No rendered image screenshots, no benchmark plots, no dashboard collage, no dark theme, no heavy gradients, no result-oriented wording, no official VGGT layout reuse.""",
            },
        ]
    )

    lines = [
        "# 原版 VGGT PPT 插图生成 Prompt",
        "",
        "## 总原则",
        "",
        "- 这份文档不是在写几个空泛关键词，而是把“你相对原版 VGGT 具体做了什么”直接写进每张图的提示词里。",
        "- 每张图都要先体现代码工作，再谈画风；风格词只能作为辅助，不能取代具体模块、文件名、公式和箭头关系。",
        "- 原版 VGGT 在图里只能作为起点或底座，主画面必须是你新增的 teacher、NPZ、wrapper、compare、metrics、probe 这些外层工作。",
        "- 所有图都只讲结构、职责、公式、文件关系，不讲效果，不讲结果，不讲优劣。",
        "",
        "## 使用方式",
        "",
        "- 每一节都先给出“这张图必须体现的具体工作”，你先用它检查图的主题有没有跑偏。",
        "- 再看“图中文字建议”，如果图片模型文字不稳定，可以先生成框架图，再在 PPT 里手工覆盖这些标签。",
        "- 最后直接复制“可直接用于生图的详细 Prompt”去生成图片；这些 prompt 已经包含你的具体文件、模块和公式，不是只有正向/负向词汇。",
        "",
    ]
    for idx, section in enumerate(sections, start=1):
        lines.extend(
            [
                f"## {idx}. {section['title']}",
                "",
                f"- 建议文件名: `{section['filename']}`",
                f"- 用途: {section['usage']}",
                "",
                "### 这张图必须体现的具体工作",
                "",
            ]
        )
        for item in section["work_items"]:
            lines.append(f"- {item}")
        lines.extend(
            [
                "",
                "### 图中文字建议",
                "",
            ]
        )
        for item in section["labels"]:
            lines.append(f"- {item}")
        lines.extend(
            [
                "",
                "### 可直接用于生图的详细 Prompt",
                "",
                "```text",
                section["prompt"].strip(),
                "```",
                "",
                "### Negative Prompt",
                "",
                "```text",
                section["negative"].strip(),
                "```",
                "",
            ]
        )
    return "\n".join(lines).strip() + "\n"


def build_visual_prompts_md_v3() -> str:
    sections = [
        {
            "title": "图 1. 当前本地架构图",
            "filename": "current_local_vggt_architecture.png",
            "usage": "用于总版和“改动与微调版”的核心架构页。写法参考你给的官方图思路：用模块层、信息流、输出头来组织，而不是把文件名堆在图上。",
            "work_items": [
                "上半部分只保留原版核心主干，作为左到右的基础通路：多视图输入、token 表征、全局聚合、帧级聚合、相机头、稠密几何头。",
                "下半部分或外层虚线框必须明确表现你新增的本地壳层：几何教师、几何缓存、伪几何数据集、监督路由器、可视化对照、结构测量、自动化探针。",
                "图的重点不是‘原版网络长什么样’，而是‘原版输出之后，你又补了哪些层，才形成现在这条完整流程’。",
                "新增训练语义必须落成具体模块：置信度加权、前景权重、跨视图深度一致性、相机/法向/重投影监督。",
                "新增分析语义必须落成具体模块：统一可视化对照、重影测量、点支撑测量、自动化扫描与汇总输出。",
            ],
            "labels": [
                "Multi-view Inputs",
                "Patch Tokens",
                "Camera Token",
                "Global Aggregation",
                "Frame Aggregation",
                "Camera Head",
                "Dense Geometry Head",
                "Geometry Teacher",
                "Geometry Cache",
                "Pseudo-geometry Dataset",
                "Supervision Router",
                "Confidence Weighting",
                "Foreground Weighting",
                "Multiview Depth Consistency",
                "Visual Compare",
                "Ghost Measure",
                "Point-support Measure",
                "Probe Automation",
                "Structured Reports",
            ],
            "prompt": """Create a clean 16:9 academic architecture diagram, warm white background, vector style, paper figure quality.

Use the composition logic of a model architecture figure like the reference image: inputs on the left, token-like intermediate representation in the middle, task heads and outputs on the right. However, the content must describe the CURRENT LOCAL PIPELINE built around original VGGT, not the native VGGT figure itself.

Top lane:
- multi-view image inputs
- patch tokens and a small camera token
- global aggregation block
- frame aggregation block
- two output heads on the right: Camera Head and Dense Geometry Head

Bottom or outer-shell lane:
- Geometry Teacher
- Geometry Cache
- Pseudo-geometry Dataset
- Supervision Router
- three visible inner branches inside the router:
  Confidence Weighting
  Foreground Weighting
  Multiview Depth Consistency
- then Visual Compare
- Ghost Measure
- Point-support Measure
- Probe Automation
- Structured Reports

Composition rules:
- original VGGT core is only the upper foundation lane, visually smaller than the added outer shell
- the added local shell should take most of the visual attention
- use modular blocks, token stacks, arrows, and grouped regions
- use a dashed outer frame to indicate the newly added local training and analysis shell

Required message:
the original backbone remains the starting point, but the main local work is the added geometry-training-analysis shell after the native outputs.

Add 2 or 3 small formula callouts near the local shell:
- P_world = R^T * ([x_cam, y_cam, z]^T - t)
- L_total = sum_k lambda_k L_k
- w_final = w_base * [1 + (beta - 1) * M_interior]

Visual style:
- soft orange for original-core blocks
- beige for teacher / cache blocks
- soft green for supervision and training blocks
- soft blue for compare and report blocks
- soft purple for probe automation blocks
- lots of whitespace, thin gray arrows, no icons, no screenshots, no photo textures""",
            "negative": """Do not write source-code filenames anywhere in the figure. Do not replicate the official VGGT layout one to one. No copied token-column arrangement from the official figure. No screenshots, no benchmark curves, no result charts, no dark theme, no dashboard collage.""",
        },
        {
            "title": "图 2. 代码改动分层图",
            "filename": "code_touch_map.png",
            "usage": "用于把‘原版内部少量改动’和‘外围新增模块’分开讲清楚，但图上只用概念职责，不出现源码文件名。",
            "work_items": [
                "左侧只放少量原版内改动，而且要写成职责语言：几何约定修正、图像读取加速、COLMAP 兼容桥、核心结构说明补充。",
                "中间保留一个小的原版核心块，表示原版主干没有被推翻。",
                "右侧放新增外层模块：几何教师、可复用几何缓存、伪几何训练入口、统一可视化对照、重影测量、自动化探针。",
                "整张图的观感必须让导师一眼明白：原版内部改得少，真正新增的工作集中在外围流程壳层。",
            ],
            "labels": [
                "Original Core Adjustments",
                "Geometry Convention Update",
                "Image Loading Acceleration",
                "COLMAP Compatibility Bridge",
                "Core Structure Clarification",
                "Original VGGT Core",
                "Geometry Teacher",
                "Reusable Geometry Cache",
                "Pseudo-geometry Finetune Entry",
                "Unified Visual Compare",
                "Ghost Measure",
                "Probe Automation",
            ],
            "prompt": """Create a clean 16:9 layered system diagram on a warm white background.

The diagram should separate two groups of work:
- a small left group for the few changes inside the original core
- a large right group for the newly added outer shell

Left group labels must be concept-based, not filenames:
- Geometry Convention Update
- Image Loading Acceleration
- COLMAP Compatibility Bridge
- Core Structure Clarification

Middle:
- a small Original VGGT Core block

Right group labels:
- Geometry Teacher
- Reusable Geometry Cache
- Pseudo-geometry Finetune Entry
- Unified Visual Compare
- Ghost Measure
- Probe Automation

Arrow logic:
- left-group adjustments point into the Original VGGT Core
- the Original VGGT Core points toward the large outer-shell group

The visual message should be:
the original core is mostly preserved, while the main local work is added as a large outer shell around it.

Use clean rounded rectangles, thin arrows, academic figure style, balanced spacing, no decorative clutter.""",
            "negative": """Do not use source-code filenames, paths, or script names. No screenshots, no result numbers, no flashy infographic style, no dark theme, no copied official VGGT architecture layout.""",
        },
        {
            "title": "图 3. Geometry teacher 与几何缓存流程图",
            "filename": "teacher_npz_flow.png",
            "usage": "用于讲‘原版输出如何变成几何监督，再变成可复用缓存’这条链，图上只保留概念模块和公式。",
            "work_items": [
                "明确画出：原版输出并不是直接进入训练，而是先被整理成几何教师。",
                "明确画出：几何教师之后有一层几何反投影与坐标整理，用来生成点图、相机元信息和世界坐标监督。",
                "明确画出：这些几何产物会被固化为可复用几何缓存，后续训练直接读取这一层。",
                "这张图的重点是‘监督底座被固定下来’，不是讲实验现象。",
            ],
            "labels": [
                "Original Outputs",
                "Depth / Confidence / Camera",
                "Geometry Teacher",
                "Unprojection",
                "Point Map",
                "Frame Meta",
                "Camera Meta",
                "Geometry Cache",
                "Dataset Read Path",
                "P_world = R^T * ([x_cam, y_cam, z]^T - t)",
            ],
            "prompt": """Create a clean 16:9 research-style flow diagram for a geometry-teacher pipeline.

Left to right blocks:
1. Original Outputs
   small labels: Depth / Confidence / Camera
2. Geometry Teacher
3. Unprojection
4. Geometry Products
   small labels: Point Map / Frame Meta / Camera Meta
5. Geometry Cache
6. Dataset Read Path

Arrow logic:
- original outputs feed the Geometry Teacher
- the Geometry Teacher feeds Unprojection
- Unprojection feeds Geometry Products
- Geometry Products feed Geometry Cache
- Geometry Cache feeds Dataset Read Path

Include a compact formula area:
- x_cam = (u + o - cx) * z / fx
- y_cam = (v + o - cy) * z / fy
- P_world = R^T * ([x_cam, y_cam, z]^T - t)

Style:
- original output block soft orange
- teacher block beige
- geometry blocks soft green
- cache block soft blue
- research-paper aesthetics, thin arrows, no screenshots, lots of whitespace

Main message:
native outputs are reorganized into stable geometry supervision and then frozen into a reusable geometry cache.""",
            "negative": """Do not use filenames, code paths, or command-line text. No charts, no rendered scenes, no benchmark tables, no dark mode, no decorative clutter.""",
        },
        {
            "title": "图 4. 监督路由图",
            "filename": "loss_routing.png",
            "usage": "用于把训练入口和监督路由关系讲清楚，图上只出现训练概念，不出现源码文件名。",
            "work_items": [
                "左侧明确是伪几何数据集，表示训练读取的是图像加几何缓存，而不是原版即时输出。",
                "中间必须是一个占主画面的监督路由器，表示训练时所有监督项都在这里统一组织。",
                "路由器内部要体现：置信度加权、前景权重、跨视图深度一致性。",
                "右侧必须扇出具体监督项：深度、点图、重投影、法向、相机、多视图深度。",
                "底部要有总损失和置信度权重公式，让这张图能直接对上训练逻辑。",
            ],
            "labels": [
                "Pseudo-geometry Dataset",
                "Image + Geometry Cache",
                "Supervision Router",
                "Confidence Weighting",
                "Foreground Weighting",
                "Multiview Depth Consistency",
                "Depth",
                "Point Map",
                "Reprojection",
                "Normal",
                "Camera",
                "Multiview Depth",
                "L_total",
            ],
            "prompt": """Create a clean 16:9 supervision-routing diagram in research-paper style.

Layout:
- left block: Pseudo-geometry Dataset
  small label: Image + Geometry Cache
- center large dominant block: Supervision Router
  inner labels:
  - Confidence Weighting
  - Foreground Weighting
  - Multiview Depth Consistency
- right side fan-out branches:
  - Depth
  - Point Map
  - Reprojection
  - Normal
  - Camera
  - Multiview Depth
- far right final merge block: L_total

Arrow logic:
- dataset feeds the Supervision Router
- the router fans out to all supervision branches
- all branches merge conceptually into L_total

Show these formulas in a small clean formula area:
- L_total = sum_k lambda_k L_k
- w_conf(x) = valid(x) * clip((c(x)-t)/(1-t), 0, 1)^gamma

Use clean rounded rectangles, balanced spacing, light academic colors, no screenshots, no icons.

The main idea is:
all local supervision terms are organized by one routing layer around the original backbone.""",
            "negative": """Do not use filenames, source-code text, result charts, dark backgrounds, neon colors, or generic corporate infographic icons.""",
        },
        {
            "title": "图 5. 前景权重语义图",
            "filename": "fg_boost_semantics.png",
            "usage": "用于把当前前景权重逻辑讲清楚，只说语义流程，不说源码文件名。",
            "work_items": [
                "必须表现这是一条权重构造路径，不是网络结构路径。",
                "必须表现顺序：有效区域、前景区域、内部区域、最终权重。",
                "必须把前景 gating 和内部增权分成两层意思，不要画成一个混在一起的热力图。",
                "必须保留三条核心公式，让图像本身就能说明当前语义。",
            ],
            "labels": [
                "Valid Geometry Region",
                "Foreground Region",
                "Interior Region",
                "Foreground Gating",
                "Interior Boost",
                "Final Supervision Weight",
                "M_valid = M_valid_all * M_fg",
                "M_interior = erode(M_fg, r)",
                "w_final = w_base * [1 + (beta - 1) * M_interior]",
            ],
            "prompt": """Create a minimal but explicit 16:9 semantic diagram for foreground-based supervision weighting.

The figure should be a left-to-right semantic pipeline:
1. Valid Geometry Region
2. Foreground Region
3. Foreground Gating
4. Interior Region
5. Interior Boost
6. Final Supervision Weight

The figure must clearly separate two meanings:
- foreground gating
- interior boost

Include these formulas directly in the figure:
- M_valid = M_valid_all * M_fg
- M_interior = erode(M_fg, r)
- w_final = w_base * [1 + (beta - 1) * M_interior]

Style:
- abstract mask / weight-map blocks
- soft orange for foreground-related blocks
- beige for morphology step
- soft green for gating step
- soft blue for final weight
- thin arrows, large whitespace, clean paper figure feel

Main message:
this branch changes supervision weights, not the original network structure.""",
            "negative": """Do not use filenames, attention heatmaps, real segmentation screenshots, experiment numbers, result plots, dark backgrounds, or photorealistic images.""",
        },
        {
            "title": "图 6. 对照-测量-探针脚手架图",
            "filename": "compare_probe_framework.png",
            "usage": "用于说明统一对照、结构测量和自动化探针是如何串起来的，图上只出现流程模块名。",
            "work_items": [
                "先有训练检查点或中间模型，再经过一次统一可视化对照。",
                "统一对照之后分出两条测量支路：重影测量和点支撑测量。",
                "测量结果再进入自动化探针层，至少体现一步探针、步长曲线探针、前景权重探针。",
                "最后输出是结构化汇总材料，而不是结果优劣判断。",
            ],
            "labels": [
                "Checkpoint",
                "Finetune Step",
                "Unified Visual Compare",
                "Ghost Measure",
                "Point-support Measure",
                "One-step Probe",
                "Stepcurve Probe",
                "Foreground-weight Probe",
                "Compare Views",
                "Summary Tables",
                "Report Pages",
            ],
            "prompt": """Create a clean 16:9 workflow diagram for the local compare-measure-probe shell.

Left to right layout:
- Checkpoint
- Finetune Step
- Unified Visual Compare
- two analysis branches:
  Ghost Measure
  Point-support Measure
- a probe layer with three modules:
  One-step Probe
  Stepcurve Probe
  Foreground-weight Probe
- final outputs:
  Compare Views
  Summary Tables
  Report Pages

Arrow logic:
- checkpoint feeds finetune step
- finetune step feeds unified visual compare
- visual compare feeds both analysis branches
- both analysis branches feed the probe layer
- the probe layer feeds final structured outputs

Style:
- checkpoint beige
- finetune soft green
- visual compare soft blue
- ghost branch soft orange
- point-support branch soft green
- probe layer soft purple
- report outputs neutral blue-gray
- thin arrows, rounded boxes, clean academic layout

Main message:
the local line adds a full compare, measurement, and probe shell around the same original backbone.""",
            "negative": """Do not use filenames, screenshots, benchmark plots, dashboard collage, dark theme, heavy gradients, or result-oriented wording.""",
        },
    ]

    lines = [
        "# 原版 VGGT PPT 插图生成 Prompt",
        "",
        "## 总原则",
        "",
        "- 这一版只用“模块职责”和“信息流关系”写图，不在图片提示词里写源码文件名。",
        "- 写法参考架构图语言：输入是什么，中间经过哪些块，右侧输出什么，新增壳层放在哪一层。",
        "- 原版 VGGT 只作为上层基础通路或底座，主画面必须是你后来补上的几何、训练、分析和自动化模块。",
        "- 所有图都只讲结构、公式、职责和关系，不讲效果，不讲结果，不讲优劣。",
        "",
        "## 使用方式",
        "",
        "- 先看“这张图必须体现的具体工作”，确认图的主题是不是落在你的改动上。",
        "- 再看“图中文字建议”，这些标签都是给图片里直接写的模块名，已经去掉源码文件名。",
        "- 最后直接复制“可直接用于生图的详细 Prompt”去生成图片；它们现在是架构图思路，不是代码文件清单。",
        "",
    ]
    for idx, section in enumerate(sections, start=1):
        lines.extend(
            [
                f"## {idx}. {section['title']}",
                "",
                f"- 建议文件名: `{section['filename']}`",
                f"- 用途: {section['usage']}",
                "",
                "### 这张图必须体现的具体工作",
                "",
            ]
        )
        for item in section["work_items"]:
            lines.append(f"- {item}")
        lines.extend(
            [
                "",
                "### 图中文字建议",
                "",
            ]
        )
        for item in section["labels"]:
            lines.append(f"- {item}")
        lines.extend(
            [
                "",
                "### 可直接用于生图的详细 Prompt",
                "",
                "```text",
                section["prompt"].strip(),
                "```",
                "",
                "### Negative Prompt",
                "",
                "```text",
                section["negative"].strip(),
                "```",
                "",
            ]
        )
    return "\n".join(lines).strip() + "\n"


def build_outline_md(slides: list[SlideSpec], deck_title: str, ppt_name: str) -> str:
    lines: list[str] = []
    lines.append(f"# {deck_title} 大纲")
    lines.append("")
    lines.append("## 总原则")
    lines.append("")
    lines.append("- 所有页面都必须围绕“我们相对原版 VGGT 改了什么”来讲。")
    lines.append("- 原版内容只能作为对照基线出现，不能成为画面主体。")
    lines.append("- 所有插图都必须服务于“新增了哪些代码层、训练层、分析层、自动化层”这个主线。")
    lines.append("- 如果某一页出现原版 VGGT，只能用于说明哪些部分保留不动。")
    lines.append("")
    lines.append("## 插图统一要求")
    lines.append("")
    lines.append("- 架构图主体必须是“相对原版新增的工作”，不是官方 VGGT 原图换字。")
    lines.append("- 代码关系图必须明确区分“直接改原版文件”和“围绕原版新增外层代码壳”。")
    lines.append("- teacher 图必须强调“把原版输出整理成几何 teacher 和 NPZ 契约”。")
    lines.append("- wrapper 图必须强调“把监督路由起来”，不能强调效果。")
    lines.append("- compare / probe 图必须强调“新增自动化脚手架”，不能强调结果优劣。")
    lines.append("")
    lines.append(f"对应文件: `{ppt_name}`")
    lines.append("用途: 便于直接给 AI 生成 PPT，也便于口头汇报时按页讲解。")
    lines.append("")
    for idx, spec in enumerate(slides, start=1):
        lines.append(f"## Slide {idx:02d}. {spec.title}")
        lines.append("")
        lines.append("### 屏幕主内容")
        lines.append("")
        for bullet in spec.bullets:
            lines.append(f"- {bullet}")
        if spec.formulas:
            lines.append("")
            lines.append("### 公式 / 规则")
            lines.append("")
            lines.append("```text")
            for formula in spec.formulas:
                lines.append(formula)
            lines.append("```")
        if spec.code_refs:
            lines.append("")
            lines.append("### 代码位置")
            lines.append("")
            for ref in spec.code_refs:
                lines.append(f"- {ref}")
        if spec.image_paths:
            lines.append("")
            lines.append("### 建议插图")
            lines.append("")
            for path in spec.image_paths:
                lines.append(f"- `{path.as_posix()}`")
            if spec.image_caption:
                lines.append(f"- 图注: {spec.image_caption}")
        if spec.speaker_notes:
            lines.append("")
            lines.append("### 讲解备注")
            lines.append("")
            lines.append(spec.speaker_notes)
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def write_deck_bundle(slides: list[SlideSpec], pptx_path: Path, md_path: Path, deck_title: str) -> tuple[Path, Path]:
    prs = build_presentation(slides)
    actual_pptx = save_with_fallback_bytes(pptx_path, lambda p: prs.save(str(p)))
    actual_md = save_with_fallback_bytes(
        md_path,
        lambda p: p.write_text(build_outline_md(slides, deck_title, actual_pptx.name), encoding="utf-8"),
    )
    return actual_pptx, actual_md


def main() -> None:
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    download_official_architecture()
    data = collect_data()
    data["assets"] = generate_visual_assets(data)
    diff_summary = collect_code_diff_summary()

    deck_jobs = [
        ("原版 VGGT 本地代码改动总览", OUT_PPTX, OUT_MD, build_combined_deck(data, diff_summary)),
        ("原版 VGGT 改动与微调代码主线", OUT_CODE_PPTX, OUT_CODE_MD, build_code_and_finetune_deck(data, diff_summary)),
        ("原版 VGGT Probe 自动化脚手架", OUT_PROBE_PPTX, OUT_PROBE_MD, build_probe_deck(data)),
        ("原版 VGGT Mask / Ghost / Point-Support 代码分支", OUT_MASK_PPTX, OUT_MASK_MD, build_mask_deck(data)),
    ]

    for deck_title, pptx_path, md_path, slides in deck_jobs:
        actual_pptx, actual_md = write_deck_bundle(slides, pptx_path, md_path, deck_title)
        print(f"Wrote PPTX: {actual_pptx}")
        print(f"Wrote outline: {actual_md}")

    prompt_path = save_with_fallback_bytes(
        OUT_PROMPTS_MD,
        lambda p: p.write_text(build_visual_prompts_md(), encoding="utf-8"),
    )
    print(f"Wrote prompts: {prompt_path}")


def geometry_formula_lines() -> list[str]:
    return [
        "x_cam = ((u + o - c_x) / f_x) * z",
        "y_cam = ((v + o - c_y) / f_y) * z",
        "P_cam = [x_cam, y_cam, z]^T",
        "P_world = R^T * (P_cam - t)",
    ]


def routing_formula_lines() -> list[str]:
    return [
        "M_valid_all(i) = 1[z_tgt(i) > 1e-6]",
        "M_sup(i) = M_valid_all(i) * [M_fg(i) + eta * (1 - M_fg(i))]",
        "c_hat(i) = clip((c(i) - t) / (1 - t), 0, 1)",
        "w_conf(i) = M_sup(i) * c_hat(i)^gamma * Q_pvq(i) * G_dyn(i)",
        "w_base(i) = M_sup(i) (uniform) | M_sup(i) * [alpha * w_conf(i) + (1 - alpha)] (mix) | w_conf(i) (conf)",
        "L_total = lambda_depth L_depth + lambda_point L_point + lambda_point_reproj L_point_reproj + lambda_point_normal_consis L_point_normal_consis + lambda_point_mv_depth L_point_mv_depth + lambda_point_mv_mask L_point_mv_mask + lambda_fg_structure_edge L_fg_structure_edge + lambda_point_mv_outside_ring L_point_mv_outside_ring + lambda_conf_eff L_conf + lambda_fg_presence L_fg_presence + lambda_geom_cons L_geom_cons + lambda_cam_eff L_cam",
    ]


def fg_formula_lines() -> list[str]:
    return [
        "M_fg_raw(i) = 1[x(i) >= 0.5] if max(x) <= 1.5 else 1[x(i) > 0]",
        "M_fg(i) = erode(M_fg_raw, r0) if fg_mask_erode_px = r0 > 0 else M_fg_raw(i)",
        "M_sup(i) = M_valid_all(i) * [M_fg(i) + eta * (1 - M_fg(i))]",
        "M_boost(i) = erode(M_fg(i), r) if region_mode = interior_only else M_fg(i)",
        "erode(x) = 1 - dilate(1 - x)",
        "w_final(i) = w_base(i) * [1 + (beta - 1) * M_boost(i)]",
        "outside FG -> 0, boundary ring -> w_base, eroded interior -> beta * w_base",
    ]


def ghost_formula_lines() -> list[str]:
    return [
        "width_ratio = W_pred / W_mask,  area_ratio = A_pred / A_mask,  center_offset_ratio = |c_pred - c_mask| / W_mask",
        "peak_count = Peaks(smooth(sum_y MidBin(y, x)))",
        "ghost_score = 1.00 * max(0, width_ratio - 1.10) + 0.40 * max(0, area_ratio - 1.30) + 0.60 * max(0, peak_count - 1) + 0.50 * max(0, center_offset_ratio - 0.22)",
        "dark_penalty = 40 * max(0, 0.045 - pred_luma_mean) + 12 * max(0, 0.10 - pred_nonblack_ratio_thr008)",
        "collapse_penalty = 6 * max(0, 0.55 - area_ratio) + 6 * max(0, 0.65 - width_ratio)",
        "ghost_visual_score = ghost_score + dark_penalty + collapse_penalty",
    ]


def support_formula_lines() -> list[str]:
    return [
        "support_map = green_channel(weight_native.png)",
        "T_support = max(threshold_floor, quantile_0.75({S(x) | S(x) > 0})) in auto mode",
        "subject_support_share = sum_x S(x) * M(x) / sum_x S(x)",
        "largest_component_share = max_c sum_{x in c} S(x) / sum_{x in active} S(x)",
        "secondary_component_mass = second_max_c sum_{x in c} S(x) / sum_{x in active} S(x)",
        "support_peak_count = Peaks(smooth(sum_y ActiveROI(y, x)))",
        "subject_psnr = -10 log10(sum_x ||pred(x) - tgt(x)||_2^2 * M(x) / (3 * sum_x M(x)))",
        "subject_l1 = sum_x ||pred(x) - tgt(x)||_1 * M(x) / (3 * sum_x M(x))",
    ]


def mask_probe_formula_lines() -> list[str]:
    return [
        "Current probe switches: use_fg_mask = on, fg_supervision_bg_floor = 0, region_mode = interior_only, fg_supervision_region_erode_px = 2",
        "M_sup(i) = M_valid_all(i) * M_fg(i)",
        "M_interior(i) = erode(M_fg(i), 2)",
        "w_final(i) = w_base(i) * [1 + (beta - 1) * M_interior(i)]",
        "alpha = 1, 2, 4  ->  beta = 2, 3, 5",
    ]


def generate_loss_routing_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.8))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.05, 0.55, 0.20, 0.16, "Pseudo-geometry\nDataset", fc="#f4efe4", ec="#9a8d74", lw=2.4, fontsize=17)
    ax.text(0.15, 0.49, "image + geometry cache", ha="center", fontsize=12, color="#6b6251")

    draw_arch_box(ax, 0.31, 0.46, 0.27, 0.26, "Supervision Router", fc="#eef5ea", ec="#5b8742", lw=2.8, fontsize=20)
    ax.text(0.445, 0.54, "confidence remap\nforeground gating\nmultiview consistency", ha="center", va="center", fontsize=13, color="#35522a")
    draw_arch_arrow(ax, 0.25, 0.63, 0.31, 0.59)

    targets = [
        ("Depth", 0.69, 0.72, "#fff1e8", "#c96310"),
        ("Point Map", 0.83, 0.72, "#eef0fb", "#6882b4"),
        ("Reprojection", 0.69, 0.57, "#eef0fb", "#6882b4"),
        ("Normal", 0.83, 0.57, "#eef5ea", "#5b8742"),
        ("MV-depth", 0.69, 0.42, "#eef5ea", "#5b8742"),
        ("Camera", 0.83, 0.42, "#f7eefb", "#7b61c8"),
    ]
    for text, x, y, fc, ec in targets:
        draw_arch_box(ax, x, y, 0.12, 0.09, text, fc=fc, ec=ec, lw=2.2, fontsize=14)
        draw_arch_arrow(ax, 0.58, 0.59, x, y + 0.045)

    draw_arch_box(ax, 0.77, 0.22, 0.18, 0.10, "L_total", fc="#eef0fb", ec="#6882b4", lw=2.4, fontsize=18)
    draw_arch_arrow(ax, 0.58, 0.59, 0.77, 0.27)

    formula = (
        "M_sup(i) = M_valid_all(i) * [M_fg(i) + eta * (1 - M_fg(i))]\n"
        "c_hat(i) = clip((c(i) - t) / (1 - t), 0, 1)\n"
        "w_conf(i) = M_sup(i) * c_hat(i)^gamma * Q_pvq(i) * G_dyn(i)\n"
        "L_total = sum_k lambda_k L_k"
    )
    ax.text(0.36, 0.16, formula, ha="left", va="center", fontsize=13, color="#222222")
    ax.text(0.50, 0.06, "原理: 先定义监督有效域，再构造基础权重，最后把多种监督项收束到统一训练入口。", ha="center", fontsize=13, color="#555555")
    return save_fig(fig, cache_png("loss_routing"))


def generate_fg_boost_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.0))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.05, 0.50, 0.15, 0.20, "Valid Region", fc="#eef5ea", ec="#5b8742", lw=2.5, fontsize=18)
    ax.text(0.125, 0.43, "M_valid_all", ha="center", fontsize=13, color="#35522a")

    draw_arch_box(ax, 0.25, 0.50, 0.15, 0.20, "Foreground\nMask", fc="#fff1e8", ec="#c96310", lw=2.5, fontsize=18)
    ax.text(0.325, 0.43, "M_fg_raw -> M_fg", ha="center", fontsize=13, color="#c96310")

    draw_arch_box(ax, 0.45, 0.50, 0.16, 0.20, "Foreground\nGating", fc="#eef5ea", ec="#5b8742", lw=2.5, fontsize=18)
    ax.text(0.53, 0.43, "M_sup = M_valid_all * [M_fg + eta * (1 - M_fg)]", ha="center", fontsize=11.5, color="#35522a")

    draw_arch_box(ax, 0.67, 0.50, 0.14, 0.20, "Interior\nRegion", fc="#f4efe4", ec="#9a8d74", lw=2.5, fontsize=18)
    ax.text(0.74, 0.43, "M_boost = erode(M_fg, r)", ha="center", fontsize=12, color="#6b6251")

    draw_arch_box(ax, 0.84, 0.50, 0.12, 0.20, "Final\nWeight", fc="#eef0fb", ec="#6882b4", lw=2.5, fontsize=19)
    ax.text(0.90, 0.43, "w_final", ha="center", fontsize=13, color="#4c5f8e")

    draw_arch_arrow(ax, 0.20, 0.60, 0.25, 0.60)
    draw_arch_arrow(ax, 0.40, 0.60, 0.45, 0.60)
    draw_arch_arrow(ax, 0.61, 0.60, 0.67, 0.60)
    draw_arch_arrow(ax, 0.81, 0.60, 0.84, 0.60)

    formula = (
        "M_fg_raw(i) = 1[x(i) >= 0.5] if max(x) <= 1.5 else 1[x(i) > 0]\n"
        "M_fg(i) = erode(M_fg_raw, r0) if fg_mask_erode_px = r0 > 0 else M_fg_raw(i)\n"
        "erode(x) = 1 - dilate(1 - x)\n"
        "w_final(i) = w_base(i) * [1 + (beta - 1) * M_boost(i)]"
    )
    ax.text(0.08, 0.22, formula, ha="left", va="center", fontsize=12.5, color="#222222")
    ax.text(0.50, 0.06, "流程: 先把监督收缩到前景，再选出内部区域，最后只对内部区域做额外增权。", ha="center", fontsize=13, color="#555555")
    return save_fig(fig, cache_png("fg_boost_semantics"))


def generate_ghost_formula_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.0))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.05, 0.58, 0.20, 0.16, "Ghost Triplet", fc="#eef0fb", ec="#6882b4", lw=2.5, fontsize=18)
    ax.text(0.15, 0.50, "mask | pred | target", ha="center", fontsize=12, color="#4c5f8e")
    draw_arch_box(ax, 0.31, 0.58, 0.16, 0.16, "Mask BBox", fc="#fff1e8", ec="#c96310", lw=2.5, fontsize=18)
    ax.text(0.39, 0.50, "W_mask, A_mask, c_mask", ha="center", fontsize=12, color="#c96310")
    draw_arch_box(ax, 0.53, 0.58, 0.16, 0.16, "Pred BBox", fc="#eef5ea", ec="#5b8742", lw=2.5, fontsize=18)
    ax.text(0.61, 0.50, "W_pred, A_pred, c_pred", ha="center", fontsize=12, color="#35522a")
    draw_arch_box(ax, 0.75, 0.58, 0.20, 0.16, "Peak Profile", fc="#f7eefb", ec="#7b61c8", lw=2.5, fontsize=18)
    ax.text(0.85, 0.50, "smooth(sum_y MidBin(y, x))", ha="center", fontsize=11.5, color="#5c4a87")

    draw_arch_arrow(ax, 0.25, 0.66, 0.31, 0.66)
    draw_arch_arrow(ax, 0.47, 0.66, 0.53, 0.66)
    draw_arch_arrow(ax, 0.69, 0.66, 0.75, 0.66)

    formula = (
        "ghost_score = 1.00 * max(0, width_ratio - 1.10)\n"
        "            + 0.40 * max(0, area_ratio - 1.30)\n"
        "            + 0.60 * max(0, peak_count - 1)\n"
        "            + 0.50 * max(0, center_offset_ratio - 0.22)\n"
        "ghost_visual_score = ghost_score + dark_penalty + collapse_penalty"
    )
    ax.text(0.08, 0.22, formula, ha="left", va="center", fontsize=12.5, color="#222222")
    ax.text(0.50, 0.06, "原理: 先量化宽度、面积、横向多峰和中心漂移，再额外惩罚过黑与塌缩。", ha="center", fontsize=13, color="#555555")
    return save_fig(fig, cache_png("ghost_scoring_formula"))


def generate_support_metrics_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.0))
    fig.patch.set_facecolor("#fcfbf7")
    ax.set_facecolor("#fcfbf7")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_axis_off()

    draw_arch_box(ax, 0.05, 0.58, 0.18, 0.16, "Support Map", fc="#eef0fb", ec="#6882b4", lw=2.5, fontsize=18)
    ax.text(0.14, 0.50, "green(weight_native)", ha="center", fontsize=12, color="#4c5f8e")
    draw_arch_box(ax, 0.29, 0.58, 0.18, 0.16, "Subject Mask", fc="#fff1e8", ec="#c96310", lw=2.5, fontsize=18)
    ax.text(0.38, 0.50, "GT mask preferred", ha="center", fontsize=12, color="#c96310")
    draw_arch_box(ax, 0.53, 0.58, 0.18, 0.16, "Active Components", fc="#eef5ea", ec="#5b8742", lw=2.5, fontsize=18)
    ax.text(0.62, 0.50, "threshold + connected components", ha="center", fontsize=11.5, color="#35522a")
    draw_arch_box(ax, 0.77, 0.58, 0.18, 0.16, "Masked Error", fc="#f7eefb", ec="#7b61c8", lw=2.5, fontsize=18)
    ax.text(0.86, 0.50, "pred vs target inside mask", ha="center", fontsize=11.5, color="#5c4a87")

    draw_arch_arrow(ax, 0.23, 0.66, 0.29, 0.66)
    draw_arch_arrow(ax, 0.47, 0.66, 0.53, 0.66)
    draw_arch_arrow(ax, 0.71, 0.66, 0.77, 0.66)

    formula = (
        "subject_support_share = sum_x S(x) * M(x) / sum_x S(x)\n"
        "largest_component_share = max_c sum_{x in c} S(x) / sum_{x in active} S(x)\n"
        "subject_psnr = -10 log10(sum_x ||pred(x) - tgt(x)||_2^2 * M(x) / (3 * sum_x M(x)))\n"
        "subject_l1 = sum_x ||pred(x) - tgt(x)||_1 * M(x) / (3 * sum_x M(x))"
    )
    ax.text(0.08, 0.22, formula, ha="left", va="center", fontsize=12.5, color="#222222")
    ax.text(0.50, 0.06, "流程: 先从 support 图取结构量，再和主体 mask 结合，最后补上 mask 内重建误差。", ha="center", fontsize=13, color="#555555")
    return save_fig(fig, cache_png("point_support_formula"))


def generate_visual_assets(data: dict) -> dict[str, Path]:
    return {
        "current_architecture": generate_current_architecture_diagram(),
        "code_touch_map": generate_code_touch_diagram(),
        "teacher_npz_flow": generate_teacher_npz_diagram(),
        "loss_routing": generate_loss_routing_diagram(),
        "fg_boost_semantics": generate_fg_boost_diagram(),
        "compare_probe_framework": generate_compare_probe_diagram(),
        "ghost_scoring_formula": generate_ghost_formula_diagram(),
        "point_support_formula": generate_support_metrics_diagram(),
    }


def render_full_image_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, spec.title)
    info_box = add_panel(slide, Inches(0.55), Inches(1.05), Inches(12.2), Inches(1.0))
    add_bullets(info_box.text_frame, wrap_lines(spec.bullets, 115), font_size=15)
    if spec.image_paths:
        img = ensure_cache_image(spec.image_paths[0], safe_stem(spec.title))
        fit_image(slide, img, Inches(0.55), Inches(2.10), Inches(12.2), Inches(3.45))

    meta_panel = add_panel(slide, Inches(0.55), Inches(5.75), Inches(12.2), Inches(1.35), fill=RGBColor(244, 239, 231))
    meta_frame = meta_panel.text_frame
    meta_frame.clear()
    meta_frame.word_wrap = True
    meta_frame.vertical_anchor = MSO_ANCHOR.TOP

    sections: list[tuple[str, list[str], bool]] = []
    if spec.formulas:
        sections.append(("公式 / 规则", spec.formulas[:8], True))
    if spec.code_refs:
        sections.append(("代码位置", spec.code_refs[:4], False))
    if spec.image_caption:
        sections.append(("图像说明", [spec.image_caption], False))

    for s_idx, (header, lines, monospace) in enumerate(sections):
        p = meta_frame.paragraphs[0] if s_idx == 0 else meta_frame.add_paragraph()
        p.text = header
        run = p.runs[0] if p.runs else p.add_run()
        run.text = header
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        p.space_after = Pt(2)
        for line in lines:
            pp = meta_frame.add_paragraph()
            pp.text = line
            rr = pp.runs[0] if pp.runs else pp.add_run()
            rr.text = line
            rr.font.name = "Consolas" if monospace or any(ch in line for ch in "=<>[]{}_") else "Microsoft YaHei"
            rr.font.size = Pt(9.5)
            rr.font.color.rgb = TEXT
            pp.space_after = Pt(1)


def build_combined_deck(data: dict, diff_summary: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})
    return [
        SlideSpec(
            title="基于原版 VGGT 的本地代码改动总览",
            bullets=[
                "这套 PPT 只讲相对原版 VGGT 改了什么，重点放在结构位置、原理、公式和流程，不讨论跑出来的效果。",
                "主线只有一句话: 原版主干基本保留，本地工作主要补在 geometry teacher、geometry cache、supervision router、visual compare 和 probe shell 这几层。",
                "讲解顺序按“整体架构 -> 直接改原版哪里 -> 新增训练语义 -> 新增分析与自动化”展开。",
            ],
            speaker_notes="先把边界压住: 这里只讲代码结构与公式，不讲结果优劣。",
        ),
        build_official_compare_slide(data),
        SlideSpec(
            title="代码改动分层: 哪些是原版内改动, 哪些是外层新增",
            bullets=[
                "直接改到原版仓库内部的入口并不多，主要是几何约定、图像读取加速、COLMAP 兼容和少量结构说明补充。",
                "主体新增工作在外层壳层: 先把原版输出整理成 geometry teacher，再固化成 cache，再由监督路由器接成统一训练入口。",
                "后面 compare、ghost、point-support 和 probe automation 也是新增壳层，而不是重写原版主干。",
            ],
            formulas=[
                "原版主干: Aggregator + Camera Head + Dense Geometry Head",
                "本地新增主链: Geometry Teacher -> Geometry Cache -> Supervision Router -> Visual Compare -> Probe Shell",
            ],
            code_refs=diff_summary["direct_original_touches"] + diff_summary["new_local_modules"][:4],
            image_paths=existing_images(assets.get("code_touch_map")),
            image_caption="左侧是少量原版内改动，右侧是围绕原版新增的训练与分析壳层。",
            layout="full_image",
        ),
        SlideSpec(
            title="Geometry Teacher 与 Geometry Cache: 先把原版输出整理成稳定监督",
            bullets=[
                "原版输出的 depth / confidence / camera 不直接喂进后续训练，而是先整理成可重复读取的几何监督底座。",
                "反投影这一步把像素、深度和相机参数变成世界坐标点图，因此 point map、camera meta 和 frame meta 都在这一层被固定下来。",
                "cache 的价值是把监督契约固定下来，后面的 finetune 和 probe 都复用同一份几何底座。",
            ],
            formulas=geometry_formula_lines(),
            code_refs=[
                "vggt/utils/geometry.py",
                "vggt_geom.py",
                "precompute_zju_vggt_geom.py",
                "modal_run_train.py",
            ],
            image_paths=existing_images(assets.get("teacher_npz_flow")),
            image_caption="流程: 原版输出先变成 geometry teacher，再固化成可重复读取的 geometry cache。",
            layout="full_image",
        ),
        SlideSpec(
            title="Supervision Router: 把多种监督项统一接入训练入口",
            bullets=[
                "训练入口的关键不是改 backbone，而是定义监督有效域、构造置信度权重，并把多种损失组织成同一个 router。",
                "路由过程先决定哪些像素可监督，再决定这些像素权重多大，最后再把 depth、point、reprojection、multiview depth、camera 等项一起求和。",
                "这样做的结果是，同一套 teacher 和数据入口可以支撑不同训练分支，而不需要反复改主干。",
            ],
            formulas=routing_formula_lines(),
            code_refs=[
                "finetune_vggt_pseudo.py: _build_conf_weight",
                "finetune_vggt_pseudo.py: _build_per_view_conf_quantile_mask",
                "finetune_vggt_pseudo.py: train-loop loss routing",
            ],
            image_paths=existing_images(assets.get("loss_routing")),
            image_caption="原理: 先构造监督权重，再把多种监督项统一接到同一个训练入口。",
            layout="full_image",
        ),
        SlideSpec(
            title="Foreground Weighting: 当前实现是前景门控加内部增权",
            bullets=[
                "这条路径改的是训练权重语义，不是网络结构。先用 foreground mask 决定监督有效域，再对前景内部区域额外增权。",
                "代码里 foreground mask 先做二值化，必要时先做预侵蚀，然后再进入监督域构造和 boost 区域构造。",
                "因此向导师汇报时，应该把它讲成“supervision weighting semantics”，而不是讲成新结构模块。",
            ],
            formulas=fg_formula_lines(),
            code_refs=[
                "finetune_vggt_pseudo.py: _erode_mask_tensor",
                "finetune_vggt_pseudo.py: _build_fg_supervision_boost_mask",
                "finetune_vggt_pseudo.py: _apply_fg_supervision_boost",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="流程: 先前景门控，再内部增权；背景、边界环带和内部区域的权重语义不同。",
            layout="full_image",
        ),
        SlideSpec(
            title="Ghost Scorer: 把重影现象翻译成结构化分数",
            bullets=[
                "ghost 不是训练损失，而是 compare 之后的后处理评分。输入固定是一张 triplet 图: mask、prediction、target。",
                "计算流程是先从 mask 和 prediction 提取包围框与面积，再统计 prediction 横向 profile 的峰数，最后叠加过黑和塌缩惩罚。",
                "所以 ghost_visual_score 本质上是在量化“过宽、过厚、多峰、偏移、太暗、塌缩”这些视觉描述。",
            ],
            formulas=ghost_formula_lines(),
            code_refs=[
                "tools/score_ghosting_from_cat_pred.py: _score_one",
                "tools/score_ghosting_from_cat_pred.py: _count_peaks",
            ],
            image_paths=existing_images(assets.get("ghost_scoring_formula")),
            image_caption="流程: triplet 拆分 -> 宽度/面积/中心偏移 -> 横向多峰 -> 惩罚项 -> ghost_visual_score。",
            layout="full_image",
        ),
        SlideSpec(
            title="Point-support 与 Masked Error: 结构量和重建量分开定义",
            bullets=[
                "point-support 图读的是 support 形状，不直接等于重建误差；masked PSNR 和 masked L1 则是主体区域内的重建误差。",
                "脚本先从 support 图的绿色通道读出 support map，再结合主体 mask 计算主体份额、连通块集中度和峰数。",
                "随后再单独在主体 mask 内比较 prediction 与 target，得到 subject_psnr 和 subject_l1。",
            ],
            formulas=support_formula_lines(),
            code_refs=[
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
            ],
            image_paths=existing_images(assets.get("point_support_formula")),
            image_caption="流程: support 结构量和 mask 内误差量分别定义，不能混成同一种指标。",
            layout="full_image",
        ),
        SlideSpec(
            title="Compare / Metrics / Probe: 围绕原版主线新增的工程壳层",
            bullets=[
                "统一 compare 负责把预测、目标和辅助图拼成稳定输出格式，后面的 ghost scorer 和 point-support 脚本都从这一层接入。",
                "one-step、stepcurve 和 mask-boost 只是三种任务编排方式，它们复用同一套 compare 和 measurement 壳层。",
                "所以这一层的贡献是把训练分支组织成可重复执行、可统一汇总的工程流程。",
            ],
            formulas=[
                "prepare task -> run task -> render compare -> score metrics -> summarize reports",
                "同一 backbone + 同一 compare / measurement shell + 不同 probe orchestration",
            ],
            code_refs=[
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
                "scripts/orig_vggt_one_step_probe/* / scripts/orig_vggt_stepcurve_probe/* / scripts/orig_vggt_mask_boost_probe/*",
            ],
            image_paths=existing_images(assets.get("compare_probe_framework")),
            image_caption="Probe 不是新模型家族，而是围绕同一条原版主线搭起来的自动化壳层。",
            layout="full_image",
        ),
        SlideSpec(
            title="汇报口径: 相对原版, 我们补的是监督与分析壳层",
            bullets=[
                "一句话总结: 原版核心结构保留，新增工作主要在几何监督、监督路由、可视化对照、结构测量和自动化脚手架。",
                "因此 PPT 应该围绕“改了哪里、变量怎么定义、公式怎么展开、流程怎么串起来”来讲，不要把重点放到结果曲线上。",
                "如果导师只记住一句话，那就是: 这套代码把原版 VGGT 拆成了一个可解释、可复用、可继续扩展的本地工程主线。",
            ],
        ),
    ]


def build_code_and_finetune_deck(data: dict, diff_summary: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})
    return [
        SlideSpec(
            title="原版 VGGT 改动与微调: 只讲当前代码主线",
            bullets=[
                "这套单独聚焦最近这条基于原版 VGGT 的代码主线，只讲代码结构、公式、原理和流程。",
                "主干不重写，重点在 geometry teacher、geometry cache 和 supervision router 这几层。",
                "汇报顺序按“原版内改动 -> teacher/cache -> supervision routing -> foreground weighting -> compare shell”展开。",
            ],
        ),
        build_official_compare_slide(data),
        SlideSpec(
            title="原版内部真正改动到哪几处",
            bullets=[
                "原版内部直接改动的入口很少，说明这次工作重点不是重做主干，而是围绕主干加训练和分析壳层。",
                "几何约定相关的修改决定了后面 point map 和世界坐标监督怎么定义；图像读取相关的修改负责把多视图数据入口跑顺。",
                "COLMAP 兼容桥和结构说明补充，解决的是接入与解释问题，而不是结构重写问题。",
            ],
            formulas=[
                "直接改原版的入口少, 但这些入口决定后续监督契约怎么成立",
                "主体新增工作集中在原版输出之后的 teacher / cache / router / compare 壳层",
            ],
            code_refs=diff_summary["direct_original_touches"] + diff_summary["secondary_touches"],
            image_paths=existing_images(assets.get("code_touch_map")),
            image_caption="原版内部改动少，外层新增多，这是当前代码主线最重要的结构特征。",
            layout="full_image",
        ),
        SlideSpec(
            title="Geometry Teacher + Geometry Cache: 监督底座如何被固定下来",
            bullets=[
                "原版 depth、confidence 和 camera 输出先被重组为 geometry teacher，再被冻结成 geometry cache。",
                "teacher 这一层负责定义点图和世界坐标监督怎么来，cache 这一层负责定义后续训练统一读取什么。",
                "所以这条链不是简单的“原版输出直接训”，而是“先定义监督底座，再复用这份底座做训练”。",
            ],
            formulas=geometry_formula_lines(),
            code_refs=[
                "vggt/utils/geometry.py",
                "vggt_geom.py",
                "precompute_zju_vggt_geom.py",
                "modal_run_train.py",
            ],
            image_paths=existing_images(assets.get("teacher_npz_flow")),
            image_caption="几何 teacher 定义监督来源，geometry cache 定义训练读取契约。",
            layout="full_image",
        ),
        SlideSpec(
            title="Supervision Router: 完整损失和基础权重如何组织",
            bullets=[
                "训练入口首先判断哪些像素有效，再把 confidence、前景信息和可选的 quantile / dynamic weight 接到一起。",
                "然后才把 depth、point、reprojection、multiview depth、camera 等损失一起并入总损失。",
                "换句话说，代码里先有监督域和权重，再有损失求和，而不是反过来。",
            ],
            formulas=routing_formula_lines(),
            code_refs=[
                "finetune_vggt_pseudo.py: _build_conf_weight",
                "finetune_vggt_pseudo.py: _build_per_view_conf_quantile_mask",
                "finetune_vggt_pseudo.py: contribution assembly",
            ],
            image_paths=existing_images(assets.get("loss_routing")),
            image_caption="监督路由先定义像素权重，再定义总损失如何汇总。",
            layout="full_image",
        ),
        SlideSpec(
            title="Foreground Weighting: 当前完整公式不是只有三行",
            bullets=[
                "前景权重路径分成四步: mask 二值化与预处理、监督有效域构造、boost 区域构造、最终权重合成。",
                "其中背景是否保留底权重，由 bg floor 控制；当前 probe 因为 bg floor 为 0，所以语义上是 FG-only supervision。",
                "汇报时最好把“门控”和“内部增权”分成两层意思说清楚，避免被误听成普通 attention 图。",
            ],
            formulas=fg_formula_lines(),
            code_refs=[
                "finetune_vggt_pseudo.py: _normalize_mask_binary_np",
                "finetune_vggt_pseudo.py: _erode_mask_tensor",
                "finetune_vggt_pseudo.py: _build_fg_supervision_boost_mask",
                "finetune_vggt_pseudo.py: _apply_fg_supervision_boost",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="当前完整语义: 前景门控定义监督域, 内部增权只提升前景内部像素的监督权重。",
            layout="full_image",
        ),
        SlideSpec(
            title="当前微调主线的另一处关键: 几何一致性项被接进同一个 router",
            bullets=[
                "multiview depth、point reprojection、geometry consistency 和 camera supervision 都不是孤立脚本，而是被接进统一训练入口。",
                "这意味着当前微调主线不是单一 depth loss，而是几何监督、重投影监督和相机监督的组合。",
                "从代码结构上看，这一层的价值是把原版输出变成可控的多项监督训练入口。",
            ],
            formulas=[
                "L_mvdepth ~ Avg rho(|z_ij(x) - D_j(pi_ij(x))|)",
                "L_cam = L_trans + w_rot * L_rot + w_fov * L_fov",
                "L_geom_cons = Avg |point_pred_from_depth - point_tgt| weighted by w_base",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py: multiview depth reprojection path",
                "finetune_vggt_pseudo.py: camera loss path",
                "finetune_vggt_pseudo.py: geometry consistency path",
            ],
        ),
        SlideSpec(
            title="这一套代码的汇报口径",
            bullets=[
                "相对原版，最该强调的不是“数值有没有变好”，而是“哪些监督定义被补齐了，哪些工程壳层被补齐了”。",
                "一句话总结是: 原版保留核心输出，本地代码把监督底座、监督路由和分析壳层全部补齐。",
                "这也是为什么这套 deck 要把公式、流程和职责拆得很开，因为真正新增的价值就在这里。",
            ],
        ),
    ]


def build_probe_deck(data: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})
    return [
        SlideSpec(
            title="One-step / Stepcurve / Mask-Boost: 这是自动化脚手架, 不是新模型",
            bullets=[
                "这套只讲 probe 自动化脚手架在代码里做了什么，不讲任何结果曲线。",
                "one-step、stepcurve 和 mask-boost 的共性是复用同一个训练入口、同一个 compare 壳层和同一组测量脚本。",
                "它们之间的差别主要在任务怎么编排、前缀怎么审计、汇总怎么输出。",
            ],
        ),
        SlideSpec(
            title="Probe 脚手架在整条主线里的位置",
            bullets=[
                "probe 层接在 finetune 之后，统一调用 compare、ghost scorer、point-support measure 和 summary 脚本。",
                "也就是说，probe 层不定义新结构，而是定义训练后如何被统一读取、统一评分和统一汇总。",
                "它的工程价值在于把零散任务变成稳定的任务编排流程。",
            ],
            formulas=[
                "checkpoint -> finetune step -> render compare -> score metrics -> summarize reports",
                "same backbone + same compare shell + different orchestration policy",
            ],
            code_refs=[
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
            ],
            image_paths=existing_images(assets.get("compare_probe_framework")),
            image_caption="probe 脚手架的位置在训练主线之外，但围绕同一条训练主线组织任务。",
            layout="full_image",
        ),
        SlideSpec(
            title="One-step: 最短 horizon 的任务组织方式",
            bullets=[
                "one-step 把“准备任务 -> 跑一次短程训练 -> 生成 compare -> 汇总输出”封装成可重复调用的最小单元。",
                "它的作用不是定义新算法，而是定义最短实验闭环怎么被自动化。",
                "因此对导师汇报时，应该讲成最小任务单元，而不是讲成新分支模型。",
            ],
            formulas=[
                "one-step = prepare task -> run task -> render compare -> summarize",
            ],
            code_refs=[
                "scripts/orig_vggt_one_step_probe/common.py",
                "scripts/orig_vggt_one_step_probe/run_task.py",
                "scripts/orig_vggt_one_step_probe/summarize_runs.py",
            ],
        ),
        SlideSpec(
            title="Stepcurve: 多 horizon 的任务组织方式",
            bullets=[
                "stepcurve 把多个 step horizon 串成一个任务列表，并要求前缀、任务来源和 compare 输出保持一致。",
                "因此 stepcurve 的重点是任务组织和前缀审计，而不是某个单点公式本身。",
                "从工程角度看，它解决的是“如何把同一条主线切成一组可对齐的阶段任务”。",
            ],
            formulas=[
                "stepcurve = horizon list + prefix audit + render compare + summarize",
            ],
            code_refs=[
                "scripts/orig_vggt_stepcurve_probe/audit_prefix.py",
                "scripts/orig_vggt_stepcurve_probe/run_task.py",
                "scripts/orig_vggt_stepcurve_probe/summarize_runs.py",
            ],
        ),
        SlideSpec(
            title="Mask-Boost Probe: 在同一脚手架上切换前景权重语义",
            bullets=[
                "mask-boost probe 不是新训练主线，而是在同一训练入口上切换 foreground weighting 相关开关。",
                "它复用同样的 compare 和 measurement 壳层，只是在训练参数层改 foreground gating 与 interior boost 的语义。",
                "所以这条分支最该汇报的是开关语义和公式定义，而不是数值表现。",
            ],
            formulas=mask_probe_formula_lines(),
            code_refs=[
                "scripts/orig_vggt_mask_boost_probe/common.py",
                "scripts/orig_vggt_mask_boost_probe/run_task.py",
                "finetune_vggt_pseudo.py: foreground weighting path",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="mask-boost probe 的本质，是在同一训练入口上切换 foreground weighting 语义。",
            layout="full_image",
        ),
        SlideSpec(
            title="这套脚手架最终产出什么",
            bullets=[
                "统一产出 compare 目录、summary JSON、summary MD 以及可直接被 PPT 引用的中间材料。",
                "因此这一层最核心的贡献，是把分支训练组织成稳定可复用的任务管线和汇总管线。",
                "这层能力本身独立于任何结果判断，可以单独向导师解释。",
            ],
            code_refs=[
                "logs/modal_phase5/reports/*.json",
                "logs/modal_phase5/reports/*.md",
                "tools/generate_orig_vggt_local_code_mods_presentation.py",
            ],
        ),
    ]


def build_mask_deck(data: dict) -> list[SlideSpec]:
    assets = data.get("assets", {})
    return [
        SlideSpec(
            title="Mask / Ghost / Point-Support: 只讲代码语义与公式",
            bullets=[
                "这套单独讲 foreground weighting、ghost scorer 和 point-support / masked error 这条代码分支。",
                "重点是变量如何定义、公式如何展开、流程如何走，不讨论任何结果曲线。",
                "向导师汇报时，这套更像“训练语义与评估脚手架说明”，不是效果汇报。",
            ],
        ),
        SlideSpec(
            title="Foreground Weighting: 完整实现如何从 mask 走到最终权重",
            bullets=[
                "完整流程分成 mask 二值化、预侵蚀、监督域构造、boost 区域构造和最终权重合成五步。",
                "bg floor 决定背景监督是否保留底权重；当前 probe 因为 bg floor = 0，所以语义上是 foreground-only supervision。",
                "因此当前实现不能简化成一句“前景增强”，它实际上是“前景门控 + 内部增权”的组合。",
            ],
            formulas=fg_formula_lines(),
            code_refs=[
                "finetune_vggt_pseudo.py: _normalize_mask_binary_np",
                "finetune_vggt_pseudo.py: _erode_mask_tensor",
                "finetune_vggt_pseudo.py: _build_fg_supervision_boost_mask",
                "finetune_vggt_pseudo.py: _apply_fg_supervision_boost",
            ],
            image_paths=existing_images(assets.get("fg_boost_semantics")),
            image_caption="mask 权重路径的关键是先定义监督域, 再定义 boost 区域, 最后合成最终权重。",
            layout="full_image",
        ),
        SlideSpec(
            title="Ghost 评分: 完整公式如何展开",
            bullets=[
                "ghost 评分输入不是任意图片，而是固定的 triplet 结构: mask、prediction、target。",
                "评分先从 mask 和 prediction 提取宽度、面积和中心，再通过横向投影统计峰数，最后叠加暗部惩罚和塌缩惩罚。",
                "因此它量化的是视觉形态是否变宽、变厚、双峰、偏移、过黑或塌缩，而不是训练 loss 本身。",
            ],
            formulas=ghost_formula_lines(),
            code_refs=[
                "tools/score_ghosting_from_cat_pred.py: _score_one",
                "tools/score_ghosting_from_cat_pred.py: _count_peaks",
            ],
            image_paths=existing_images(assets.get("ghost_scoring_formula")),
            image_caption="ghost_visual_score 由基础形态分数、暗部惩罚和塌缩惩罚三部分组成。",
            layout="full_image",
        ),
        SlideSpec(
            title="Point-support / Masked Error: 结构指标和误差指标如何并列定义",
            bullets=[
                "point-support 读的是 support 结构量，masked PSNR / masked L1 读的是主体区域内的重建误差，这两类量不能混读。",
                "脚本先从 support 图取绿色通道，再结合主体 mask 统计主体份额、连通块集中度和峰数。",
                "随后再单独在主体 mask 内比较 prediction 与 target，得到 subject_psnr 和 subject_l1。",
            ],
            formulas=support_formula_lines(),
            code_refs=[
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
                "scripts/orig_vggt_viewcount/render_raw_compare.py",
            ],
            image_paths=existing_images(assets.get("point_support_formula")),
            image_caption="support 结构量和主体 mask 内误差量分开定义, 但在同一测量脚本中统一输出。",
            layout="full_image",
        ),
        SlideSpec(
            title="这条代码分支放回整条主线里怎么理解",
            bullets=[
                "foreground weighting 改的是训练权重语义，ghost scorer 和 point-support measure 改的是 compare 后的测量语义。",
                "它们共同组成的是一层围绕原版主线的训练与评估壳层，而不是新的 backbone 结构。",
                "因此这套汇报的重点，应该始终是“代码如何定义语义和指标”，而不是“结果看起来怎样”。",
            ],
            formulas=[
                "training semantics: foreground weighting",
                "evaluation semantics: ghost scoring + point-support + masked error",
            ],
            code_refs=[
                "finetune_vggt_pseudo.py",
                "tools/score_ghosting_from_cat_pred.py",
                "scripts/orig_vggt_stepcurve_probe/measure_point_support.py",
            ],
            image_paths=existing_images(assets.get("compare_probe_framework")),
            image_caption="foreground weighting 属于训练壳层, ghost 和 point-support 属于 compare 之后的测量壳层。",
            layout="full_image",
        ),
        SlideSpec(
            title="向导师的一句话总结",
            bullets=[
                "Mask 这条线真正新增的是监督权重语义，Ghost 这条线真正新增的是视觉形态评分脚本，Point-support 这条线真正新增的是结构测量和 mask 内误差测量。",
                "三者共同组成了一条围绕原版主线的训练语义与评估语义壳层。",
                "所以这套 deck 应该讲公式、变量和流程，不应该讲效果。",
            ],
        ),
    ]


def build_visual_prompts_md() -> str:
    return textwrap.dedent(
        """
        # 原版 VGGT PPT 插图生成 Prompt

        ## 总原则

        - 提示词只用模块职责、变量定义、公式和信息流来写图，不在图里堆源文件名。
        - 重点始终是“相对原版，我们补了哪些训练与分析层”，不是把官方图换一遍文字。
        - 所有图都只讲结构、原理、公式和流程，不讲实验效果，不讲优劣判断。

        ## 1. 当前本地架构图

        - 建议文件名: `current_local_vggt_architecture.png`
        - 重点:
        - 上层保留原版核心主干: Multi-view Inputs, Patch Tokens, Camera Token, Global Aggregation, Frame Aggregation, Camera Head, Dense Geometry Head
        - 下层是本地新增壳层: Geometry Teacher, Geometry Cache, Pseudo-geometry Dataset, Supervision Router, Visual Compare, Ghost Measure, Point-support Measure, Probe Automation, Structured Reports
        - 必须让人一眼看懂: 原版核心保留, 主体新增工作发生在原版输出之后

        ```text
        Create a clean 16:9 academic architecture diagram with a warm white background and paper-figure quality.

        Use the composition idea of an official model architecture figure: native core on the upper lane, added local shell on the lower or outer lane. The image must describe the CURRENT LOCAL PIPELINE around original VGGT, not the native VGGT figure itself.

        Upper lane:
        - Multi-view Inputs
        - Patch Tokens and a small Camera Token
        - Global Aggregation
        - Frame Aggregation
        - Camera Head
        - Dense Geometry Head

        Lower or outer-shell lane:
        - Geometry Teacher
        - Geometry Cache
        - Pseudo-geometry Dataset
        - Supervision Router
        - Visual Compare
        - Ghost Measure
        - Point-support Measure
        - Probe Automation
        - Structured Reports

        Add compact formula callouts:
        - P_world = R^T * (P_cam - t)
        - w_conf(i) = M_sup(i) * c_hat(i)^gamma
        - w_final(i) = w_base(i) * [1 + (beta - 1) * M_boost(i)]
        ```

        ## 2. Geometry Teacher / Cache 流程图

        - 建议文件名: `teacher_npz_flow.png`
        - 重点:
        - 原版输出先变成 Geometry Teacher
        - 再经过反投影生成 Point Map / Frame Meta / Camera Meta
        - 再被冻结成 Geometry Cache 供后续训练读取

        ```text
        Create a clean 16:9 flow diagram for a geometry-teacher pipeline.

        Pipeline:
        - Original Outputs: Depth / Confidence / Camera
        - Geometry Teacher
        - Unprojection
        - Geometry Products: Point Map / Frame Meta / Camera Meta
        - Geometry Cache
        - Dataset Read Path

        Show formulas:
        - x_cam = ((u + o - c_x) / f_x) * z
        - y_cam = ((v + o - c_y) / f_y) * z
        - P_world = R^T * (P_cam - t)
        ```

        ## 3. Supervision Router 图

        - 建议文件名: `loss_routing.png`
        - 重点:
        - 中心必须是一个显眼的 Supervision Router
        - 左侧是 Pseudo-geometry Dataset
        - 右侧扇出 Depth, Point Map, Reprojection, Normal, Camera, Multiview Depth
        - 图上必须出现展开公式

        ```text
        Create a clean 16:9 supervision-routing diagram in research style.

        Layout:
        - left: Pseudo-geometry Dataset
        - center: a large Supervision Router
        - right: Depth, Point Map, Reprojection, Normal, Camera, Multiview Depth
        - far right: L_total

        Include expanded formulas:
        - M_sup(i) = M_valid_all(i) * [M_fg(i) + eta * (1 - M_fg(i))]
        - c_hat(i) = clip((c(i) - t) / (1 - t), 0, 1)
        - w_conf(i) = M_sup(i) * c_hat(i)^gamma * Q_pvq(i) * G_dyn(i)
        - L_total = sum_k lambda_k L_k
        ```

        ## 4. Foreground Weighting 语义图

        - 建议文件名: `fg_boost_semantics.png`
        - 重点:
        - 这是权重语义图, 不是网络结构图
        - 顺序必须是 Valid Region -> Foreground Mask -> Foreground Gating -> Interior Region -> Final Weight
        - 图上必须保留二值化、侵蚀和最终权重公式

        ```text
        Create a minimal but explicit 16:9 semantic diagram for foreground-based supervision weighting.

        Pipeline:
        - Valid Region
        - Foreground Mask
        - Foreground Gating
        - Interior Region
        - Final Weight

        Include formulas:
        - M_fg_raw(i) = 1[x(i) >= 0.5] if max(x) <= 1.5 else 1[x(i) > 0]
        - M_fg(i) = erode(M_fg_raw, r0) if pre-erosion is enabled else M_fg_raw(i)
        - M_sup(i) = M_valid_all(i) * [M_fg(i) + eta * (1 - M_fg(i))]
        - M_boost(i) = erode(M_fg(i), r) if region_mode = interior_only else M_fg(i)
        - w_final(i) = w_base(i) * [1 + (beta - 1) * M_boost(i)]
        ```

        ## 5. Ghost 评分公式图

        - 建议文件名: `ghost_scoring_formula.png`
        - 重点:
        - 输入必须是 triplet: mask, prediction, target
        - 中间量必须有 Mask BBox, Pred BBox, Peak Profile
        - 底部必须写完整 ghost_score 与 ghost_visual_score 公式

        ```text
        Create a clean 16:9 formula-driven workflow diagram for ghost scoring.

        Pipeline:
        - Ghost Triplet
        - Mask BBox
        - Pred BBox
        - Peak Profile
        - Final ghost_visual_score

        Include formulas:
        - width_ratio = W_pred / W_mask
        - area_ratio = A_pred / A_mask
        - center_offset_ratio = |c_pred - c_mask| / W_mask
        - ghost_score = 1.00 * max(0, width_ratio - 1.10) + 0.40 * max(0, area_ratio - 1.30) + 0.60 * max(0, peak_count - 1) + 0.50 * max(0, center_offset_ratio - 0.22)
        - dark_penalty = 40 * max(0, 0.045 - pred_luma_mean) + 12 * max(0, 0.10 - pred_nonblack_ratio_thr008)
        - collapse_penalty = 6 * max(0, 0.55 - area_ratio) + 6 * max(0, 0.65 - width_ratio)
        - ghost_visual_score = ghost_score + dark_penalty + collapse_penalty
        ```

        ## 6. Point-support / Masked Error 公式图

        - 建议文件名: `point_support_formula.png`
        - 重点:
        - 左边必须是 Support Map 和 Subject Mask
        - 中间是 Active Components
        - 右边是 Masked Error
        - 图上要同时放结构指标和误差指标的公式

        ```text
        Create a clean 16:9 formula-driven diagram for point-support and masked reconstruction metrics.

        Pipeline:
        - Support Map
        - Subject Mask
        - Active Components
        - Masked Error

        Include formulas:
        - subject_support_share = sum_x S(x) * M(x) / sum_x S(x)
        - largest_component_share = max_c sum_{x in c} S(x) / sum_{x in active} S(x)
        - secondary_component_mass = second_max_c sum_{x in c} S(x) / sum_{x in active} S(x)
        - support_peak_count = Peaks(smooth(sum_y ActiveROI(y, x)))
        - subject_psnr = -10 log10(sum_x ||pred(x) - tgt(x)||_2^2 * M(x) / (3 * sum_x M(x)))
        - subject_l1 = sum_x ||pred(x) - tgt(x)||_1 * M(x) / (3 * sum_x M(x))
        ```
        """
    ).strip() + "\n"


def build_mask_boost_ghost_explainer_md() -> str:
    return textwrap.dedent(
        """
        # 原版 VGGT Mask / Ghost / Point-Support 公式展开说明

        ## 范围

        - 这份说明只讲当前代码里这三条语义线怎么定义: foreground weighting、ghost scoring、point-support / masked error。
        - 重点是变量定义、完整公式、原理和流程，不讨论任何结果表现。

        ## 1. Foreground Weighting

        ### 1.1 变量定义

        - `x(i)`: 数据集原始 mask 像素值
        - `M_fg_raw(i)`: 原始前景二值 mask
        - `M_fg(i)`: 进入监督路由之前实际使用的前景 mask
        - `M_valid_all(i)`: 由目标深度定义的几何有效域
        - `eta`: `fg_supervision_bg_floor`
        - `beta`: `fg_supervision_boost`
        - `r0`: `fg_mask_erode_px`
        - `r`: `fg_supervision_region_erode_px`

        ### 1.2 完整公式

        ```text
        M_fg_raw(i) =
            1[x(i) >= 0.5],   if max(x) <= 1.5
            1[x(i) > 0],      if max(x) > 1.5
        ```

        ```text
        M_fg(i) =
            erode(M_fg_raw, r0)(i),   if r0 > 0
            M_fg_raw(i),              otherwise
        ```

        ```text
        M_valid_all(i) = 1[z_tgt(i) > 1e-6]
        ```

        ```text
        M_sup(i) =
            M_valid_all(i) * M_fg(i),                                  if eta = 0
            M_valid_all(i) * [M_fg(i) + eta * (1 - M_fg(i))],          if eta > 0
        ```

        ```text
        c_hat(i) =
            clip((c(i) - t) / (1 - t), 0, 1),   if t > 0
            c(i),                               if t = 0
        ```

        ```text
        w_conf0(i) = M_sup(i) * c_hat(i)^gamma
        w_conf(i)  = w_conf0(i) * Q_pvq(i) * G_dyn(i)
        ```

        ```text
        w_base(i) =
            M_sup(i),                                       if mode = uniform
            M_sup(i) * [alpha * w_conf(i) + (1 - alpha)],  if mode = mix
            w_conf(i),                                      if mode = conf
        ```

        ```text
        M_boost(i) =
            erode(M_fg(i), r),   if region_mode = interior_only
            M_fg(i),             if region_mode = all
        ```

        ```text
        erode(x) = 1 - dilate(1 - x)
        ```

        ```text
        w_final(i) = w_base(i) * [1 + (beta - 1) * M_boost(i)]
        ```

        ### 1.3 当前 probe 的特例

        ```text
        use_fg_mask = on
        fg_supervision_bg_floor = 0
        fg_supervision_region_mode = interior_only
        fg_supervision_region_erode_px = 2
        fg_supervision_boost in {2, 3, 5}
        ```

        因此当前 probe 化简成:

        ```text
        M_sup(i)      = M_valid_all(i) * M_fg(i)
        M_interior(i) = erode(M_fg(i), 2)
        w_final(i)    = w_base(i) * [1 + (beta - 1) * M_interior(i)]
        alpha = 1, 2, 4 -> beta = 2, 3, 5
        ```

        ## 2. Ghost Scoring

        ### 2.1 中间量

        ```text
        width_ratio         = W_pred / W_mask
        area_ratio          = A_pred / A_mask
        center_offset_ratio = |c_pred - c_mask| / W_mask
        peak_count          = Peaks(smooth(sum_y MidBin(y, x)))
        ```

        ### 2.2 完整公式

        ```text
        ghost_score =
            1.00 * max(0, width_ratio - 1.10)
          + 0.40 * max(0, area_ratio - 1.30)
          + 0.60 * max(0, peak_count - 1)
          + 0.50 * max(0, center_offset_ratio - 0.22)
        ```

        ```text
        dark_penalty =
            40 * max(0, 0.045 - pred_luma_mean)
          + 12 * max(0, 0.10 - pred_nonblack_ratio_thr008)
        ```

        ```text
        collapse_penalty =
            6 * max(0, 0.55 - area_ratio)
          + 6 * max(0, 0.65 - width_ratio)
        ```

        ```text
        ghost_visual_score = ghost_score + dark_penalty + collapse_penalty
        ```

        ```text
        ghost_soft_score =
            0.80 * max(0, width_ratio - 1.00)
          + 0.25 * max(0, area_ratio - 1.10)
          + 0.50 * max(0, peak_count - 1)
          + 0.35 * max(0, center_offset_ratio - 0.10)
        ```

        ## 3. Point-support 与 Masked Error

        ### 3.1 结构指标

        ```text
        S(x) = green_channel(weight_native.png)
        T_support = max(threshold_floor, quantile_0.75({S(x) | S(x) > 0}))
        Active(x) = 1[S(x) >= T_support]
        ```

        ```text
        subject_support_share = sum_x S(x) * M(x) / sum_x S(x)
        outside_subject_support_share = sum_x S(x) * (1 - M(x)) / sum_x S(x)
        largest_component_share = max_c sum_{x in c} S(x) / sum_{x in active} S(x)
        secondary_component_mass = second_max_c sum_{x in c} S(x) / sum_{x in active} S(x)
        support_peak_count = Peaks(smooth(sum_y ActiveROI(y, x)))
        ```

        ### 3.2 mask 内误差指标

        ```text
        subject_psnr = -10 log10(sum_x ||pred(x) - tgt(x)||_2^2 * M(x) / (3 * sum_x M(x)))
        subject_l1   = sum_x ||pred(x) - tgt(x)||_1 * M(x) / (3 * sum_x M(x))
        ```

        ## 4. 原理与流程总结

        - foreground weighting 先定义监督域，再定义 boost 区域，最后合成最终权重。
        - ghost scoring 先做几何外形统计，再做暗部与塌缩惩罚。
        - point-support 脚本同时输出结构量和 mask 内误差量，两者不能混读。
        """
    ).strip() + "\n"


def main() -> None:
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    download_official_architecture()
    data = collect_data()
    data["assets"] = generate_visual_assets(data)
    diff_summary = collect_code_diff_summary()

    deck_jobs = [
        ("原版 VGGT 本地代码改动总览", OUT_PPTX, OUT_MD, build_combined_deck(data, diff_summary)),
        ("原版 VGGT 改动与微调代码主线", OUT_CODE_PPTX, OUT_CODE_MD, build_code_and_finetune_deck(data, diff_summary)),
        ("原版 VGGT Probe 自动化脚手架", OUT_PROBE_PPTX, OUT_PROBE_MD, build_probe_deck(data)),
        ("原版 VGGT Mask / Ghost / Point-Support 代码分支", OUT_MASK_PPTX, OUT_MASK_MD, build_mask_deck(data)),
    ]

    for deck_title, pptx_path, md_path, slides in deck_jobs:
        actual_pptx, actual_md = write_deck_bundle(slides, pptx_path, md_path, deck_title)
        print(f"Wrote PPTX: {actual_pptx}")
        print(f"Wrote outline: {actual_md}")

    prompt_path = save_with_fallback_bytes(
        OUT_PROMPTS_MD,
        lambda p: p.write_text(build_visual_prompts_md(), encoding="utf-8"),
    )
    print(f"Wrote prompts: {prompt_path}")

    explainer_path = save_with_fallback_bytes(
        MASK_BOOST_EXPLAINER_MD,
        lambda p: p.write_text(build_mask_boost_ghost_explainer_md(), encoding="utf-8"),
    )
    print(f"Wrote explainer: {explainer_path}")


if __name__ == "__main__":
    main()
